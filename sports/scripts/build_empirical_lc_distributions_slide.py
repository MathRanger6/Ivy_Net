#!/usr/bin/env python3
"""Build two-slide PD17 empirical L_C deck (slides 2–3 in HAND block).

Slide text = literal LaTeX in Calibri — highlight → Insert → Equation → LaTeX to math.

Run (repo root):
  python sports/scripts/build_empirical_lc_distributions_slide.py
  python sports/scripts/build_empirical_lc_distributions_slide.py --slides-only
  python sports/scripts/build_empirical_lc_distributions_slide.py --gamma 3
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_EMPIRICAL_LC_DECK, EMPIRICAL_PD17, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "empirical_lc_distributions.py"
FIG_1D = EMPIRICAL_PD17 / "EMPIRICAL_L_C_distribution.png"
FIG_2D = EMPIRICAL_PD17 / "EMPIRICAL_L_C_vs_Tj_2d.png"
META = EMPIRICAL_PD17 / "EMPIRICAL_L_C_meta.json"
OUT_PPTX = AUTO_EMPIRICAL_LC_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.2)
SIDEBAR_W = Inches(3.35)
FIG_W = CONTENT_W - SIDEBAR_W - COL_GAP

CLAIM_1D = (
    "Claim (PD17): Real rosters carry a spread of team congestion "
    r"L_C — descriptive score-side input before we pin \rho, \theta, and \gamma in sim."
)

CLAIM_2D = (
    "Claim (PD17): On actual team-seasons, stronger realized rosters "
    r"(\hat{T}_{j}) tend to co-occur with higher L_C — empirical Sketch A."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figures(*, gamma: float | None = None) -> None:
    print("Regenerating empirical L_C figures ...")
    cmd = [sys.executable, str(DIAG_SCRIPT)]
    if gamma is not None:
        cmd.extend(["--gamma", str(gamma)])
    subprocess.run(cmd, cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def _add_header(slide, *, title: str, meta: dict) -> float:
    """Title + subtitle; return body start y in inches."""
    title_top = 0.24
    title_h = 0.5
    title_box = slide.shapes.add_textbox(MARGIN, Inches(title_top), CONTENT_W, Inches(title_h))
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + 0.04
    sub_h = 0.27
    sub_box = slide.shapes.add_textbox(MARGIN, Inches(sub_top), CONTENT_W, Inches(sub_h))
    seasons = meta.get("seasons", "2011-2021")
    theta = meta.get("theta")
    gamma = meta.get("gamma", 10)
    kn = meta.get("theta_K_over_N", {})
    k_over_n = kn.get("K_over_N")
    sub_parts = [
        rf"MBB {seasons} · team smooth L_C · PPM z within season · min 20 min",
    ]
    if theta is not None and k_over_n is not None:
        sub_parts.append(
            rf"\theta = F^{{-1}}_{{\hat{{A}}}}(1-K/N) = {theta:.3f} "
            rf"(K/N={k_over_n:.4f}) · \gamma={gamma:g} placeholder"
        )
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        " · ".join(sub_parts),
        font_size=HAND_SUBTITLE_PT,
    )
    return sub_top + sub_h + 0.1


def _add_sidebar_figure_slide(
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    fig_path: Path,
    claim: str,
    meta: dict,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body_top = _add_header(slide, title=title, meta=meta)

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_h = footer_top - Inches(body_top) - Inches(0.08)

    side_box = slide.shapes.add_textbox(MARGIN, Inches(body_top), SIDEBAR_W, body_h)
    side_tf = side_box.text_frame
    side_tf.word_wrap = True
    fill_bullets_raw_latex(side_tf, bullets, font_size=HAND_BODY_PT)

    _add_picture_fitted(
        slide,
        fig_path,
        MARGIN + SIDEBAR_W + COL_GAP,
        Inches(body_top),
        FIG_W,
        body_h,
    )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], claim, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slides-only",
        action="store_true",
        help="Skip PNG regen (text/layout only)",
    )
    parser.add_argument(
        "--gamma",
        type=float,
        default=None,
        metavar="G",
        help="Pass through to empirical_lc_distributions.py (viability sharpness γ)",
    )
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figures(gamma=args.gamma)
    meta = _load_meta()
    lc = meta.get("L_C", {})
    tj = meta.get("T_j_hat", {})
    kn = meta.get("theta_K_over_N", {})
    theta = meta.get("theta")
    gamma = meta.get("gamma", 10)

    bullets_1d = [
        r"Team L_C = mean_{j} \sigma(\gamma(\hat{A}_{j} - \theta)) on each roster.",
        r"Team-smooth — one L_C per team-season, shared by all roster players.",
        rf"\gamma = {gamma:g} (539 placeholder until PD17 \gamma / \lambda calibration).",
    ]
    if lc:
        bullets_1d.append(
            rf"This run: mean L_C={lc.get('mean', 0):.3f}, "
            rf"sd={lc.get('std', 0):.3f} ({lc.get('n', 0):,} team-seasons)."
        )
    if theta is not None and kn:
        bullets_1d.append(
            rf"\theta from draft rate: K/N={kn.get('K_over_N', 0):.4f} "
            rf"\Rightarrow \theta={theta:.3f} on \hat{{A}}_{{i}}."
        )

    bullets_2d = [
        r"Each cell: count of team-seasons at (\hat{T}_{j}, L_C).",
        r"Color = count (empirical Sketch A — no \rho arms).",
        r"\hat{T}_{j} = mean \hat{A}_{i} on roster (realized talent, not T_{j^*}).",
        r"Look for upward co-movement: higher \hat{T}_{j} with higher L_C.",
    ]
    if lc and tj:
        bullets_2d.append(
            rf"Panel: mean \hat{{T}}_{{j}}={tj.get('mean', 0):.3f}, "
            rf"mean L_C={lc.get('mean', 0):.3f}."
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_sidebar_figure_slide(
        prs,
        title=r"PD17 — Empirical team L_C distribution",
        bullets=bullets_1d,
        fig_path=FIG_1D,
        claim=CLAIM_1D,
        meta=meta,
    )
    _add_sidebar_figure_slide(
        prs,
        title=r"PD17 — Empirical Sketch A: \hat{T}_{j} vs L_C",
        bullets=bullets_2d,
        fig_path=FIG_2D,
        claim=CLAIM_2D,
        meta=meta,
    )

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (2 slides — PD17 empirical L_C block)")


if __name__ == "__main__":
    main()
