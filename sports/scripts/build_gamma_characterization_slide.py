#!/usr/bin/env python3
"""Build two-slide γ characterization deck (sweep + λ_crit explainer).

Run (repo root):
  python sports/scripts/build_gamma_characterization_slide.py

Outputs:
  HEROs_and_PASSes/slides/auto/CHAR_gamma_characterization_AUTO.pptx
    slide 1 — γ sweep on sort-and-chop (λ curves by γ)
    slide 2 — λ_crit ≈ 4/γ educational figure
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from hero_gallery_paths import AUTO_GAMMA_DECK, SORT_CHOP_LAMBDA, ensure_hero_dirs

OUT_PPTX = AUTO_GAMMA_DECK
GAMMA_SCRIPT = SCRIPTS / "gamma_sweep_diagnostic.py"
LAMBDA_FIG_SCRIPT = SCRIPTS / "build_lambda_gamma_threshold_figure.py"
FIG_SWEEP = SORT_CHOP_LAMBDA / "GAMMA_sweep_lambda_curves_key_arms.png"
FIG_SWEEP_FULL = SORT_CHOP_LAMBDA / "GAMMA_sweep_lambda_curves.png"
FIG_THRESHOLD = SORT_CHOP_LAMBDA / "LAMBDA_threshold_gamma_viability.png"

from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_W = Inches(4.55)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

CLAIM_SWEEP = (
    "Claim: $\\gamma$ (sigmoid slope in $L_C$) shifts when congestion bites — "
    "higher $\\gamma$ lowers $\\lambda_{\\mathrm{crit}} \\approx 4/\\gamma$ on sort-and-chop; "
    "below crit, $\\lambda$ arms can look identical."
)

CLAIM_THRESHOLD = (
    "Claim: $\\lambda_{\\mathrm{crit}} \\approx 4/\\gamma$ is the first $\\lambda$ where "
    "score reorder beats talent-only on sort-and-chop — not when the curve "
    "visually departs from $\\lambda=0$ in coarse bins."
)

SLIDE_VARIANTS = [
    {
        "title": r"Phase B — Characterize $\gamma$ (viability sharpness) — selection curves",
        "figure": FIG_SWEEP,
        "claim": CLAIM_SWEEP,
        "params_kind": "sweep",
        "notes": [
            r"Sort-and-chop benchmark: same rosters per $\gamma$; only $\gamma$ and $\lambda$ in score change.",
            r"Figure uses key arms $\lambda \in \{0, 0.55, 1.0\}$ — full five-arm grid in \texttt{GAMMA\_sweep\_lambda\_curves.png}.",
            r"$\gamma=5$: $\lambda_{\mathrm{crit}}\approx 0.8$ — hump needs larger $\lambda$.",
            r"$\gamma=10$ (539): $\lambda_{\mathrm{crit}}\approx 0.4$ — curve morph near $\lambda \gtrsim 0.4$.",
            r"Below crit: monotone elite-edge shape; above crit: inverted-$U$ / hump.",
        ],
    },
    {
        "title": r"Phase B — $\lambda_{\mathrm{crit}} \approx 4/\gamma$ (sort-and-chop algebra)",
        "figure": FIG_THRESHOLD,
        "claim": CLAIM_THRESHOLD,
        "params_kind": "threshold",
        "notes": [
            r"Soft viability: $\sigma(\gamma(A-\theta))$ — not a hard step; $\gamma$ sets knee steepness.",
            r"On sort-and-chop, $A_i \parallel L_C$ within teams → small $\lambda$ may not reorder top-$K$.",
            r"$\lambda_{\mathrm{crit}}$: first $\lambda$ where $S_i = A_i - \lambda L_C$ changes who gets selected.",
            r"Visual curve departure in $16$ bins can lag $\lambda_{\mathrm{crit}}$ — use finer bins or check reorder.",
            r"Ties to Alex Q: why $\lambda=0$ and $\lambda=0.25$ match until $\lambda \gtrsim 4/\gamma$.",
        ],
    },
]


def _regenerate_figures() -> None:
    for script, label in (
        (GAMMA_SCRIPT, "γ sweep"),
        (LAMBDA_FIG_SCRIPT, "λ_crit explainer"),
    ):
        print(f"Regenerating {label} ...")
        subprocess.run([sys.executable, str(script)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2


def _add_params_column(slide, *, left, top, width, height, kind: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    head = tf.paragraphs[0]
    populate_paragraph_with_latex(
        head,
        r"Characterize $\gamma$ — one-at-a-time (539 sort-and-chop benchmark)",
        font_size=11,
    )
    head.font.bold = True
    head.space_after = 6

    eq = tf.add_paragraph()
    eq.text = ""
    populate_paragraph_with_latex(
        eq,
        r"$\sigma(x) = 1/(1+e^{-x})$, $x = \gamma(A-\theta)$",
        font_size=11,
    )
    eq.space_after = 4

    if kind == "sweep":
        bullets = [
            r"$\theta=0.72$ fixed; sort-and-chop assign",
            r"arms: $\gamma \in \{5, 10, 20\}$",
            r"within each $\gamma$: $\lambda \in \{0, 0.25, 0.55, 0.75, 1.0\}$",
            r"VISUALIZE = mean $Y_{\mathrm{selected}}$ vs pool mean ($16$ bins)",
        ]
    else:
        bullets = [
            r"Educational: how $\gamma$ shapes $L_C$ along ability axis",
            r"vertical lines: $\lambda_{\mathrm{crit}} = 4/\gamma$ for each $\gamma$",
            r"539 defaults: $\theta=0.72$, $\gamma=10$ → crit $= 0.4$",
            r"Pairs with sort-and-chop $\lambda$ threshold memo (doc 06)",
        ]

    for line in bullets:
        para = tf.add_paragraph()
        para.text = ""
        populate_paragraph_with_latex(para, line, font_size=10)
        para.space_after = 2


def _add_slide(prs: Presentation, variant: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.28)
    title_h = Inches(0.46)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_with_latex(
        title_box.text_frame.paragraphs[0],
        variant["title"],
        font_size=20,
    )
    title_box.text_frame.paragraphs[0].font.bold = True

    body_top = title_top + title_h + Inches(0.1)
    footer_h = Inches(0.52)
    footer_top = SLIDE_H - MARGIN - footer_h
    notes_h = Inches(1.35)
    params_h = footer_top - body_top - notes_h - Inches(0.08)

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_params_column(
        slide,
        left=left_x,
        top=body_top,
        width=LEFT_W,
        height=params_h,
        kind=variant["params_kind"],
    )
    _add_picture_fitted(slide, variant["figure"], right_x, body_top, RIGHT_W, params_h)

    notes_box = slide.shapes.add_textbox(
        left_x,
        footer_top - notes_h - Inches(0.06),
        LEFT_W + RIGHT_W + COL_GAP,
        notes_h,
    )
    fill_bullets_latex(notes_box.text_frame, variant["notes"], font_size=9)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_with_latex(
        foot.text_frame.paragraphs[0], variant["claim"], font_size=11
    )
    foot.text_frame.paragraphs[0].font.bold = True
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    ensure_hero_dirs()
    _regenerate_figures()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for variant in SLIDE_VARIANTS:
        _add_slide(prs, variant)

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(SLIDE_VARIANTS)} slides)")


if __name__ == "__main__":
    main()
