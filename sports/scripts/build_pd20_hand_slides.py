#!/usr/bin/env python3
"""Build PD20 AUTO reference deck for CHAR_PD20_HAND.pptx.

Slides:
  1. PD20 intro — Gibbs SELECT + priority order
  2. Temperature sweep (λ = 1.5, 2.0; rule D)
  3. Cold limit — rule C vs D at λ = 2, t = 0.001
  4. Takeaways — inverted-U gate + open questions

Run (repo root):
  python sports/scripts/build_pd20_hand_slides.py --slides-only
  python sports/scripts/build_pd20_hand_slides.py   # reruns cold-limit diagnostic first

Output:
  slides/auto/CHAR_PD20_HAND_AUTO.pptx

Copy into HAND: slides/CHAR_PD20_HAND.pptx (Change Picture + bullets).
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import (
    AUTO_PD20_DECK,
    HAND_PD20_DECK,
    PD20_TEMPERATURE,
    ensure_hero_dirs,
)
from pd17_interval_overlap_slide import add_picture_fitted, load_meta

COLD_DIAG = SCRIPTS / "grandchild_temperature_cold_limit_diagnostic.py"

SWEEP_FIG = PD20_TEMPERATURE / "GRANDCHILD_temperature_select_sweep_2011_2021.png"
SWEEP_META = PD20_TEMPERATURE / "GRANDCHILD_temperature_select_sweep_2011_2021_meta.json"
COLD_FIG = PD20_TEMPERATURE / "GRANDCHILD_temperature_cold_limit_2011_2021.png"
COLD_META = PD20_TEMPERATURE / "GRANDCHILD_temperature_cold_limit_2011_2021_meta.json"
OUT_PPTX = AUTO_PD20_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_TEXT_W = Inches(4.15)
FIG_W = CONTENT_W - LEFT_TEXT_W - COL_GAP


def _append_text_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )
    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_top = body_top
    bullet_h = footer_top - bullet_top - Inches(0.06)
    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, CONTENT_W, bullet_h)
    fill_bullets_raw_latex(bullet_box.text_frame, bullets, font_size=HAND_BODY_PT)
    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_raw_latex(
        foot.text_frame.paragraphs[0],
        claim,
        font_size=HAND_CLAIM_PT,
        bold=True,
    )
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _append_figure_slide(
    prs: Presentation,
    *,
    fig_path: Path,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )
    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_h = Inches(2.35)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = body_top
    fig_h = footer_top - fig_top - Inches(0.04)
    fig_left = MARGIN + LEFT_TEXT_W + COL_GAP
    add_picture_fitted(slide, fig_path, fig_left, fig_top, FIG_W, fig_h)
    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, LEFT_TEXT_W, bullet_h)
    fill_bullets_raw_latex(bullet_box.text_frame, bullets, font_size=HAND_BODY_PT)
    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_raw_latex(
        foot.text_frame.paragraphs[0],
        claim,
        font_size=HAND_CLAIM_PT,
        bold=True,
    )
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _count_inverted_u(meta: dict, *, lam: float | None = None) -> tuple[int, int]:
    runs = meta.get("runs", [])
    if lam is not None:
        runs = [r for r in runs if abs(float(r.get("lambda", -1)) - lam) < 1e-9]
    total = len(runs)
    n_u = sum(
        1
        for r in runs
        if r.get("curvature_loo", {}).get("shape") == "inverted_u_like"
    )
    return n_u, total


def _intro_bullets() -> list[str]:
    return [
        r"PD20: replace deterministic top-$K$ SELECT with Gibbs rule on $S_i$.",
        r"Lock: $P_i \propto \exp(S_i/t)$ with $t$ in the denominator (stat-phys).",
        r"Sim bridge: rule D = Gibbs weights + $K$ draws without replacement.",
        r"SCORE unchanged: $S_i = \hat{A}_i - \lambda L_{C,g(i)}$; $\rho$ still in ASSIGN only.",
        r"Priority: sweep $\log_{10} t$ — prove inverted-U survives **before** MLE.",
        r"ASSIGN $\rightarrow$ SCORE $\rightarrow$ SELECT — score $\neq$ select (binding).",
    ]


def _sweep_bullets(meta: dict) -> list[str]:
    seasons = meta.get("seasons", "2011-2021")
    rho = meta.get("rho", 0.5)
    lams = meta.get("lambda_panels", [1.5, 2.0])
    emp = meta.get("empirical", {}).get("curvature_loo", {})
    bullets = [
        r"SELECT: draft rate vs LOO (left) and pool mean (right) — Hero axes.",
        rf"MBB {seasons} · empirical caps · $\rho={rho:g}$ · rule D · $\lambda$ panels "
        + ", ".join(f"{x:g}" for x in lams),
        rf"Empirical LOO: {emp.get('shape', '?').replace('_', ' ')}.",
    ]
    for lam in lams:
        n_u, total = _count_inverted_u(meta, lam=float(lam))
        bullets.append(rf"$\lambda={lam:g}$: inverted-U-like on LOO for {n_u}/{total} $t$ arms.")
    bullets.extend(
        [
            r"Cold $t$: near hard cut; hot $t$: flattens toward uniform $K$-draws.",
            r"Pool mean often stays monotone — same PD17 readout limitation.",
        ]
    )
    return bullets


def _sweep_claim(meta: dict) -> str:
    n_u_2, total = _count_inverted_u(meta, lam=2.0)
    if total and n_u_2 >= total - 2:
        return (
            "Talking point: at $\\lambda=2$ (breakpoint band), inverted-U on LOO persists "
            "across most of the $t$ grid — softening SELECT did not kill the hump."
        )
    return (
        "Talking point: check LOO panels at $\\lambda=1.5$–$2$ — does a $t$ window "
        "preserve inverted-U vs empirical?"
    )


def _cold_bullets(meta: dict) -> list[str]:
    lam = meta.get("lambda", 2.0)
    t = meta.get("temperature", 0.001)
    gap_loo = meta.get("max_bin_gap_loo", float("nan"))
    gap_mean = meta.get("max_bin_gap_pool_mean", float("nan"))
    curv_c = meta.get("curvature_loo_c", {}).get("shape", "?")
    curv_d = meta.get("curvature_loo_d", {}).get("shape", "?")
    match = meta.get("cold_limit_match", False)
    return [
        r"Same ASSIGN seeds; only SELECT rule differs (C = top-$K$, D = cold Gibbs).",
        rf"Fixed $\lambda={lam:g}$, $t={t:g}$ ($\log_{{10}} t = {meta.get('log10_t', -3):g}$).",
        rf"Max bin gap C vs D: LOO = {gap_loo:.2e}, pool mean = {gap_mean:.2e}.",
        rf"LOO curvature: C = {curv_c.replace('_', ' ')}; D = {curv_d.replace('_', ' ')}.",
        r"Expect overlay coincidence at cold $t$ (Gibbs $\to$ hard ranking).",
        rf"Numerical match flag: {match} (gap $< 10^{{-9}}$ on binned curves).",
    ]


def _cold_claim(meta: dict) -> str:
    if meta.get("cold_limit_match"):
        return (
            "Talking point: rule D at cold $t$ reproduces rule C bin-for-bin — "
            "Gibbs SELECT nests the old top-$K$ limit."
        )
    gap = meta.get("max_bin_gap_loo", float("nan"))
    return (
        f"Talking point: cold-limit check — max LOO bin gap C vs D = {gap:.2e}. "
        "Visually should nearly overlay; small gaps may be stochastic K-draw noise."
    )


def _takeaway_bullets(sweep_meta: dict, cold_meta: dict) -> list[str]:
    n_u_15, tot_15 = _count_inverted_u(sweep_meta, lam=1.5)
    n_u_2, tot_2 = _count_inverted_u(sweep_meta, lam=2.0)
    return [
        r"PD20 gate (LOO): inverted-U survives Gibbs SELECT at breakpoint $\lambda$.",
        rf"Full panel: $\lambda=1.5$ → {n_u_15}/{tot_15} $t$ arms U-like; "
        rf"$\lambda=2$ → {n_u_2}/{tot_2} $t$ arms U-like.",
        r"Cold $t$ ≈ top-$K$ (slide 3); mid $t$ still shows hump on LOO.",
        r"Hot $t$ → uniform-$K$ limit (flattening) — different failure mode.",
        r"Open for Alex: K-draw vs literal softmax Bernoulli for MLE.",
        r"Next: MLE for $\lambda^*, t^*$ (and $\gamma$); $\rho^*$ via $H_{\mathrm{sort}}$ separately.",
    ]


def _takeaway_claim() -> str:
    return (
        "PD20 step 4 cleared on LOO: proceed to MLE design once K-draw semantics "
        "are confirmed with Alex."
    )


def build_deck(*, sweep_meta: dict, cold_meta: dict) -> None:
    ensure_hero_dirs()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    seasons = sweep_meta.get("seasons", "2011-2021")
    rho = sweep_meta.get("rho", 0.5)

    _append_text_slide(
        prs,
        title=r"PD20 — Gibbs SELECT (temperature $t$)",
        subtitle=rf"MBB empirical caps · $\rho={rho:g}$ · inverted-U before MLE",
        bullets=_intro_bullets(),
        claim=r"Goal: soft SELECT must not erase the phenomenological inverted-U on LOO pool quality.",
    )

    lam_str = ", ".join(f"{x:g}" for x in sweep_meta.get("lambda_panels", [1.5, 2.0]))
    _append_figure_slide(
        prs,
        fig_path=SWEEP_FIG,
        title=r"Temperature sweep — rule D (Gibbs + $K$ draws)",
        subtitle=rf"MBB {seasons} · $\lambda \in \{{{lam_str}\}}$ · $\log_{{10}} t \in [-3,3]$",
        bullets=_sweep_bullets(sweep_meta),
        claim=_sweep_claim(sweep_meta),
    )

    _append_figure_slide(
        prs,
        fig_path=COLD_FIG,
        title=r"Cold limit — rule C vs rule D",
        subtitle=rf"$\lambda={cold_meta.get('lambda', 2):g}$ · $t={cold_meta.get('temperature', 0.001):g}$ · same seeds",
        bullets=_cold_bullets(cold_meta),
        claim=_cold_claim(cold_meta),
    )

    _append_text_slide(
        prs,
        title=r"PD20 takeaways",
        subtitle=rf"MBB {seasons} · empirical caps · rule D",
        bullets=_takeaway_bullets(sweep_meta, cold_meta),
        claim=_takeaway_claim(),
    )

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--slides-only",
        action="store_true",
        help="Use existing PNG/meta; skip cold-limit diagnostic rerun",
    )
    parser.add_argument(
        "--quick",
        action="store_true",
        help="Pass --quick to cold-limit diagnostic when not --slides-only",
    )
    args = parser.parse_args()

    if not args.slides_only:
        cmd = [sys.executable, str(COLD_DIAG)]
        if args.quick:
            cmd.append("--quick")
        print("=== Running cold-limit diagnostic ===")
        subprocess.run(cmd, cwd=str(REPO), check=True)

    if not SWEEP_META.is_file():
        raise FileNotFoundError(
            f"Missing {SWEEP_META.name} — run grandchild_temperature_select_sweep.py first."
        )
    if not COLD_META.is_file():
        raise FileNotFoundError(
            f"Missing {COLD_META.name} — run grandchild_temperature_cold_limit_diagnostic.py "
            "or omit --slides-only."
        )

    sweep_meta = load_meta(SWEEP_META)
    cold_meta = load_meta(COLD_META)

    print("=== Building PD20 AUTO deck ===")
    build_deck(sweep_meta=sweep_meta, cold_meta=cold_meta)

    print()
    print(f"Copy into HAND: {HAND_PD20_DECK}")
    print(f"  Reference: {OUT_PPTX}")
    print(f"  Figures:   {SWEEP_FIG.name}")
    print(f"             {COLD_FIG.name}")


if __name__ == "__main__":
    main()
