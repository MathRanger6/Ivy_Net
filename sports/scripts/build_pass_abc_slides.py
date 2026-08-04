#!/usr/bin/env python3
"""Build 3-slide PowerPoint deck for Pass A / B / C gallery figures.

Layout: 16:9 widescreen — title, full-width figure, bullets + footer below plot.
Math in bullets/footers uses $...$ (rendered via gallery_mathtext).

Run (repo root):
  python sports/scripts/build_pass_abc_slides.py

Output:
  3-Master_Plan/re_entry/HEROs_and_PASSes/slides/PASS_ABC_Gallery_Slides.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from hero_gallery_paths import PASS_A, PASS_B, PASS_C_RHO, SLIDES, ensure_hero_dirs

OUT_PPTX = SLIDES / "PASS_ABC_Gallery_Slides.pptx"

from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex

# Standard 16:9 (PowerPoint default widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN

SLIDES = [
    {
        "title": "Pass A — Empirical MBB: Roster Pressure in the Data",
        "image": PASS_A / "PASS_A_empirical_talent_vs_roster_side_by_side.png",
        "bullets": [
            "Real NCAA panel (2011–2021): mean NBA draft rate by ventile.",
            r"Left — talent alone: ability (ppm $z$ within season). Monotone up.",
            r"Right — roster context: leave-one-out teammate quality ($\mathrm{poolq\_loo}$). Inverted-U.",
            r"No $\lambda$ in the empirical story — stylized fact Pass B/C echo.",
            r"Pipeline: VISUALIZE only (bin real $Y_{\mathrm{draft}}$ on ability | $\mathrm{poolq\_loo}$).",
        ],
        "footer": r"Assign $\rightarrow$ Score $\rightarrow$ Select in the real world; we read outcomes only.",
    },
    {
        "title": "Pass B — Generative: Congestion in Score Bends the Curve",
        "image": PASS_B / "PASS_B_generative_lambda_knockout_side_by_side.png",
        "bullets": [
            r"Synthetic league: ASSIGN $\rightarrow$ SCORE $\rightarrow$ SELECT $\rightarrow$ VISUALIZE ($539$ preset).",
            r"Left — $\lambda = 0$: score $S_i = A_i$ only (talent-only ranking). Roughly monotone.",
            r"Right — $\lambda > 0$: score $S_i = A_i - w \cdot L_C$ (viable-peer congestion in SCORE).",
            r"VISUALIZE on pool mean (team ability, includes self) — not $\mathrm{poolq\_loo}$.",
            r"Claim: roster pressure in the advancement rule bends the curve (qualitative POC).",
        ],
        "footer": r"Pass B knob: $\lambda$ / $w$ in SCORE. Does not vary assignment $\rho$.",
    },
    {
        "title": "Pass C — Generative: Assortativity Shapes the Sorted World",
        "image": PASS_C_RHO / "PASS_C_rho_ablation_selection_by_pool_mean.png",
        "bullets": [
            r"Same $539$ score as Pass B right arm — $S_i = A_i - w \cdot L_C$ held fixed.",
            r"One draw of talent; only ASSIGN changes (soft $\rho$ ladder + sort-and-chop).",
            r"$\rho$ controls how sharply players match team targets — roster sorting.",
            r"VISUALIZE on pool mean: assortativity moves the inverted-U readout.",
            r"Story: $\lambda$ puts roster pressure in score; $\rho$ delivers pressured environments.",
        ],
        "footer": r"Pass C knob: $\rho$ in ASSIGN. Score + top-$K$ fixed across arms.",
    },
]


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    """Insert image scaled to fit box; return bottom edge (top + height)."""
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return top + Inches(0.45)

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    return pic.top + pic.height


def _add_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title_top = Inches(0.28)
    title_h = Inches(0.52)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    ttf = title_box.text_frame
    populate_paragraph_with_latex(ttf.paragraphs[0], spec["title"], font_size=22)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].alignment = PP_ALIGN.LEFT

    img_top = title_top + title_h + Inches(0.12)
    footer_h = Inches(0.38)
    text_bottom_reserve = Inches(2.05)
    max_img_h = SLIDE_H - img_top - text_bottom_reserve - footer_h

    img_bottom = _add_picture_fitted(
        slide,
        spec["image"],
        MARGIN,
        img_top,
        CONTENT_W,
        max_img_h,
    )

    text_top = img_bottom + Inches(0.14)
    footer_top = SLIDE_H - MARGIN - footer_h

    mid = len(spec["bullets"]) // 2 + len(spec["bullets"]) % 2
    col_w = (CONTENT_W - Inches(0.25)) // 2
    left_col = slide.shapes.add_textbox(MARGIN, text_top, col_w, footer_top - text_top)
    fill_bullets_latex(left_col.text_frame, spec["bullets"][:mid], font_size=11)

    right_col = slide.shapes.add_textbox(
        MARGIN + col_w + Inches(0.25), text_top, col_w, footer_top - text_top
    )
    fill_bullets_latex(right_col.text_frame, spec["bullets"][mid:], font_size=11)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    ftf = foot.text_frame
    populate_paragraph_with_latex(ftf.paragraphs[0], spec["footer"], font_size=10)
    ftf.paragraphs[0].font.italic = True
    ftf.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    ensure_hero_dirs()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for spec in SLIDES:
        _add_slide(prs, spec)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (16:9 — figure top, wording below)")


if __name__ == "__main__":
    main()
