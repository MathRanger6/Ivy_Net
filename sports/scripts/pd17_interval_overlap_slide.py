"""Shared python-pptx layout for PD17 interval-overlap reference slides."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_TEXT_W = Inches(4.15)
FIG_W = CONTENT_W - LEFT_TEXT_W - COL_GAP


def load_meta(path: Path) -> dict:
    if path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height) -> None:
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return
    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def build_interval_overlap_slide(
    *,
    fig_path: Path,
    out_pptx: Path,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_h = Inches(2.35)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = body_top
    fig_h = footer_top - fig_top - Inches(0.04)
    fig_left = MARGIN + LEFT_TEXT_W + COL_GAP

    add_picture_fitted(slide, fig_path, fig_left, fig_top, FIG_W, fig_h)

    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, LEFT_TEXT_W, bullet_h)
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True
    fill_bullets_raw_latex(bullet_tf, bullets, font_size=HAND_BODY_PT)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], claim, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"Wrote {out_pptx}")


def build_figure_focus_slide(
    *,
    fig_path: Path,
    out_pptx: Path,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str = "",
) -> None:
    """Large centered figure — title/subtitle top, compact bullets + optional claim below."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.42) if claim else Inches(0.0)
    footer_top = SLIDE_H - MARGIN - footer_h
    bullet_h = Inches(1.05)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = sub_top + sub_h + Inches(0.1)
    fig_h = bullet_top - fig_top - Inches(0.08)

    add_picture_fitted(slide, fig_path, MARGIN, fig_top, CONTENT_W, fig_h)

    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, CONTENT_W, bullet_h)
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True
    fill_bullets_raw_latex(bullet_tf, bullets, font_size=HAND_BODY_PT)

    if claim:
        foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
        foot_tf = foot.text_frame
        foot_tf.word_wrap = True
        populate_paragraph_raw_latex(foot_tf.paragraphs[0], claim, font_size=HAND_CLAIM_PT, bold=True)
        foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"Wrote {out_pptx}")


def build_text_reference_slide(
    *,
    out_pptx: Path,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
    fig_path: Path | None = None,
) -> None:
    """Glossary / explainer layout — full-width bullets, optional right figure."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)

    if fig_path is not None and fig_path.is_file():
        bullet_w = LEFT_TEXT_W
        fig_left = MARGIN + LEFT_TEXT_W + COL_GAP
        fig_w = FIG_W
    else:
        bullet_w = CONTENT_W
        fig_left = None
        fig_w = None

    bullet_top = body_top
    bullet_h = footer_top - bullet_top - Inches(0.06)
    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, bullet_w, bullet_h)
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True
    fill_bullets_raw_latex(bullet_tf, bullets, font_size=HAND_BODY_PT)

    if fig_path is not None and fig_left is not None and fig_w is not None:
        fig_h = footer_top - body_top - Inches(0.04)
        add_picture_fitted(slide, fig_path, fig_left, body_top, fig_w, fig_h)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], claim, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"Wrote {out_pptx}")
