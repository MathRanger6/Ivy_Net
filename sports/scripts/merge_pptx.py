#!/usr/bin/env python3
"""Merge .pptx decks in order (append all slides).

Default (--python): re-embeds pictures; copies groups intact (do not flatten).

Recommended on Mac when PowerPoint is installed:
  --native   uses AppleScript + Microsoft PowerPoint (best fidelity)

Run (repo root):
  python sports/scripts/merge_pptx.py OUT.pptx part1.pptx part2.pptx ...
  python sports/scripts/merge_pptx.py --native OUT.pptx part1.pptx part2.pptx ...
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from copy import deepcopy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1]
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            sp_tree.remove(child)


def _copy_shape(shape, dest_slide) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        stream = BytesIO(shape.image.blob)
        dest_slide.shapes.add_picture(
            stream, shape.left, shape.top, shape.width, shape.height
        )
        return

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # Keep nested lines/freeforms together — flattening breaks Model slide.
        el = deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(el, "p:extLst")
        return

    el = deepcopy(shape.element)
    dest_slide.shapes._spTree.insert_element_before(el, "p:extLst")


def _replace_slide_content(source_slide, dest_slide) -> None:
    _clear_slide(dest_slide)
    for shape in source_slide.shapes:
        _copy_shape(shape, dest_slide)


def _copy_slide(source_slide, dest_prs: Presentation):
    layout = dest_prs.slide_layouts[6]
    dest_slide = dest_prs.slides.add_slide(layout)
    _replace_slide_content(source_slide, dest_slide)


def _append_slides_python(src_path: Path, dest_prs: Presentation) -> int:
    src = Presentation(str(src_path))
    for slide in src.slides:
        _copy_slide(slide, dest_prs)
    return len(src.slides)


def _merge_python(out_path: Path, parts: list[Path]) -> int:
    base = Presentation(str(parts[0]))
    dest = Presentation()
    dest.slide_width = base.slide_width
    dest.slide_height = base.slide_height

    total = 0
    for part in parts:
        added = _append_slides_python(part, dest)
        print(f"  + {added} slide(s) from {part.name} (python)")
        total += added

    out_path.parent.mkdir(parents=True, exist_ok=True)
    dest.save(str(out_path))
    return total


def _merge_native(out_path: Path, parts: list[Path]) -> None:
    script = Path(__file__).resolve().parents[2] / "scripts" / "merge_pptx_native.applescript"
    if not script.is_file():
        raise FileNotFoundError(script)
    cmd = ["osascript", str(script), str(out_path.resolve())]
    cmd.extend(str(p.resolve()) for p in parts)
    print("  Using Microsoft PowerPoint (AppleScript merge)...")
    subprocess.run(cmd, check=True)


def update_pass_slides(
    dest_path: Path,
    src_path: Path,
    *,
    dest_start: int = 0,
    src_count: int | None = None,
) -> int:
    """Replace slides in dest with matching slides from src; leave other dest slides untouched."""
    if not dest_path.is_file():
        raise FileNotFoundError(dest_path)
    if not src_path.is_file():
        raise FileNotFoundError(src_path)

    src = Presentation(str(src_path))
    dest = Presentation(str(dest_path))

    count = len(src.slides) if src_count is None else src_count
    if count <= 0:
        raise ValueError("src_count must be positive")
    if count > len(src.slides):
        raise ValueError(
            f"{src_path.name} has {len(src.slides)} slide(s); requested {count}"
        )
    if dest_start + count > len(dest.slides):
        raise ValueError(
            f"{dest_path.name} has {len(dest.slides)} slide(s); "
            f"cannot replace {count} starting at slide {dest_start + 1}"
        )

    for offset in range(count):
        _replace_slide_content(src.slides[offset], dest.slides[dest_start + offset])

    dest_path.parent.mkdir(parents=True, exist_ok=True)
    dest.save(str(dest_path))
    return count


def merge_pptx(
    out_path: Path,
    parts: list[Path],
    *,
    skip_missing: bool,
    use_native: bool,
) -> None:
    if not parts:
        raise SystemExit("merge_pptx: no input parts")

    resolved: list[Path] = []
    for part in parts:
        if part.is_file():
            resolved.append(part)
            continue
        if skip_missing:
            print(f"  skip missing: {part}")
            continue
        raise FileNotFoundError(part)

    if not resolved:
        raise SystemExit("merge_pptx: no input files found on disk")

    if use_native:
        _merge_native(out_path, resolved)
        print(f"Wrote {out_path} (native PowerPoint merge)")
        return

    total = _merge_python(out_path, resolved)
    print(f"Wrote {out_path} ({total} slides, python merge)")


def main() -> None:
    parser = argparse.ArgumentParser(description="Merge PowerPoint decks in order.")
    parser.add_argument("out_pptx", type=Path, help="Output .pptx path")
    parser.add_argument("parts", nargs="*", type=Path, help="Input decks in order")
    parser.add_argument(
        "--skip-missing",
        action="store_true",
        help="Skip input files that do not exist",
    )
    parser.add_argument(
        "--python",
        action="store_true",
        help="Force python-pptx merge (pictures re-embedded; groups copied intact)",
    )
    parser.add_argument(
        "--native",
        action="store_true",
        help="Use Microsoft PowerPoint via AppleScript (best on Mac)",
    )
    parser.add_argument(
        "--update-passes",
        action="store_true",
        help="Replace first N slides in OUT from PART1; leave remaining slides untouched",
    )
    parser.add_argument(
        "--pass-count",
        type=int,
        default=None,
        help="With --update-passes: how many slides to copy (default: all slides in PART1)",
    )
    args = parser.parse_args()

    if args.update_passes:
        if len(args.parts) != 1:
            raise SystemExit("--update-passes requires exactly one source deck (PASS_ABC)")
        replaced = update_pass_slides(
            args.out_pptx,
            args.parts[0],
            src_count=args.pass_count,
        )
        kept = len(Presentation(str(args.out_pptx)).slides) - replaced
        print(
            f"Updated {args.out_pptx.name}: replaced slide(s) 1–{replaced}; "
            f"{kept} slide(s) after that left unchanged"
        )
        return

    if not args.parts:
        raise SystemExit("merge_pptx: no input parts")

    use_native = args.native
    if not args.python and not args.native:
        # Default: native on macOS when not forced to python
        use_native = sys.platform == "darwin"

    merge_pptx(
        args.out_pptx,
        args.parts,
        skip_missing=args.skip_missing,
        use_native=use_native,
    )


if __name__ == "__main__":
    main()
