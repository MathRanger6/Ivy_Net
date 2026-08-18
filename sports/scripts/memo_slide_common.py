"""Memo-style AUTO slides — large type, conversational narrative (PD20–22 item 15).

Unlike HAND diagnostic decks (small sidebar + figure), memo slides use labeled
story blocks: Why this came up → What we ran → What showed up → What you can say.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from gallery_mathtext import RAW_LATEX_FONT, populate_paragraph_raw_latex

# Memo typography — medium-large, readable at a glance in a meeting room.
MEMO_TITLE_PT = 28
MEMO_QUESTION_PT = 15
MEMO_LABEL_PT = 14
MEMO_BODY_PT = 17
MEMO_SO_WHAT_PT = 17
MEMO_SECTION_SPACE_AFTER_PT = 14
MEMO_LABEL_BODY_GAP_PT = 2

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.48)
CONTENT_W = SLIDE_W - 2 * MARGIN

MEMO_HAND_CUE_PT = 13
MEMO_HAND_CUE_COLOR = (0x55, 0x55, 0x55)  # dark gray footer cue

NARRATIVE_SECTIONS = (
    ("Why this came up", "why"),
    ("What we ran", "what"),
    ("What showed up", "saw"),
    ("What you can say", "so_what"),
)


def _add_label_paragraph(text_frame, label: str, *, bold: bool = True) -> None:
    para = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
    if text_frame.text and para.text:
        para = text_frame.add_paragraph()
    populate_paragraph_raw_latex(para, label, font_size=MEMO_LABEL_PT, bold=bold)
    para.space_after = Pt(MEMO_LABEL_BODY_GAP_PT)
    para.space_before = Pt(0)
    para.line_spacing = 1.15


def _add_body_paragraph(text_frame, text: str, *, font_size: int = MEMO_BODY_PT, bold: bool = False) -> None:
    para = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
    if text_frame.text and para.text:
        para = text_frame.add_paragraph()
    populate_paragraph_raw_latex(para, text, font_size=font_size, bold=bold)
    para.space_after = Pt(MEMO_SECTION_SPACE_AFTER_PT)
    para.space_before = Pt(0)
    para.line_spacing = 1.25


def fill_narrative_sections(
    text_frame,
    *,
    why: str,
    what: str,
    saw: str,
    so_what: str,
) -> None:
    """Four labeled prose blocks — no bullet glyphs."""
    text_frame.clear()
    text_frame.word_wrap = True
    blocks = {
        "why": why,
        "what": what,
        "saw": saw,
        "so_what": so_what,
    }
    for label, key in NARRATIVE_SECTIONS:
        _add_label_paragraph(text_frame, f"{label}:")
        body_size = MEMO_SO_WHAT_PT if key == "so_what" else MEMO_BODY_PT
        body_bold = key == "so_what"
        _add_body_paragraph(text_frame, blocks[key], font_size=body_size, bold=body_bold)


def append_bridge_slide(
    prs: Presentation,
    *,
    title: str,
    lead: str,
    blocks: list[tuple[str, str]],
) -> None:
    """Big-picture / transition slide — title, one-line lead, labeled prose blocks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.26)
    title_h = Inches(0.56)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=MEMO_TITLE_PT,
        bold=True,
    )

    lead_top = title_top + title_h + Inches(0.06)
    lead_h = Inches(0.48)
    lead_box = slide.shapes.add_textbox(MARGIN, lead_top, CONTENT_W, lead_h)
    lead_tf = lead_box.text_frame
    lead_tf.word_wrap = True
    populate_paragraph_raw_latex(
        lead_tf.paragraphs[0],
        lead,
        font_size=MEMO_QUESTION_PT,
        bold=False,
    )

    body_top = lead_top + lead_h + Inches(0.08)
    body_h = SLIDE_H - MARGIN - body_top
    body_box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, body_h)
    tf = body_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for label, text in blocks:
        _add_label_paragraph(tf, f"{label}:")
        _add_body_paragraph(tf, text)


def append_hand_companion_slide(
    prs: Presentation,
    *,
    hand_slide: int,
    act: str,
    title: str,
    question: str,
    why: str,
    what: str,
    saw: str,
    so_what: str,
) -> None:
    """Narrative companion for one CHAR_PD20_HAND slide — footer cues paste-after."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    cue_h = Inches(0.38)
    cue_top = SLIDE_H - MARGIN - cue_h
    body_bottom = cue_top - Inches(0.06)

    title_top = Inches(0.22)
    title_h = Inches(0.52)
    meta_h = Inches(0.28)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=MEMO_TITLE_PT - 2,
        bold=True,
    )

    meta_top = title_top + title_h
    meta_box = slide.shapes.add_textbox(MARGIN, meta_top, CONTENT_W, meta_h)
    populate_paragraph_raw_latex(
        meta_box.text_frame.paragraphs[0],
        rf"{act} \text{{---}} narrative for HAND slide {hand_slide}",
        font_size=MEMO_HAND_CUE_PT,
        bold=False,
    )

    q_top = meta_top + meta_h + Inches(0.02)
    q_h = Inches(0.44)
    q_box = slide.shapes.add_textbox(MARGIN, q_top, CONTENT_W, q_h)
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    populate_paragraph_raw_latex(
        q_tf.paragraphs[0],
        question,
        font_size=MEMO_QUESTION_PT - 1,
        bold=False,
    )

    body_top = q_top + q_h + Inches(0.04)
    body_h = body_bottom - body_top
    body_box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, body_h)
    fill_narrative_sections(
        body_box.text_frame,
        why=why,
        what=what,
        saw=saw,
        so_what=so_what,
    )

    cue_box = slide.shapes.add_textbox(MARGIN, cue_top, CONTENT_W, cue_h)
    cue_para = cue_box.text_frame.paragraphs[0]
    cue_para.alignment = PP_ALIGN.CENTER
    populate_paragraph_raw_latex(
        cue_para,
        rf"$\downarrow$ Paste \texttt{{CHAR\_PD20\_HAND}} slide {hand_slide} immediately after this slide",
        font_size=MEMO_HAND_CUE_PT,
        bold=True,
    )


def append_narrative_memo_slide(
    prs: Presentation,
    *,
    title: str,
    question: str,
    why: str,
    what: str,
    saw: str,
    so_what: str,
) -> None:
    """One conversational memo slide answering a single question."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.26)
    title_h = Inches(0.56)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=MEMO_TITLE_PT,
        bold=True,
    )

    q_top = title_top + title_h + Inches(0.04)
    q_h = Inches(0.52)
    q_box = slide.shapes.add_textbox(MARGIN, q_top, CONTENT_W, q_h)
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    populate_paragraph_raw_latex(
        q_tf.paragraphs[0],
        question,
        font_size=MEMO_QUESTION_PT,
        bold=False,
    )

    body_top = q_top + q_h + Inches(0.06)
    body_h = SLIDE_H - MARGIN - body_top
    body_box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, body_h)
    fill_narrative_sections(
        body_box.text_frame,
        why=why,
        what=what,
        saw=saw,
        so_what=so_what,
    )


def append_memo_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str = "",
) -> None:
    """Legacy bullet memo slide — prefer append_narrative_memo_slide for item 15."""
    from gallery_mathtext import HAND_BULLET_LEAD, _set_char_bullet

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_top = Inches(0.28)
    title_h = Inches(0.62)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=MEMO_TITLE_PT,
        bold=True,
    )
    sub_top = title_top + title_h + Inches(0.06)
    sub_h = Inches(0.34)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=MEMO_QUESTION_PT,
    )
    footer_h = Inches(0.52) if claim else Inches(0.0)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.14)
    bullet_box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, footer_top - body_top - Inches(0.08))
    tf = bullet_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        populate_paragraph_raw_latex(para, HAND_BULLET_LEAD + line.strip(), font_size=MEMO_BODY_PT)
        para.line_spacing = 1.35
        para.space_after = Pt(10)
        _set_char_bullet(para)
    if claim:
        foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
        populate_paragraph_raw_latex(
            foot.text_frame.paragraphs[0],
            claim,
            font_size=MEMO_BODY_PT,
            bold=True,
        )


def new_memo_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def save_memo_deck(prs: Presentation, out_pptx: Path) -> None:
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"Wrote {out_pptx}")
