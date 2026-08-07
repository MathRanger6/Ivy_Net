#!/usr/bin/env python3
"""Build two-slide ρ characterization deck (with / without sort-and-chop on figure).

Run (repo root):
  python sports/scripts/build_rho_characterization_slide.py

Regenerates Pass C PNG variants, then writes:
  HEROs_and_PASSes/slides/auto/CHAR_rho_characterization_AUTO.pptx
    slide 1 — sort-and-chop visible
    slide 2 — sort-and-chop suppressed (soft ρ arms only)
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from gallery_knobs import RHO_HIGH, RHO_LOW, RHO_MODERATE, RHO_VERY_HIGH, HERO_SEED
from hero_gallery_paths import AUTO_RHO_DECK, PASS_C_RHO, ensure_hero_dirs

META = PASS_C_RHO / "PASS_C_rho_ablation_meta.json"
OUT_PPTX = AUTO_RHO_DECK
PASS_C_SCRIPT = SCRIPTS / "pass_c_rho_ablation_bundle.py"

FIG_WITH = PASS_C_RHO / "PASS_C_rho_ablation_selection_by_pool_mean_with_sortchop.png"
FIG_WITHOUT = PASS_C_RHO / "PASS_C_rho_ablation_selection_by_pool_mean.png"

from gallery_mathtext import add_plain_latex_block, fill_bullets_latex, populate_paragraph_with_latex

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_W = Inches(4.55)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

CLAIM = (
    "Claim: at fixed score and top-$K$, $\\rho$ in assignment matters — "
    "low-$\\rho$ mixing yields a weak monotone tilt (no inverted-$U$); "
    "high-$\\rho$ assortativity delivers the peer-pressure inverted-$U$ readout."
)

RHO_INTUITION = (
    "Why $\\rho$ matters here: assignment assortativity sets how spread out team "
    "$L_C$ is across the league — at $\\rho \\approx 0$ (arm $0.001$) rosters mix "
    "and $L_C$ is nearly uniform; selection still tilts up pool-mean bins because "
    "global top-$K$ favors better players on stronger teams (see doc 08). "
    "At high $\\rho$ weak teams pile up near $L_C \\approx 0$ while elite rosters "
    "face real congestion — the inverted-$U$ emerges."
)

SLIDE_VARIANTS = [
    {
        "subtitle": "sort-and-chop shown",
        "figure": FIG_WITH,
        "show_sortchop_on_figure": True,
    },
    {
        "subtitle": "sort-and-chop suppressed",
        "figure": FIG_WITHOUT,
        "show_sortchop_on_figure": False,
    },
]


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {"selection": {"w": 0.55}, "assignment_sigma": 0.65, "preset": "539"}


def _regenerate_figures() -> None:
    print("Regenerating Pass C figures (soft-only + with sort-and-chop) ...")
    subprocess.run(
        [sys.executable, str(PASS_C_SCRIPT)],
        cwd=str(REPO),
        check=True,
    )


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2


def _add_params_column(
    slide,
    meta: dict,
    *,
    left,
    top,
    width,
    height,
    show_sortchop_on_figure: bool,
) -> None:
    w = float(meta.get("selection", {}).get("w", 0.55))
    sigma = float(meta.get("assignment_sigma", 0.65))
    preset = meta.get("preset", "539")

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    head = tf.paragraphs[0]
    populate_paragraph_with_latex(
        head,
        rf"Characterize $\rho$ — one-at-a-time at {preset} baseline "
        r"(score + top-$K$ fixed)",
        font_size=11,
    )
    head.font.bold = True
    head.space_after = 6

    assign_latex = (
        r"\pi_{ij} \propto \exp\left(-\rho \cdot \frac{(A_i-T_{j^*})^2}{2\sigma^2}\right)"
    )
    add_plain_latex_block(tf, assign_latex, font_size=10, label="ASSIGN (knob ρ) — LaTeX:")

    bullets = [
        rf"$\sigma={sigma:g}$ fixed",
        rf"arms: $\rho \in {{{RHO_LOW:g}, {RHO_MODERATE:g}, {RHO_HIGH:g}, {RHO_VERY_HIGH:g}}}$"
        + (" + sort-and-chop" if show_sortchop_on_figure else ""),
        rf"$S_i = A_i - {w:g} \cdot L_C$ (crowding in score held constant across arms)",
        r"top-$K$ by $S_i$",
        r"VISUALIZE = mean $Y_{\mathrm{selected}}$ vs pool mean ($16$ bins)",
    ]
    for line in bullets:
        para = tf.add_paragraph()
        para.text = ""
        populate_paragraph_with_latex(para, line, font_size=10)
        para.space_after = 2
        para.level = 0


def _add_oat_notes(
    slide,
    *,
    left,
    top,
    width,
    height,
    show_sortchop_on_figure: bool,
) -> None:
    bullets = [
        rf"OAT: one draw of $A_i$, $T_{{j^*}}$ (seed ${HERO_SEED}$); only ASSIGN changes across arms.",
        r"Low $\rho$: rosters mix — $L_C$ nearly uniform; curve tilts up weakly (no inverted-$U$).",
        r"High $\rho$: assortative matching — flat bottom bins, interior peak, top-bin dip.",
    ]
    if show_sortchop_on_figure:
        bullets.append(
            r"Sort-and-chop benchmark: hard rank match (same score rule as soft arms)."
        )
    else:
        bullets.append(
            r"Sort-and-chop arm omitted from plot (still in CSV bundle; toggle restores it)."
        )

    box = slide.shapes.add_textbox(left, top, width, height)
    fill_bullets_latex(box.text_frame, bullets, font_size=9)


def _add_slide(prs: Presentation, meta: dict, variant: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.28)
    title_h = Inches(0.46)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_with_latex(
        title_box.text_frame.paragraphs[0],
        rf"Phase B — Characterize $\rho$ (assignment assortativity) — {variant['subtitle']}",
        font_size=20,
    )
    title_box.text_frame.paragraphs[0].font.bold = True

    body_top = title_top + title_h + Inches(0.1)
    intuition_h = Inches(0.38)
    footer_h = Inches(0.52)
    footer_top = SLIDE_H - MARGIN - footer_h
    intuition_top = footer_top - intuition_h - Inches(0.04)
    notes_h = Inches(1.35)
    params_h = intuition_top - body_top - notes_h - Inches(0.08)

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_params_column(
        slide,
        meta,
        left=left_x,
        top=body_top,
        width=LEFT_W,
        height=params_h,
        show_sortchop_on_figure=variant["show_sortchop_on_figure"],
    )

    _add_picture_fitted(
        slide,
        variant["figure"],
        right_x,
        body_top,
        RIGHT_W,
        params_h,
    )

    _add_oat_notes(
        slide,
        left=left_x,
        top=intuition_top - notes_h - Inches(0.06),
        width=LEFT_W + RIGHT_W + COL_GAP,
        height=notes_h,
        show_sortchop_on_figure=variant["show_sortchop_on_figure"],
    )

    note_box = slide.shapes.add_textbox(MARGIN, intuition_top, CONTENT_W, intuition_h)
    populate_paragraph_with_latex(
        note_box.text_frame.paragraphs[0],
        RHO_INTUITION,
        font_size=10,
    )
    note_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_with_latex(foot.text_frame.paragraphs[0], CLAIM, font_size=11)
    foot.text_frame.paragraphs[0].font.bold = True
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    ensure_hero_dirs()
    _regenerate_figures()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for variant in SLIDE_VARIANTS:
        _add_slide(prs, meta, variant)

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(SLIDE_VARIANTS)} slides)")


if __name__ == "__main__":
    main()
