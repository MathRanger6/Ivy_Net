#!/usr/bin/env python3
"""Build PD17 sim rho vs empirical coverage overlay slide.

Run (repo root):
  python sports/scripts/build_empirical_rho_coverage_slide.py
  python sports/scripts/build_empirical_rho_coverage_slide.py --slides-only
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_EMPIRICAL_RHO_COVERAGE_DECK, EMPIRICAL_PD17, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "empirical_rho_coverage_overlay.py"
FIG = EMPIRICAL_PD17 / "EMPIRICAL_rho_coverage_overlay.png"
META = EMPIRICAL_PD17 / "EMPIRICAL_rho_coverage_overlay_meta.json"
OUT_PPTX = AUTO_EMPIRICAL_RHO_COVERAGE_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_TEXT_W = Inches(4.15)
FIG_W = CONTENT_W - LEFT_TEXT_W - COL_GAP

CLAIM = (
    "Claim (PD17): Real rosters sit at high coverage peak — sim \\rho must be tuned "
    r"so soft assign matches that overlap, not sort-and-chop."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figure() -> None:
    print("Regenerating empirical vs sim rho coverage overlay ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def _readout_bullets(meta: dict) -> list[str]:
    emp = meta.get("empirical", {})
    legacy = meta.get("sim_legacy_arms", [])
    sweep = meta.get("sim_sweep_arms", [])
    chop = meta.get("sim_sort_chop_peak")
    seed = meta.get("seed")
    mode = meta.get("sim_panel_mode", "three_panel")

    bullets = [
        r"Three panels: empirical (left) · four \rho arms (center) · \rho=1\to32 sweep (right).",
        r"Same A_i / T_{j^*} draw across sim panels; seed fixed.",
    ]
    if emp:
        bullets.append(
            rf"Empirical max coverage={emp.get('coverage_max', 0):,}; "
            rf"sort-and-chop max={emp.get('coverage_disjoint_max', 0)}."
        )
    if legacy and sweep:
        rho_four = r"\rho\in\{0.001, 1, 8, 32\}"
        bullets.append(
            rf"Center: {rho_four}; right sweep peaks "
            rf"\rho=1 \rightarrow {sweep[0].get('coverage_peak', 0)}, "
            rf"\rho=32 \rightarrow {sweep[-1].get('coverage_peak', 0)}."
        )
    if chop is not None:
        bullets.append(rf"Sim sort-and-chop peak={chop} on both sim panels.")
    if seed is not None:
        bullets.append(rf"PPM z vs [0,1] ability — compare shape, seed {seed}.")
    if mode == "two_panel":
        bullets[0] = r"Two-panel mode: empirical | \rho sweep only (no center four-arm)."
    return bullets


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--slides-only", action="store_true")
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figure()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        r"PD17 — Empirical vs sim interval overlap (\rho calibration target)",
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    preset = meta.get("preset", "539")
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        rf"Assign block capstone · empirical MBB vs {preset} sim · same coverage diagnostic",
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_h = Inches(2.45)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = body_top
    fig_h = footer_top - fig_top - Inches(0.04)
    fig_left = MARGIN + LEFT_TEXT_W + COL_GAP

    _add_picture_fitted(slide, FIG, fig_left, fig_top, FIG_W, fig_h)

    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, LEFT_TEXT_W, bullet_h)
    fill_bullets_raw_latex(bullet_box.text_frame, _readout_bullets(meta), font_size=HAND_BODY_PT)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — rho coverage overlay)")


if __name__ == "__main__":
    main()
