#!/usr/bin/env python3
"""Build one-slide λ characterization deck (Phase B template).

Run (repo root):
  python sports/scripts/build_lambda_characterization_slide.py

Regenerates Pass B λ ablation figure, then writes:
  HEROs_and_PASSes/slides/auto/CHAR_lambda_characterization_AUTO.pptx
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from hero_gallery_paths import AUTO_LAMBDA_DECK, PASS_B, ensure_hero_dirs

META = PASS_B / "PASS_B_lambda_ablation_meta.json"
OUT_PPTX = AUTO_LAMBDA_DECK
PASS_B_SCRIPT = SCRIPTS / "pass_b_lambda_ablation_bundle.py"
FIG = PASS_B / "PASS_B_lambda_ablation_selection_by_pool_mean.png"

from gallery_knobs import LAMBDA_HIGH, LAMBDA_LOW, LAMBDA_MODERATE
from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_W = Inches(4.55)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

CLAIM = (
    "Claim: at fixed $\\rho$ and top-$K$, $\\lambda$ in score matters — "
    "$\\lambda=0$ yields a roughly monotone talent ladder; "
    "$\\lambda>0$ puts roster pressure in the advancement rule and bends the curve."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {
        "preset": "539",
        "assignment": {"rho_fixed": 8.0, "assignment_sigma": 0.65},
    }


def _regenerate_figure() -> None:
    print("Regenerating Pass B λ ablation figure ...")
    subprocess.run([sys.executable, str(PASS_B_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2


def _add_params_column(slide, meta: dict, *, left, top, width, height) -> None:
    rho = float(meta.get("assignment", {}).get("rho_fixed", 8.0))
    sigma = float(meta.get("assignment", {}).get("assignment_sigma", 0.65))
    preset = meta.get("preset", "539")

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    head = tf.paragraphs[0]
    populate_paragraph_with_latex(
        head,
        rf"Characterize $\lambda$ — one-at-a-time at {preset} baseline "
        r"(assign + top-$K$ fixed)",
        font_size=11,
    )
    head.font.bold = True
    head.space_after = 6

    eq = tf.add_paragraph()
    eq.text = ""
    eq.space_after = 4
    populate_paragraph_with_latex(
        eq,
        r"$S_i = A_i - \lambda \cdot L_C$",
        font_size=12,
    )

    bullets = [
        rf"$\rho={rho:g}$ fixed (soft assign; same roster all arms)",
        rf"$\sigma={sigma:g}$ fixed",
        rf"arms: $\lambda \in {{0, {LAMBDA_LOW:g}, {LAMBDA_MODERATE:g}, {LAMBDA_HIGH:g}}}$",
        r"top-$K$ by $S_i$",
        r"VISUALIZE = mean $Y_{\mathrm{selected}}$ vs pool mean ($16$ bins)",
    ]
    for line in bullets:
        para = tf.add_paragraph()
        para.text = ""
        populate_paragraph_with_latex(para, line, font_size=10)
        para.space_after = 2
        para.level = 0


def _add_oat_notes(slide, *, left, top, width, height) -> None:
    bullets = [
        r"OAT: one draw of $A_i$, $T_{j^*}$ + one assignment (seed $42$); only SCORE changes.",
        r"$\lambda=0$: $S_i = A_i$ only — selection vs pool mean stays roughly monotone.",
        r"$\lambda>0$: congestion in score — inverted-$U$ / roster-pressure readout emerges.",
        r"Pairs with Pass C: $\lambda$ bends score; $\rho$ shapes the sorted roster world.",
    ]
    box = slide.shapes.add_textbox(left, top, width, height)
    fill_bullets_latex(box.text_frame, bullets, font_size=9)


def _add_slide(prs: Presentation, meta: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.28)
    title_h = Inches(0.46)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_with_latex(
        title_box.text_frame.paragraphs[0],
        r"Phase B — Characterize $\lambda$ (weight on $L_C$ in score)",
        font_size=20,
    )
    title_box.text_frame.paragraphs[0].font.bold = True

    body_top = title_top + title_h + Inches(0.1)
    footer_h = Inches(0.52)
    footer_top = SLIDE_H - MARGIN - footer_h
    notes_h = Inches(1.35)
    params_h = footer_top - body_top - notes_h - Inches(0.08)

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_params_column(
        slide, meta, left=left_x, top=body_top, width=LEFT_W, height=params_h
    )
    _add_picture_fitted(slide, FIG, right_x, body_top, RIGHT_W, params_h)
    _add_oat_notes(
        slide,
        left=left_x,
        top=footer_top - notes_h - Inches(0.06),
        width=LEFT_W + RIGHT_W + COL_GAP,
        height=notes_h,
    )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_with_latex(foot.text_frame.paragraphs[0], CLAIM, font_size=11)
    foot.text_frame.paragraphs[0].font.bold = True
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    ensure_hero_dirs()
    _regenerate_figure()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _add_slide(prs, meta)
    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
