#!/usr/bin/env python3
"""Build PowerPoint from reigning hero star-sweep manifest.

Each slide: title, Δ vs reigning lock (EW16 · 09–21), prose diff, HERO PNG, CLI, shape.

Run (repo root, after reigning_hero_star_sweep.py):
  python sports/scripts/build_reigning_hero_star_sweep_slides.py
  python sports/scripts/build_reigning_hero_star_sweep_slides.py --manifest path/to/manifest.json
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from build_hero_permutation_slides import (  # noqa: E402
    CONTENT_W,
    MARGIN,
    SLIDE_H,
    SLIDE_W,
    _add_mono_box,
    _add_picture,
    _add_title,
    _run_counts,
    _shape_line,
)
from hero_gallery_paths import REIGNING_HERO_STAR_SWEEPS  # noqa: E402
from reigning_hero_star_diff import (  # noqa: E402
    REIGNING_LOCK_TAG,
    SEASON_WINDOW_LABELS,
    diff_from_reigning_lock,
)

DEFAULT_MANIFEST = REIGNING_HERO_STAR_SWEEPS / "manifest.json"
DEFAULT_OUT = REIGNING_HERO_STAR_SWEEPS / "HERO_star_sweep_slides_AUTO.pptx"
LOCK_PNG = (
    REPO
    / "3-Master_Plan/re_entry/HEROs_and_PASSes/population_sandbox/hero"
    / "HERO_ew16_allt_min20_mg10_09_21_last_ps_perm_loo_ever_lastps_ew16.png"
)


def _add_prose_box(slide, top: float, lines: list[str]) -> float:
    h = Inches(0.72)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(11)
        p.space_after = Pt(2)
    return top + h + Inches(0.06)


def _add_lock_slide(prs: Presentation, lock: dict, *, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lock_bins = int(lock.get("n_bins", 16))
    lock_win = str(lock.get("season_window", "09_21"))
    win_label = SEASON_WINDOW_LABELS.get(lock_win, lock_win)
    y = _add_title(
        slide,
        f"Slide {slide_num} — Reigning hero (reference lock)",
        subtitle=f"Slide 12 · tag `{REIGNING_LOCK_TAG}` · EW{lock_bins} · {win_label}",
    )
    explain = [
        "Named hero from 09–21 permutation deck (HAND matrix slide 12).",
        "Star sweep varies only EW bin count and season window; everything else matches this lock.",
        "Sweep grid omits EW16 — nearest neighbors are EW12 and EW20 at 09–21.",
        "β₂ ≈ +0.00172 (flat / not concave) on 2009–2021.",
    ]
    y = _add_prose_box(slide, y, explain)
    img_bottom = _add_picture(
        slide,
        LOCK_PNG,
        MARGIN,
        y,
        CONTENT_W,
        SLIDE_H - y - Inches(0.5) - MARGIN,
    )
    if not LOCK_PNG.is_file():
        note = slide.shapes.add_textbox(MARGIN, img_bottom + Inches(0.05), CONTENT_W, Inches(0.35))
        note.text_frame.text = (
            f"(Lock PNG not in repo — regenerate: pass_a with --output-tag {REIGNING_LOCK_TAG})"
        )
        note.text_frame.paragraphs[0].font.size = Pt(9)
        note.text_frame.paragraphs[0].font.italic = True


def _add_run_slide(
    prs: Presentation,
    run: dict,
    lock: dict,
    *,
    slide_num: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tag = run.get("output_tag", "")
    label = run.get("label", tag)
    diff = diff_from_reigning_lock(run, lock)
    run["diff_from_reigning_lock"] = diff

    y = _add_title(slide, f"Slide {slide_num} — {label}", subtitle=diff["subtitle"])
    y = _add_prose_box(slide, y, diff["prose_lines"])

    cmd_h = Inches(1.35)
    shape_h = Inches(0.35)
    footer_reserve = cmd_h + shape_h + Inches(0.18)
    max_img_h = SLIDE_H - y - footer_reserve - MARGIN

    png_rel = run.get("hero_png")
    png_path = REPO / png_rel if png_rel else Path()
    img_bottom = _add_picture(slide, png_path, MARGIN, y, CONTENT_W, max_img_h)

    cmd_top = img_bottom + Inches(0.08)
    cmd_text = run.get("command") or run.get("command_one_line") or "(no command recorded)"
    _add_mono_box(slide, MARGIN, cmd_top, CONTENT_W, cmd_h, cmd_text, font_size=8)

    shape_top = cmd_top + cmd_h + Inches(0.04)
    held = "Held vs lock: " + " · ".join(diff["held_constant"][:4]) + " · …"
    shape_box = slide.shapes.add_textbox(MARGIN, shape_top, CONTENT_W, shape_h)
    tf = shape_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _shape_line(run.get("shape") or {})
    p.font.size = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = held
    p2.font.size = Pt(9)
    p2.font.italic = True


def _add_intro_slide(prs: Presentation, manifest: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lock = manifest.get("reigning_lock") or {}
    lock_bins = int(lock.get("n_bins", 16))
    lock_win = str(lock.get("season_window", "09_21"))
    n_runs, n_planned, complete = _run_counts(manifest)
    status = f"{n_runs} run slides" if complete else f"{n_runs}/{n_planned} runs (incomplete)"
    _add_title(
        slide,
        "Slide 1 — Reigning hero star sweep",
        subtitle=f"{status} · {manifest.get('date', '')} · vs lock EW{lock_bins} · {lock_win}",
    )
    bins = ", ".join(str(b) for b in (manifest.get("n_bins_grid") or []))
    wins = ", ".join(str(w) for w in (manifest.get("season_windows") or []))
    body_top = Inches(1.05)
    box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        "Purpose: sensitivity of named reigning hero to EW ventile count and season window.",
        "",
        f"Reigning lock (slide 12): EW{lock_bins} · {SEASON_WINDOW_LABELS.get(lock_win, lock_win)} · "
        "poolq_LOO · ever · last-ps · ALLT · equal_width · min20 · mg10 · winsor 1–99.",
        "",
        "Varied in this sweep:",
        f"  • n_bins ∈ {{{bins}}}  (lock EW16 not in grid — use EW12/EW20 as neighbors)",
        f"  • season_window ∈ {{{wins}}}",
        "",
        "Each candidate slide explains Δ vs the lock; β₂ comparisons apply within the same era.",
        "Regenerate runs: python sports/scripts/reigning_hero_star_sweep.py",
        "Rebuild deck: python sports/scripts/build_reigning_hero_star_sweep_slides.py",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)


def _write_summary_md(manifest: dict, out_path: Path) -> None:
    lock = manifest.get("reigning_lock") or {}
    rows = ["# Reigning hero star sweep — Δ vs lock", ""]
    rows.append(
        f"Lock: EW{lock.get('n_bins', 16)} · {lock.get('season_window', '09_21')} · "
        f"`{REIGNING_LOCK_TAG}`"
    )
    rows.append("")
    rows.append("| Slide | Tag | n_bins | Season | β₂ | Diff vs lock |")
    rows.append("|-------|-----|--------|--------|-----|--------------|")
    for run in manifest.get("runs") or []:
        diff = run.get("diff_from_reigning_lock") or diff_from_reigning_lock(run, lock)
        shape = run.get("shape") or {}
        b2 = shape.get("beta_sq")
        b2s = f"{b2:+.5g}" if b2 is not None else "—"
        terse = "; ".join(diff.get("terse") or []) or "—"
        rows.append(
            f"| {run.get('slide_index', '?')} | `{run.get('output_tag', '')}` | "
            f"{run.get('n_bins', '?')} | {run.get('season_window', '')} | {b2s} | {terse} |"
        )
    out_path.write_text("\n".join(rows) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build reigning hero star-sweep PowerPoint.")
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    parser.add_argument(
        "--no-lock-slide",
        action="store_true",
        help="Skip slide 2 lock reference (intro + 20 runs only).",
    )
    parser.add_argument(
        "--write-manifest",
        action="store_true",
        default=True,
        help="Write slide_index and diff_from_reigning_lock back to manifest (default: on).",
    )
    parser.add_argument(
        "--no-write-manifest",
        action="store_false",
        dest="write_manifest",
        help="Do not rewrite manifest.json after build.",
    )
    args = parser.parse_args()

    manifest_path = args.manifest if args.manifest.is_absolute() else REPO / args.manifest
    if not manifest_path.is_file():
        raise SystemExit(
            f"Missing manifest {manifest_path}\n"
            "Run: python sports/scripts/reigning_hero_star_sweep.py"
        )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    runs = manifest.get("runs") or []
    if not runs:
        raise SystemExit("Manifest has no runs — run star sweep first (not --dry-run only).")

    lock = manifest.get("reigning_lock") or {
        "n_bins": 16,
        "season_window": "09_21",
        "roster_x": "poolq_loo",
        "y_draft_mode": "ever",
        "panel_rows": "last-ps",
        "poolq_binning": "equal_width",
        "dft": False,
    }

    out_path = args.out if args.out.is_absolute() else REPO / args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _add_intro_slide(prs, manifest)

    if not args.no_lock_slide:
        _add_lock_slide(prs, lock, slide_num=len(prs.slides) + 1)

    for run in runs:
        slide_num = len(prs.slides) + 1
        _add_run_slide(prs, run, lock, slide_num=slide_num)
        run["slide_index"] = slide_num

    prs.save(str(out_path))

    summary_md = out_path.parent / "HERO_star_sweep_slides_summary.md"
    manifest["runs"] = runs
    _write_summary_md(manifest, summary_md)

    if args.write_manifest:
        manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
        print(f"Updated {manifest_path.relative_to(REPO)}", flush=True)

    n_content = len(runs) + (0 if args.no_lock_slide else 1)
    print(
        f"Wrote {out_path.relative_to(REPO)} · {len(prs.slides)} slides "
        f"(intro + {n_content} content)",
        flush=True,
    )
    print(f"Wrote {summary_md.relative_to(REPO)}", flush=True)


if __name__ == "__main__":
    main()
