#!/usr/bin/env python3
"""Build PowerPoint from hero_permutation_slides/manifest.json.

Layout (matches Charles manual deck):
  - Title line (setting label + diff from FIXED HERO)
  - HERO PNG (large)
  - Command text box (multi-line CLI)
  - Shape readout (bin0 / peak / LPM β₂)

Run (repo root, after hero_permutation_sweep.py):
  python sports/scripts/build_hero_permutation_slides.py
  python sports/scripts/build_hero_permutation_slides.py --manifest path/to/manifest.json
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

SANDBOX = REPO / "3-Master_Plan/re_entry/HEROs_and_PASSes/population_sandbox"
DEFAULT_MANIFEST = SANDBOX / "hero_permutation_slides/manifest.json"
DEFAULT_OUT = SANDBOX / "hero_permutation_slides/HERO_permutation_slides_AUTO.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN


def _add_title(slide, text: str, *, subtitle: str | None = None) -> float:
    top = Inches(0.22)
    h = Inches(0.55 if subtitle else 0.48)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(20)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.italic = True
    return top + h + Inches(0.08)


def _add_picture(slide, img_path: Path, left, top, max_w, max_h) -> float:
    if not img_path.is_file():
        tb = slide.shapes.add_textbox(left, top, max_w, Inches(0.4))
        tb.text_frame.text = f"[Missing: {img_path.name}]"
        return top + Inches(0.45)
    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_w)
    if pic.height > max_h:
        scale = max_h / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_w - pic.width) // 2
    return pic.top + pic.height


def _add_mono_box(slide, left, top, width, height, text: str, *, font_size: int = 9) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Menlo"
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.LEFT


def _shape_line(shape: dict) -> str:
    if not shape:
        return "Shape: (no CSV summary — re-run sweep)"
    peak_v = shape.get("peak_vent")
    b0 = shape.get("bin0_rate_pct")
    peak = shape.get("peak_rate_pct")
    last = shape.get("last_bin_rate_pct")
    b2 = shape.get("beta_sq")
    flag = " ⚠ bin1 LOWEST ventile is PEAK" if shape.get("bin0_is_peak") else ""
    b2s = f" · LPM β₂={b2:+.4g}" if b2 is not None else ""
    return (
        f"Shape: bin1={b0}% · peak=ventile {int(peak_v)+1 if peak_v is not None else '?'} @ {peak}% "
        f"· last={last}%{b2s}{flag}"
    )


def _add_slide(prs: Presentation, run: dict, *, baseline_tag: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    idx = run.get("slide_index", "?")
    tag = run.get("output_tag", "")
    label = run.get("label", tag)
    diffs = run.get("diff_from_baseline") or []
    if tag == baseline_tag or not diffs:
        subtitle = "Baseline — canonical deck HERO"
    else:
        subtitle = "Δ vs FIXED HERO: " + "; ".join(diffs)

    y_after_title = _add_title(slide, f"Slide {idx} — {label}", subtitle=subtitle)

    img_top = y_after_title
    cmd_h = Inches(1.55)
    shape_h = Inches(0.35)
    footer_reserve = cmd_h + shape_h + Inches(0.2)
    max_img_h = SLIDE_H - img_top - footer_reserve - MARGIN

    png_rel = run.get("hero_png")
    png_path = REPO / png_rel if png_rel else Path()
    img_bottom = _add_picture(slide, png_path, MARGIN, img_top, CONTENT_W, max_img_h)

    cmd_top = img_bottom + Inches(0.1)
    cmd_text = run.get("command") or run.get("command_one_line") or "(no command recorded)"
    _add_mono_box(slide, MARGIN, cmd_top, CONTENT_W, cmd_h, cmd_text, font_size=8)

    shape_top = cmd_top + cmd_h + Inches(0.05)
    shape_box = slide.shapes.add_textbox(MARGIN, shape_top, CONTENT_W, shape_h)
    sf = shape_box.text_frame
    sf.paragraphs[0].text = _shape_line(run.get("shape") or {})
    sf.paragraphs[0].font.size = Pt(11)


def _add_intro_slide(prs: Presentation, manifest: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    season_win = manifest.get("season_window") or manifest.get("season_windows") or "11_21"
    if isinstance(season_win, list):
        season_win = ", ".join(season_win)
    _add_title(
        slide,
        "HERO permutation sweep",
        subtitle=(
            f"Tier {manifest.get('tier', '?')} · season {season_win} · "
            f"{manifest.get('n_specs', '?')} slides · {manifest.get('date', '')}"
        ),
    )
    body_top = Inches(1.1)
    box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        "Purpose: compare roster-x × Y-mode × panel-rows (and optional +DFT / EW20).",
        "",
        "Three knobs are orthogonal:",
        "  • roster-x: poolq_loo (LOO teammates) vs poolq (team mean, incl. self)",
        "  • y-draft-mode: ever (Y=1 all PS for draftees) vs season (Y=1 last PS only)",
        "  • panel-rows: all-ps vs last-ps (cross-section)",
        "",
        "Slide 2 = FIXED HERO baseline. Later slides = other combinations.",
        "Command box on each slide = exact CLI to reproduce.",
        "",
        "Regenerate: python sports/scripts/hero_permutation_sweep.py --tier core|extended|full|real_full …",
        "Rebuild deck: python sports/scripts/build_hero_permutation_slides.py",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build HERO permutation PowerPoint.")
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()

    manifest_path = args.manifest if args.manifest.is_absolute() else REPO / args.manifest
    if not manifest_path.is_file():
        raise SystemExit(
            f"Missing manifest {manifest_path}\n"
            "Run: python sports/scripts/hero_permutation_sweep.py --tier core"
        )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    runs = manifest.get("runs") or []
    if not runs:
        raise SystemExit("Manifest has no runs — run sweep first (not --dry-run only).")

    baseline_tag = (manifest.get("baseline") or {}).get("output_tag", "FIXED_HERO")
    out_path = args.out if args.out.is_absolute() else REPO / args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _add_intro_slide(prs, manifest)
    for run in runs:
        _add_slide(prs, run, baseline_tag=baseline_tag)

    prs.save(str(out_path))
    print(f"Wrote {out_path.relative_to(REPO)} · {len(runs)+1} slides (intro + {len(runs)} runs)")


if __name__ == "__main__":
    main()
