#!/usr/bin/env python3
"""Build two-slide θ characterization deck (OAT + θ×K/N panel).

Run (repo root):
  python sports/scripts/build_theta_characterization_slide.py

Outputs:
  HEROs_and_PASSes/slides/auto/CHAR_theta_characterization_AUTO.pptx
    slide 1 — θ OAT at K/N=10%
    slide 2 — θ × K/N peak-bin panel
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from hero_gallery_paths import AUTO_THETA_DECK, THETA, ensure_hero_dirs

OUT_PPTX = AUTO_THETA_DECK
THETA_OAT_SCRIPT = SCRIPTS / "theta_oat_diagnostic.py"
THETA_KN_SCRIPT = SCRIPTS / "theta_kn_sweep_diagnostic.py"
FIG_OAT = THETA / "THETA_OAT_selection_by_pool_mean.png"
FIG_KN = THETA / "THETA_KN_sweep_peak_bin.png"
META_OAT = THETA / "THETA_OAT_meta.json"

from gallery_knobs import LAMBDA_FIXED_RHO, LAMBDA_MODERATE
from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_W = Inches(4.55)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

CLAIM_OAT = (
    "Claim: at fixed $\\rho$, $\\lambda$, and $K/N=10\\%$, $\\theta$ (viability cutline) "
    "moves the readout — low $\\theta$ keeps a mid-pool hump; high $\\theta$ "
    "top-saturates the curve."
)

CLAIM_KN = (
    "Claim: $\\theta$ and $K/N$ co-vary — at MBB-like selectivity ($1\\%$) raising $\\theta$ "
    "shifts peak bin up the pool ladder; at $40\\%$ selectivity curves are top-saturated "
    "regardless of $\\theta$."
)

SLIDE_VARIANTS = [
    {
        "title": r"Phase B — Characterize $\theta$ (viability cutline) — OAT at $K/N=10\%$",
        "figure": FIG_OAT,
        "claim": CLAIM_OAT,
        "params_kind": "oat",
        "notes": [
            r"OAT: soft $\rho=8$, $\lambda=0.55$, $\gamma=10$ fixed; only $\theta$ in $\sigma(\gamma(A-\theta))$ changes.",
            r"$\theta=0.50$: mid-pool hump (peak bin $\approx 13$).",
            r"$\theta=0.72$ (539 default): similar hump, slightly higher peak rate.",
            r"$\theta=0.90$: top bins dominate — curve becomes monotone-increasing at the elite edge.",
        ],
    },
    {
        "title": r"Phase B — Characterize $\theta$ — panel with $K/N$ (selectivity)",
        "figure": FIG_KN,
        "claim": CLAIM_KN,
        "params_kind": "kn_panel",
        "notes": [
            r"Grid: $\theta \in \{0.50, 0.72, 0.90\}$ $\times$ $K/N \in \{1\%, 10\%, 40\%\}$ on $N=5600$.",
            r"At $K/N=1\%$: peak bin rises $6 \to 9 \to 12$ as $\theta$ increases.",
            r"At $K/N=10\%$: $\theta=0.90$ pushes peak to bin $16$ (top saturation).",
            r"At $K/N=40\%$: all arms peak at bin $16$ — selectivity dominates $\theta$.",
            r"Do not fix $\theta = f(K/N)$ without Alex — panel is descriptive.",
        ],
    },
]


def _load_meta() -> dict:
    if META_OAT.is_file():
        return json.loads(META_OAT.read_text(encoding="utf-8"))
    return {"preset": "539", "gamma": 10.0}


def _regenerate_figures() -> None:
    for script, label in (
        (THETA_KN_SCRIPT, "θ×K/N panel"),
        (THETA_OAT_SCRIPT, "θ OAT"),
    ):
        print(f"Regenerating {label} ...")
        subprocess.run([sys.executable, str(script)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2


def _add_params_column(slide, meta: dict, *, left, top, width, height, kind: str) -> None:
    gamma = float(meta.get("gamma", meta.get("viability_sharpness", 10.0)))
    preset = meta.get("preset", "539")

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    head = tf.paragraphs[0]
    subtitle = (
        r"(assign + score + top-$K$ fixed)"
        if kind == "oat"
        else r"($K/N$ varies; assign + score rule fixed)"
    )
    populate_paragraph_with_latex(
        head,
        rf"Characterize $\theta$ — one-at-a-time at {preset} baseline {subtitle}",
        font_size=11,
    )
    head.font.bold = True
    head.space_after = 6

    eq = tf.add_paragraph()
    eq.text = ""
    populate_paragraph_with_latex(
        eq,
        r"$L_C = \mathrm{mean}_{j \neq i}\,\sigma\!\big(\gamma(A_j - \theta)\big)$",
        font_size=11,
    )
    eq.space_after = 4

    if kind == "oat":
        bullets = [
            rf"$\rho={LAMBDA_FIXED_RHO:g}$, $\lambda={LAMBDA_MODERATE:g}$, $\gamma={gamma:g}$ fixed",
            r"arms: $\theta \in \{0.50, 0.72, 0.90\}$",
            r"$K/N = 10\%$ ($K=560$, $N=5600$)",
            r"VISUALIZE = mean $Y_{\mathrm{selected}}$ vs pool mean ($16$ bins)",
        ]
    else:
        bullets = [
            rf"$\rho={LAMBDA_FIXED_RHO:g}$, $\lambda={LAMBDA_MODERATE:g}$, $\gamma={gamma:g}$ fixed",
            r"$\theta \in \{0.50, 0.72, 0.90\}$",
            r"$K/N \in \{1\%, 10\%, 40\%\}$ — heatmap = peak pool-mean bin",
            r"Readout: where the inverted-$U$ peak sits vs selectivity",
        ]

    for line in bullets:
        para = tf.add_paragraph()
        para.text = ""
        populate_paragraph_with_latex(para, line, font_size=10)
        para.space_after = 2


def _add_slide(prs: Presentation, meta: dict, variant: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.28)
    title_h = Inches(0.46)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_with_latex(
        title_box.text_frame.paragraphs[0],
        variant["title"],
        font_size=20,
    )
    title_box.text_frame.paragraphs[0].font.bold = True

    body_top = title_top + title_h + Inches(0.1)
    footer_h = Inches(0.52)
    footer_top = SLIDE_H - MARGIN - footer_h
    notes_h = Inches(1.35)
    params_h = footer_top - body_top - notes_h - Inches(0.08)

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_params_column(
        slide,
        meta,
        left=left_x,
        top=body_top,
        width=LEFT_W,
        height=params_h,
        kind=variant["params_kind"],
    )
    _add_picture_fitted(slide, variant["figure"], right_x, body_top, RIGHT_W, params_h)

    notes_box = slide.shapes.add_textbox(
        left_x,
        footer_top - notes_h - Inches(0.06),
        LEFT_W + RIGHT_W + COL_GAP,
        notes_h,
    )
    fill_bullets_latex(notes_box.text_frame, variant["notes"], font_size=9)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_with_latex(
        foot.text_frame.paragraphs[0], variant["claim"], font_size=11
    )
    foot.text_frame.paragraphs[0].font.bold = True
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def main() -> None:
    ensure_hero_dirs()
    _regenerate_figures()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for variant in SLIDE_VARIANTS:
        _add_slide(prs, meta, variant)

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} ({len(SLIDE_VARIANTS)} slides)")


if __name__ == "__main__":
    main()
