#!/usr/bin/env python3
"""Build PowerPoint from reigning hero basic data plots (porch BDPs).

Each slide: title, lock context, prose, PNG, regenerate CLI, stats footer.

Run (repo root, after reigning_hero_basic_plots.py):
  python sports/scripts/build_reigning_hero_basic_plots_slides.py
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from build_hero_permutation_slides import (  # noqa: E402
    CONTENT_W,
    MARGIN,
    SLIDE_H,
    SLIDE_W,
    _add_mono_box,
    _add_picture,
    _add_title,
)
from hero_gallery_paths import REIGNING_HERO_BASIC_PLOTS  # noqa: E402
from reigning_hero_star_diff import REIGNING_LOCK_TAG  # noqa: E402

DEFAULT_OUT = REIGNING_HERO_BASIC_PLOTS / "REIGNING_BDP_slides_AUTO.pptx"
MANIFEST = REIGNING_HERO_BASIC_PLOTS / "manifest.json"
SPEC = "mg10 min20 09_21"

LOCK_LINES = [
    "Reigning hero (slide 12): poolq_LOO · ever · last-ps · EW16 · ALLT · 2009–2021.",
    "Filters: min20 · mg10 · PPM z · winsor 1–99 on poolq_LOO (LOO plots only).",
    f"HERO tag: `{REIGNING_LOCK_TAG}`.",
]


def _add_prose_box(slide, top: float, lines: list[str], *, height: float = 0.78) -> float:
    h = Inches(height)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}" if line and not line.startswith(" ") else line
        p.font.size = Pt(11)
        p.space_after = Pt(2)
    return top + h + Inches(0.06)


def _load_json(path: Path) -> dict[str, Any]:
    if path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def _summary_ai_tj(meta: dict) -> str:
    ai = meta.get("A_i_hat") or {}
    tj = meta.get("T_j_hat") or {}
    dft = meta.get("dft_overlay") or {}
    parts = [
        f"Â: n={ai.get('n', '?'):,} · mean={ai.get('mean', 0):+.3f} · sd={ai.get('std', 0):.3f}",
        f"T̂_j: n={tj.get('n', '?'):,} · mean={tj.get('mean', 0):+.3f} · sd={tj.get('std', 0):.3f}",
    ]
    if dft.get("A_i_hat"):
        dai = dft["A_i_hat"]
        parts.append(f"+DFT overlay: n={dai.get('n', '?'):,} player-seasons on draft-ever teams")
    return " · ".join(parts)


def _summary_loo(meta: dict) -> str:
    ai = meta.get("A_i_hat") or {}
    loo = meta.get("poolq_loo") or meta.get("T_j_hat") or {}
    winsor = meta.get("poolq_winsor_quantiles")
    w = f"winsor {winsor[0]:.0%}–{winsor[1]:.0%}" if winsor else "no winsor"
    return (
        f"Â: n={ai.get('n', '?'):,} · sd={ai.get('std', 0):.3f} · "
        f"poolq_LOO: n={loo.get('n', '?'):,} · sd={loo.get('std', 0):.3f} · {w}"
    )


def _summary_tj_loo(meta: dict) -> str:
    tj = meta.get("T_j_hat") or {}
    loo = meta.get("poolq_loo") or {}
    winsor = meta.get("poolq_winsor_quantiles")
    w = f"winsor {winsor[0]:.0%}–{winsor[1]:.0%}" if winsor else "no winsor"
    return (
        f"T̂_j: sd={tj.get('std', 0):.3f} · poolq_LOO: sd={loo.get('std', 0):.3f} · {w} · "
        "tail spikes at winsor clip = artifact"
    )


def _summary_overlap(meta: dict) -> str:
    return (
        f"n={meta.get('n_team_seasons', '?'):,} team-seasons · "
        f"H_sort={meta.get('H_sort', 0):.4f} · "
        f"coverage mean={meta.get('coverage_mean', 0):.0f} · "
        f"disjoint mean={meta.get('coverage_disjoint_mean', 0):.2%}"
    )


def _summary_draft_rate(meta: dict) -> str:
    metric = str(meta.get("metric", "")).upper()
    vents = meta.get("player_ventiles") or []
    if not vents:
        return f"{metric} · panel n={meta.get('panel_n', '?'):,}"
    lo = vents[0]
    hi = vents[-1]
    return (
        f"{metric} · all-ps n={meta.get('panel_n', '?'):,} · "
        f"draft rate vent1={100*lo.get('draft_rate', 0):.2f}% · "
        f"vent16={100*hi.get('draft_rate', 0):.2f}% · orange=+DFT overlay"
    )


def _summary_team_size(meta: dict) -> str:
    ts = meta.get("team_season_summary") or {}
    return (
        f"|T_j| mean={ts.get('mean', 0):.2f} · median={ts.get('median', 0):.0f} · "
        f"n={ts.get('n', '?'):,} team-seasons · all-ps · orange=+DFT"
    )


def _summary_team_games(meta: dict) -> str:
    return (
        f"n={meta.get('n_team_seasons', '?'):,} team-seasons · "
        f"median={meta.get('games_n_median', 0):.0f} games · "
        f"mean={meta.get('games_n_mean', 0):.1f} · after box QC"
    )


def _summary_loo_dist(meta: dict) -> str:
    loo = meta.get("poolq_loo") or {}
    loo_d = meta.get("poolq_loo_dft") or {}
    winsor = meta.get("poolq_winsor")
    w = (
        f"winsor {winsor[0]:.0%}–{winsor[1]:.0%}"
        if winsor
        else "no poolq_LOO winsor"
    )
    line = (
        f"last-ps n={loo.get('n', '?'):,} · med={loo.get('median', 0):.3f} · "
        f"μ={loo.get('mean', 0):.3f} · {w}"
    )
    if loo_d:
        line += f" · +DFT med={loo_d.get('median', 0):.3f}"
    return line


def _summary_draft_rate_loo(meta: dict) -> str:
    coef = meta.get("lpm_quadratic") or {}
    b2 = coef.get("beta_poolq_loo_sq", 0)
    return (
        f"EW{meta.get('n_bins', 16)} · n={meta.get('n_panel_rows', '?'):,} · "
        f"LPM β₂={b2:+.5f} · ever-Y · last-ps"
    )


def _summary_draft_rate_loo_tj(meta: dict) -> str:
    loo = meta.get("lpm_poolq_loo") or {}
    tj = meta.get("lpm_t_j") or {}
    return (
        f"EW{meta.get('n_bins', 16)} · n={meta.get('n_panel_rows', '?'):,} · "
        f"LOO β₂={loo.get('beta_poolq_loo_sq', 0):+.5f} · "
        f"T̂_j β₂={tj.get('beta_t_j_sq', 0):+.5f}"
    )


def _summary_minutes(meta: dict) -> str:
    ps = meta.get("player_season_minutes") or {}
    ts = meta.get("team_mean_minutes") or {}
    ps_d = meta.get("player_season_minutes_dft") or {}
    line = (
        f"player PS: n={ps.get('n', '?'):,} · med={ps.get('median', 0):.0f} min · "
        f"team mean: n={ts.get('n', '?'):,} · med={ts.get('median', 0):.0f} min · all-ps post-min20"
    )
    if ps_d:
        line += f" · +DFT player med={ps_d.get('median', 0):.0f} min"
    return line


def _summary_hsort_dist(meta: dict) -> str:
    pooled = meta.get("H_sort_pooled_empirical")
    sim_n = meta.get("H_sort_sim_rho0_n")
    sim_mean = meta.get("H_sort_sim_rho0_mean")
    line = f"pooled empirical H_sort={pooled:.4f}" if pooled is not None else "pooled empirical H_sort=?"
    if sim_n:
        line += f" · sim ρ=0 mean={sim_mean:.4f} (n={sim_n})"
    return line


def _summary_ability_residuals(meta: dict) -> str:
    h = meta.get("H_sort_pooled")
    std_t = meta.get("std_team_residual")
    std_g = meta.get("std_grand_residual")
    dft = meta.get("dft") or {}
    line = (
        f"n={meta.get('n', '?'):,} · H_sort={h:.4f} · "
        f"σ(team resid)={std_t:.3f} · σ(grand resid)={std_g:.3f}"
    )
    if dft:
        line += f" · +DFT H_sort={dft.get('H_sort_pooled', 0):.4f}"
    return line


def _summary_ability_loo_residuals(meta: dict) -> str:
    dft = meta.get("dft") or {}
    line = (
        f"n={meta.get('n', '?'):,} · σ(LOO resid)={meta.get('std_loo_residual', 0):.3f} · "
        f"frac_var={meta.get('frac_var_loo_residual', 0):.3f}"
    )
    if dft:
        line += f" · +DFT σ={dft.get('std_loo_residual', 0):.3f}"
    return line


def _summary_ability_tj_vs_loo_residuals(meta: dict) -> str:
    return (
        f"n={meta.get('n', '?'):,} · σ(team resid)={meta.get('std_team_residual', 0):.3f} · "
        f"σ(LOO resid)={meta.get('std_loo_residual', 0):.3f} · "
        f"corr={meta.get('corr_team_loo_residual', 0):.4f}"
    )


def _summary_tj_minus_loo(meta: dict) -> str:
    dft = meta.get("dft") or {}
    line = (
        f"n={meta.get('n', '?'):,} · μ(T̂_j−LOO)={meta.get('mean_tj_minus_loo', 0):+.4f} · "
        f"σ={meta.get('std_tj_minus_loo', 0):.3f} · H_sort^team={meta.get('H_sort_team', 0):.4f} · "
        f"H_sort^LOO={meta.get('H_sort_loo', 0):.4f}"
    )
    if dft:
        line += f" · +DFT σ={dft.get('std_tj_minus_loo', 0):.3f}"
    return line


def _plot_catalog(spec_slug: str) -> list[dict[str, Any]]:
    """Ordered porch plots — stems match reigning_hero_basic_plots.py."""
    return [
        {
            "key": "overlap",
            "title": "Team interval overlap",
            "subtitle": "HAND analog ~PD20 slide 21 · all-ps",
            "prose": [
                "Full roster panel (all player-seasons) — needed for team-season interval geometry.",
                "Same filter chain as lock: mg10 · min20 · 09–21 · PPM z · ALLT.",
                "Complements last-ps HERO aperture — shows roster spread before LOO cross-section.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_team_interval_overlap_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_team_interval_overlap_{spec_slug}_meta.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only overlap",
            "summary_fn": _summary_overlap,
            "panel_rows": "all-ps",
            "hand_analog": "PD20 ~slide 21",
        },
        {
            "key": "ai_tj",
            "title": r"Â_i and T̂_j distributions",
            "subtitle": "BDP slide 7 analog · last-ps · ALLT + orange +DFT",
            "prose": [
                "Last-ps cross-section (one row per athlete) — matches HERO scoring aperture.",
                "Left: Â_i (player PPM z). Right: T̂_j (team-season mean). No poolq winsor.",
                "Orange overlay = +DFT subset (player-seasons on draft-ever teams).",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Ai_Tj_{spec_slug}_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Ai_Tj_{spec_slug}_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only ai_tj",
            "summary_fn": _summary_ai_tj,
            "panel_rows": "last-ps",
            "hand_analog": "BDP slide 7",
        },
        {
            "key": "ai_loo",
            "title": r"Â_i and poolq_LOO distributions",
            "subtitle": "Mirror of Â|T̂_j · last-ps · winsor 1–99 on LOO",
            "prose": [
                "Right panel = LOO teammate mean (poolq_LOO) with reigning winsor 1–99.",
                "Same last-ps aperture as HERO lock; shows individual vs local congestion.",
                "Tail structure on LOO reflects winsor clip — not LOO algebra.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Ai_poolq_loo_{spec_slug}_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Ai_poolq_loo_{spec_slug}_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only ai_loo",
            "summary_fn": _summary_loo,
            "panel_rows": "last-ps",
            "hand_analog": "BDP slide 7 (LOO variant)",
        },
        {
            "key": "tj_loo",
            "title": r"T̂_j vs poolq_LOO density overlay",
            "subtitle": "New porch diagnostic · last-ps · winsor on LOO only",
            "prose": [
                "Compares team mean T̂_j (broad) vs LOO teammate mean (tight) on shared PPM z axis.",
                "LOO winsor 1–99 — tail spikes at ±0.6 are clip artifacts, not biology.",
                "Motivation: visualize gap between team talent and local congestion before HERO bins.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Tj_vs_poolq_loo_{spec_slug}_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Tj_vs_poolq_loo_{spec_slug}_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only tj_loo",
            "summary_fn": _summary_tj_loo,
            "panel_rows": "last-ps",
            "hand_analog": "(new)",
        },
        {
            "key": "tj_loo_nowinsor",
            "title": r"T̂_j vs poolq_LOO density overlay (no winsor)",
            "subtitle": "Mirror of T̂_j vs LOO · last-ps · raw poolq_LOO",
            "prose": [
                "Same hist-density overlay as prior slide — poolq_LOO winsor 1–99 removed.",
                "Shows natural LOO tails without clip spikes at ±0.6.",
                "Reigning HERO lock still uses winsor; this is porch diagnostic only.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_Tj_vs_poolq_loo_{spec_slug}_ppm_lastps_NOWINSOR.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_Tj_vs_poolq_loo_{spec_slug}_ppm_lastps_NOWINSOR.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only tj_loo_nowinsor",
            "summary_fn": _summary_tj_loo,
            "panel_rows": "last-ps",
            "hand_analog": "(new)",
        },
        {
            "key": "hsort_dist",
            "title": r"Empirical $H_{\mathrm{sort}}$ distribution",
            "subtitle": "Realized sorting index · all-ps · vs sim ρ=0 reference",
            "prose": [
                r"$H_{\mathrm{sort}}$ = variance explained by team assignment (not interval overlap).",
                "Per-season empirical bars + pooled line; orange = Grandchild sim ρ=0 distribution.",
                "Pooled reigning lock ≈ 0.064 — weak homophily on realized rosters.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Hsort_dist_{spec_slug}_ppm_allps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_Hsort_dist_{spec_slug}_ppm_allps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only hsort_dist",
            "summary_fn": _summary_hsort_dist,
            "panel_rows": "all-ps",
            "hand_analog": "PD17 H_sort readout",
        },
        {
            "key": "ability_residuals",
            "title": r"Ability residuals — team vs grand mean",
            "subtitle": r"$H_{\mathrm{sort}}$ decomposition · all-ps · ALLT + purple +DFT",
            "prose": [
                r"Left: $A_i - \hat{T}_j$ (within-team demeaned perf — $H_{\mathrm{sort}}$ numerator).",
                r"Right: $A_i - \bar{A}$ (grand-mean demeaned — denominator).",
                "Pooled over 09–21; purple overlay = +DFT (draft-ever teams).",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_residuals_{spec_slug}_ppm_allps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_residuals_{spec_slug}_ppm_allps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only ability_residuals",
            "summary_fn": _summary_ability_residuals,
            "panel_rows": "all-ps",
            "hand_analog": "BDP H_sort porch",
        },
        {
            "key": "ability_loo_residuals",
            "title": r"Ability minus LOO ($A_i - \mathrm{poolq}^{\mathrm{LOO}}_i$)",
            "subtitle": "LOO congestion residual · all-ps · ALLT + purple +DFT",
            "prose": [
                r"Histogram + ECDF of $A_i - \mathrm{poolq}^{\mathrm{LOO}}_i$ (teammate mean excl. $i$).",
                "Analogous to team residual panel but LOO replaces team mean in the decomposition.",
                "σ ≈ 1.05 > team residual σ — LOO axis carries more player-level spread.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_loo_residuals_{spec_slug}_ppm_allps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_loo_residuals_{spec_slug}_ppm_allps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only ability_loo_residuals",
            "summary_fn": _summary_ability_loo_residuals,
            "panel_rows": "all-ps",
            "hand_analog": "BDP LOO porch",
        },
        {
            "key": "ability_tj_vs_loo_residuals",
            "title": r"Team vs LOO ability residuals (overlay)",
            "subtitle": "Superposed KDEs · all-ps · orange team · blue LOO",
            "prose": [
                r"Orange: $A_i - \hat{T}_j$ (team mean incl. $i$). Blue: $A_i - \mathrm{poolq}^{\mathrm{LOO}}_i$.",
                "Same player-rows; corr ≈ 0.9997 — nearly collinear but LOO slightly wider.",
                "Motivation: visualize self-inclusion gap before T̂_j − LOO slide.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_Tj_vs_loo_residuals_{spec_slug}_ppm_allps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_ability_Tj_vs_loo_residuals_{spec_slug}_ppm_allps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only ability_tj_vs_loo_residuals",
            "summary_fn": _summary_ability_tj_vs_loo_residuals,
            "panel_rows": "all-ps",
            "hand_analog": "BDP LOO porch",
        },
        {
            "key": "tj_minus_loo",
            "title": r"$\hat{T}_j - \mathrm{poolq}^{\mathrm{LOO}}_i$ (self-inclusion gap)",
            "subtitle": "Team mean minus LOO teammate mean · all-ps · ALLT + purple +DFT",
            "prose": [
                r"Distribution of $\hat{T}_j - \mathrm{poolq}^{\mathrm{LOO}}_i$ — how much self pulls team mean above LOO.",
                r"Reports $H_{\mathrm{sort}}^{\mathrm{team}}$ vs $H_{\mathrm{sort}}^{\mathrm{LOO}}$ (LOO can go negative).",
                "Tight spike near 0 (σ ≈ 0.09); axis-mismatch diagnostic for SCORE vs HERO readout.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_Tj_minus_loo_{spec_slug}_ppm_allps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_Tj_minus_loo_{spec_slug}_ppm_allps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only tj_minus_loo",
            "summary_fn": _summary_tj_minus_loo,
            "panel_rows": "all-ps",
            "hand_analog": "BDP LOO porch",
        },
        {
            "key": "loo_dist",
            "title": r"Player poolq_LOO distribution",
            "subtitle": "Alex Aug 2026 · last-ps · winsor 1–99 · ALLT + orange +DFT",
            "prose": [
                "Standalone distribution of player LOO teammate quality (poolq_LOO).",
                "Histogram + ECDF on reigning last-ps aperture — not paired with Â_i.",
                "Orange = +DFT overlay (draft-ever teams).",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_poolq_loo_dist_{spec_slug}_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_poolq_loo_dist_{spec_slug}_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only loo_dist",
            "summary_fn": _summary_loo_dist,
            "panel_rows": "last-ps",
            "hand_analog": "(new · Alex request)",
        },
        {
            "key": "loo_dist_nowinsor",
            "title": r"Player poolq_LOO distribution (no winsor)",
            "subtitle": "Mirror of LOO dist · last-ps · raw poolq_LOO · ALLT + orange +DFT",
            "prose": [
                "Same hist + ECDF pair as prior slide — poolq_LOO winsor 1–99 removed.",
                "Shows natural LOO tails without clip spikes at ±0.6.",
                "Reigning HERO lock still uses winsor; this is porch diagnostic only.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_poolq_loo_dist_{spec_slug}_ppm_lastps_nowinsor.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_poolq_loo_dist_{spec_slug}_ppm_lastps_nowinsor.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only loo_dist_nowinsor",
            "summary_fn": _summary_loo_dist,
            "panel_rows": "last-ps",
            "hand_analog": "(new · Alex request)",
        },
        {
            "key": "draft_rate_loo",
            "title": r"P(Y=1) vs player poolq_LOO",
            "subtitle": "Reigning lock EW16 · ever-Y · last-ps",
            "prose": [
                "Mean Y_draft by equal-width poolq_LOO bins (16) — same binning as reigning HERO.",
                "Bar height = empirical draft rate; labels = bin n; LPM β₂ annotation.",
                "Complements HERO PNG in ../hero/ — porch copy for slide 12 deck.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_draft_rate_poolq_loo_{spec_slug}_ew16_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_draft_rate_poolq_loo_{spec_slug}_ew16_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only draft_rate_loo",
            "summary_fn": _summary_draft_rate_loo,
            "panel_rows": "last-ps",
            "hand_analog": "HERO / Pass A roster panel",
        },
        {
            "key": "draft_rate_loo_tj",
            "title": r"P(Y=1) — player LOO vs T̂_j (side-by-side)",
            "subtitle": "Reigning lock EW16 · ever-Y · last-ps · same panel both panels",
            "prose": [
                "2×2: top row EW16 equal-width; bottom row quantile 16 (same n per bin).",
                "Left: HERO poolq_LOO · Right: team T̂_j — same panel, both binning schemes.",
                "LPM β₂ on continuous x (not binned); compare shape sensitivity to binning.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_draft_rate_poolq_loo_vs_Tj_{spec_slug}_ew16_ppm_lastps.png",
            "meta": REIGNING_HERO_BASIC_PLOTS
            / f"REIGNING_BDP_draft_rate_poolq_loo_vs_Tj_{spec_slug}_ew16_ppm_lastps.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only draft_rate_loo_tj",
            "summary_fn": _summary_draft_rate_loo_tj,
            "panel_rows": "last-ps",
            "hand_analog": "HERO vs F-HERO team axis",
        },
        {
            "key": "draft_rate_apgms",
            "title": "Draft rate vs APGMS ventiles",
            "subtitle": "Fixed_Ai slides 11–14 analog · all-ps · ALLT + orange +DFT",
            "prose": [
                "Exposure diagnostic: average player-games per season (APGMS) vs draft rate.",
                "All-ps panel (full minutes exposure) — not last-ps HERO aperture.",
                "Quantile ventiles on APGMS; independent of HERO poolq_LOO binning.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_draft_rate_APGMS_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_draft_rate_APGMS_{spec_slug}.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only draft_rate",
            "summary_fn": _summary_draft_rate,
            "panel_rows": "all-ps",
            "hand_analog": "Fixed_Ai 11–14",
        },
        {
            "key": "draft_rate_argms",
            "title": "Draft rate vs ARGMS ventiles",
            "subtitle": "Fixed_Ai slides 11–14 analog · all-ps · ALLT + orange +DFT",
            "prose": [
                "Exposure diagnostic: average roster-games per season (ARGMS) vs draft rate.",
                "Team-level minutes exposure complement to APGMS player ventiles.",
                "Same 09–21 lock filters; orange +DFT overlay on draft-ever teams.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_draft_rate_ARGMS_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_draft_rate_ARGMS_{spec_slug}.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only draft_rate",
            "summary_fn": _summary_draft_rate,
            "panel_rows": "all-ps",
            "hand_analog": "Fixed_Ai 11–14",
        },
        {
            "key": "team_games",
            "title": "Team games per season",
            "subtitle": "PD22 memo slide 27 analog · after box QC (mg10)",
            "prose": [
                "Distinct game_id counts per (team_id, season) in ESPN box — after dash-QC and mg10.",
                "Dual panel: linear counts + log-y share (same layout as PD22 backup slide).",
                "Sparse one-game team-seasons dropped by mg10; mass clusters near full-season ~30 games.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_team_games_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_team_games_{spec_slug}.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only team_games",
            "summary_fn": _summary_team_games,
            "panel_rows": "team-season (box)",
            "hand_analog": "PD22 memo slide 27",
        },
        {
            "key": "minutes",
            "title": "Season minutes — player vs team mean",
            "subtitle": "Exposure porch · all-ps · min20 floor marked",
            "prose": [
                "2×2: player histogram (+ log-y inset) and ECDF; team mean ditto.",
                "Blue = full panel; orange line = +DFT (draft-ever teams only).",
                "Stats boxes on ECDF panels show w/o DFT and +DFT medians.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_minutes_player_team_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_minutes_player_team_{spec_slug}.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only minutes",
            "summary_fn": _summary_minutes,
            "panel_rows": "all-ps",
            "hand_analog": "PD22 minutes (porch variant)",
        },
        {
            "key": "team_size",
            "title": "Team roster size |T_j|",
            "subtitle": "BDP slide 18 analog · all-ps · ALLT + orange +DFT",
            "prose": [
                "Roster depth distribution — full panel (all-ps), not last-ps cross-section.",
                "Validates team-season cardinality under min20 · mg10 lock filters.",
                "Orange +DFT overlay: teams with at least one draftee in sample.",
            ],
            "png": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_team_size_{spec_slug}.png",
            "meta": REIGNING_HERO_BASIC_PLOTS / f"REIGNING_BDP_team_size_{spec_slug}.json",
            "command": "python sports/scripts/reigning_hero_basic_plots.py --only team_size",
            "summary_fn": _summary_team_size,
            "panel_rows": "all-ps",
            "hand_analog": "BDP slide 18",
        },
    ]


def _manifest_entry(entry: dict, meta: dict, *, slide_index: int | None = None) -> dict:
    out = {
        "key": entry["key"],
        "title": entry["title"],
        "subtitle": entry["subtitle"],
        "panel_rows": entry["panel_rows"],
        "hand_analog": entry["hand_analog"],
        "prose": entry["prose"],
        "command": entry["command"],
        "png": str(entry["png"].relative_to(REPO)) if entry["png"].is_absolute() else entry["png"].name,
        "meta_json": entry["meta"].name,
        "png_exists": entry["png"].is_file(),
        "meta_exists": entry["meta"].is_file(),
    }
    if slide_index is not None:
        out["slide_index"] = slide_index
    if meta:
        out["meta_date"] = meta.get("date")
    return out


def _add_intro_slide(prs: Presentation, *, n_plots: int, n_present: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    status = f"{n_present}/{n_plots} PNGs on disk" if n_present < n_plots else f"{n_plots} plots"
    _add_title(
        slide,
        "Slide 1 — Reigning hero basic data plots (porch)",
        subtitle=f"{status} · {date.today().isoformat()} · lock {SPEC}",
    )
    body_top = Inches(1.05)
    box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        "Purpose: porch diagnostics for the named reigning hero before star sweep / F-HERO.",
        "",
        *LOCK_LINES,
        "",
        "Population split (Charles Aug 2026):",
        "  • last-ps: Â, LOO, T̂_j vs LOO — matches HERO scoring cross-section",
        "  • all-ps: interval overlap, APGMS/ARGMS, roster size — full panel geometry",
        "",
        "Regenerate plots: python sports/scripts/reigning_hero_basic_plots.py",
        "Rebuild deck: python sports/scripts/build_reigning_hero_basic_plots_slides.py",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)


def _add_plot_slide(prs: Presentation, entry: dict, meta: dict, *, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = _add_title(
        slide,
        f"Slide {slide_num} — {entry['title']}",
        subtitle=entry["subtitle"],
    )
    y = _add_prose_box(slide, y, entry["prose"])

    cmd_h = Inches(0.55)
    footer_h = Inches(0.42)
    footer_reserve = cmd_h + footer_h + Inches(0.15)
    max_img_h = SLIDE_H - y - footer_reserve - MARGIN

    img_bottom = _add_picture(slide, entry["png"], MARGIN, y, CONTENT_W, max_img_h)

    cmd_top = img_bottom + Inches(0.08)
    _add_mono_box(slide, MARGIN, cmd_top, CONTENT_W, cmd_h, entry["command"], font_size=8)

    footer_top = cmd_top + cmd_h + Inches(0.04)
    summary_fn = entry["summary_fn"]
    stats = summary_fn(meta) if meta else "(no meta JSON — re-run plot)"
    panel_note = f"Rows: {entry['panel_rows']} · HAND: {entry['hand_analog']}"
    box = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = stats
    p.font.size = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = panel_note
    p2.font.size = Pt(9)
    p2.font.italic = True


def _write_summary_md(catalog: list[dict], manifest: dict, out_path: Path) -> None:
    rows = [
        "# Reigning hero BDP porch — slide index",
        "",
        f"Lock: `{REIGNING_LOCK_TAG}` · `{SPEC}`",
        "",
        "| Slide | Key | Title | Rows | HAND analog | PNG |",
        "|-------|-----|-------|------|-------------|-----|",
    ]
    for item in manifest.get("plots") or []:
        png = "✓" if item.get("png_exists") else "—"
        rows.append(
            f"| {item.get('slide_index', '?')} | `{item.get('key', '')}` | "
            f"{item.get('title', '')} | {item.get('panel_rows', '')} | "
            f"{item.get('hand_analog', '')} | {png} |"
        )
    out_path.write_text("\n".join(rows) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build reigning hero BDP porch PowerPoint.")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    parser.add_argument(
        "--write-manifest",
        action="store_true",
        default=True,
        help="Write manifest.json with slide_index (default: on).",
    )
    parser.add_argument(
        "--no-write-manifest",
        action="store_false",
        dest="write_manifest",
        help="Do not rewrite manifest.json.",
    )
    args = parser.parse_args()

    spec_slug = SPEC.replace(" ", "_")
    catalog = _plot_catalog(spec_slug)
    n_present = sum(1 for e in catalog if e["png"].is_file())

    out_path = args.out if args.out.is_absolute() else REPO / args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _add_intro_slide(prs, n_plots=len(catalog), n_present=n_present)

    manifest_plots: list[dict] = []
    for entry in catalog:
        slide_num = len(prs.slides) + 1
        meta = _load_json(entry["meta"])
        _add_plot_slide(prs, entry, meta, slide_num=slide_num)
        manifest_plots.append(_manifest_entry(entry, meta, slide_index=slide_num))

    prs.save(str(out_path))

    manifest = {
        "created": date.today().isoformat(),
        "deck": "reigning_hero_basic_data_plots",
        "reigning_spec": SPEC,
        "reigning_lock_tag": REIGNING_LOCK_TAG,
        "n_plots": len(catalog),
        "n_png_present": n_present,
        "out_pptx": out_path.name,
        "plots": manifest_plots,
    }
    summary_md = out_path.parent / "REIGNING_BDP_slides_summary.md"
    _write_summary_md(catalog, manifest, summary_md)

    if args.write_manifest:
        MANIFEST.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
        print(f"Wrote {MANIFEST.relative_to(REPO)}", flush=True)

    print(
        f"Wrote {out_path.relative_to(REPO)} · {len(prs.slides)} slides "
        f"(intro + {len(catalog)} plots · {n_present}/{len(catalog)} PNGs found)",
        flush=True,
    )
    print(f"Wrote {summary_md.relative_to(REPO)}", flush=True)
    if n_present < len(catalog):
        print(
            "WARNING: some PNGs missing — run: python sports/scripts/reigning_hero_basic_plots.py",
            flush=True,
        )


if __name__ == "__main__":
    main()
