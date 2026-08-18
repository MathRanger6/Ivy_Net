"""Mathtext for Pass A/B/C gallery PNGs (matplotlib) and slides (python-pptx runs).

Matplotlib: wrap math in $...$ in titles/labels.
Slides: same $...$ convention; math segments render as Cambria Math + sub/superscripts.
"""

from __future__ import annotations

import re

_GREEK = {
    r"\lambda": "λ",
    r"\rho": "ρ",
    r"\sigma": "σ",
    r"\pi": "π",
    r"\beta": "β",
    r"\propto": "∝",
    r"\sim": "∼",
    r"\approx": "≈",
    r"\exp": "exp",
    r"\in": "∈",
    r"\cdot": "·",
    r"\,": "·",
    r"\Rightarrow": "⇒",
    r"\rightarrow": "→",
    r"\uparrow": "↑",
    r"\!": "",
}


def configure_matplotlib_mathtext() -> None:
    import matplotlib as mpl

    mpl.rcParams.update(
        {
            "mathtext.fontset": "cm",
            "axes.unicode_minus": False,
        }
    )


def _strip_math_commands(s: str) -> str:
    s = re.sub(r"\\mathrm\{([^}]*)\}", r"\1", s)
    s = re.sub(r"\\texttt\{([^}]*)\}", r"\1", s)
    s = re.sub(r"\\left|\\right", "", s)
    s = re.sub(r"\\frac\{([^}]*)\}\{([^}]*)\}", r"(\1)/(\2)", s)
    s = s.replace(r"\{", "{").replace(r"\}", "}")
    for tok, ch in _GREEK.items():
        s = s.replace(tok, ch)
    return s


def _set_baseline(run, *, subscript: bool = False, superscript: bool = False) -> None:
    """Office Open XML baseline (1/1000 percent); font.subscript is a no-op in pptx 1.0.x."""
    if not subscript and not superscript:
        return
    r_pr = run._r.get_or_add_rPr()
    if subscript:
        r_pr.set("baseline", "-25000")
    elif superscript:
        r_pr.set("baseline", "30000")


def _add_run(
    paragraph,
    text: str,
    *,
    size: int,
    italic: bool = False,
    subscript: bool = False,
) -> None:
    from pptx.util import Pt

    if not text:
        return
    run = paragraph.add_run()
    run.text = text
    run_size = max(size - 2, 8) if subscript else size
    run.font.size = Pt(run_size)
    run.font.name = "Cambria Math"
    run.font.italic = italic
    _set_baseline(run, subscript=subscript)


def _add_superscript_run(paragraph, text: str, *, size: int) -> None:
    from pptx.util import Pt

    if not text:
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(max(size - 2, 8))
    run.font.name = "Cambria Math"
    _set_baseline(run, superscript=True)


def _populate_math_runs(paragraph, math: str, *, size: int) -> None:
    """Render a single $...$ math segment into pptx runs."""
    math = _strip_math_commands(math)
    i = 0
    buf: list[str] = []
    italic = True

    def flush() -> None:
        nonlocal buf, italic
        if buf:
            _add_run(paragraph, "".join(buf), size=size, italic=italic)
            buf = []

    while i < len(math):
        ch = math[i]
        if ch == "_":
            flush()
            i += 1
            if i < len(math) and math[i] == "{":
                j = math.index("}", i)
                sub_text = math[i + 1 : j]
                i = j + 1
            elif i < len(math):
                sub_text = math[i]
                i += 1
            else:
                sub_text = ""
            _add_run(paragraph, sub_text, size=size, italic=True, subscript=True)
            continue
        if ch == "^":
            flush()
            i += 1
            if i < len(math) and math[i] == "{":
                j = math.index("}", i)
                sup_text = math[i + 1 : j]
                i = j + 1
            elif i < len(math):
                sup_text = math[i]
                i += 1
            else:
                sup_text = ""
            _add_superscript_run(paragraph, sup_text, size=size)
            continue
        if ch.isalpha() and ch.isupper() and len(buf) == 0:
            italic = True
        buf.append(ch)
        i += 1
    flush()


def populate_paragraph_with_latex(paragraph, text: str, *, font_size: int = 11) -> None:
    """Parse plain text with $...$ math into a python-pptx paragraph."""
    from pptx.util import Pt

    paragraph.text = ""
    parts = text.split("$")
    for idx, part in enumerate(parts):
        if not part:
            continue
        if idx % 2 == 1:
            _populate_math_runs(paragraph, part, size=font_size)
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(font_size)


def fill_bullets_latex(text_frame, lines: list[str], *, font_size: int = 11) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        populate_paragraph_with_latex(para, line, font_size=font_size)
        para.level = 0
        para.space_after = 3
        para.space_before = 0


# PD17 hand-deck workflow: literal LaTeX in body font → highlight → Insert → Equation → LaTeX
RAW_LATEX_FONT = "Calibri"

# Typography matched to CHAR_*_HAND.pptx (Phase B / ρ / λ characterization decks).
HAND_TITLE_PT = 20
HAND_SUBTITLE_PT = 10
HAND_BODY_PT = 10
HAND_CLAIM_PT = 11
HAND_BULLET_LEAD = "  "  # two spaces between • and bullet text (Charles HAND convention)
HAND_BULLET_LINE_SPACING = 1.5  # line spacing on sidebar bullet boxes (Charles HAND)


def _set_char_bullet(paragraph, *, char: str = "•") -> None:
    """Explicit round bullet — matches HAND intro / sidebar lists."""
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        existing = p_pr.find(qn(tag))
        if existing is not None:
            p_pr.remove(existing)
    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", char)
    p_pr.insert(0, bu_char)


def _sanitize_memo_text(text: str) -> str:
    """Strip control chars that corrupt LaTeX tokens (e.g. \\r in \\rho → _x000D_ho in PPT)."""
    if "\r" in text:
        text = text.replace("\r", "")
    return text


def populate_paragraph_raw_latex(
    paragraph,
    text: str,
    *,
    font_size: int = 11,
    bold: bool = False,
) -> None:
    """Write literal LaTeX (e.g. \\hat{A}_{i}) — no mathtext runs, baseline 0.

    Do NOT wrap in $...$ — Charles converts via highlight → Insert → Equation →
    LaTeX to math. Use braces on subscripts (\\hat{A}_{i}) so PowerPoint Mac
    does not auto-format _i / _j as subscript when the file opens.
    """
    from pptx.util import Pt

    text = _sanitize_memo_text(text)
    paragraph.text = ""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = RAW_LATEX_FONT
    run.font.bold = bold
    run.font.italic = False
    r_pr = run._r.get_or_add_rPr()
    if "baseline" in r_pr.attrib:
        del r_pr.attrib["baseline"]
    r_pr.set("baseline", "0")


def fill_bullets_raw_latex(
    text_frame,
    lines: list[str],
    *,
    font_size: int = HAND_BODY_PT,
    bullet_from: int = 0,
    space_after: int = 0,
) -> None:
    """Bullet list with literal LaTeX tokens (PD17 hand-deck workflow).

    bullet_from: paragraph index to start explicit • bullets (0 = all lines).

    PD17 conventions (Charles HAND decks):
    - line_spacing = 1.5 on every bullet paragraph
    - no $...$ wrappers in slide strings — plain \\hat{A}_{i}, \\rho, etc.
    """
    from pptx.util import Pt

    text_frame.clear()
    text_frame.word_wrap = True
    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        text = line.rstrip()
        if i >= bullet_from:
            text = HAND_BULLET_LEAD + text.lstrip()
        populate_paragraph_raw_latex(para, text, font_size=font_size)
        para.level = 0
        para.space_after = Pt(space_after)
        para.space_before = Pt(0)
        para.line_spacing = HAND_BULLET_LINE_SPACING
        if i >= bullet_from:
            _set_char_bullet(para)


def add_plain_latex_block(
    text_frame,
    latex: str,
    *,
    font_size: int = 11,
    label: str | None = None,
) -> None:
    """Raw LaTeX as plain text — highlight → Insert → Equation → LaTeX to math."""
    from pptx.util import Pt

    if label:
        label_para = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
        if text_frame.text and label_para.text:
            label_para = text_frame.add_paragraph()
        label_para.text = label
        label_para.font.size = Pt(font_size)
        label_para.font.bold = True
        label_para.space_after = 2

    para = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
    if text_frame.text and para.text:
        para = text_frame.add_paragraph()
    para.text = ""
    run = para.add_run()
    run.text = latex
    run.font.size = Pt(font_size)
    run.font.name = RAW_LATEX_FONT
    para.space_after = 4
