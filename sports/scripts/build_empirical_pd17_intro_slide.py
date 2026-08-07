#!/usr/bin/env python3
"""Build PD17 deck opener — empirical MBB glossary (text-only).

Replaces Phase B fake-league intro on CHAR_PD17_HAND slide 1.
Literal LaTeX in Calibri — no $...$; highlight → Insert → Equation → LaTeX to math.

Run (repo root):
  python sports/scripts/build_empirical_pd17_intro_slide.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_BULLET_LEAD,
    HAND_BULLET_LINE_SPACING,
    HAND_CLAIM_PT,
    HAND_TITLE_PT,
    populate_paragraph_raw_latex,
    _set_char_bullet,
)
from hero_gallery_paths import AUTO_EMPIRICAL_PD17_INTRO_DECK, ensure_hero_dirs

OUT_PPTX = AUTO_EMPIRICAL_PD17_INTRO_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.28)
LEFT_W = Inches(6.35)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

TITLE = r"PD17 — Empirical MBB characterization (real rosters, not curve fitting)"

GLOSSARY_HEAD = r"Symbols (mini glossary)"
GLOSSARY_HEADERS = ("Symbol", "Plain name", "On this deck")
GLOSSARY_ROWS = [
    (r"\hat{A}_{i}", "Player ability", r"PPM z within season on filtered panel (2011–2021)"),
    (r"\hat{T}_{j}", "Realized team talent", r"Mean \hat{A}_{i} on roster — not T_{j^*}"),
    (r"T_{j^*}", "Sim assignment target", r"Pass B only — not observed in NCAA data"),
    (r"\rho", "Assignment assortativity", r"Sim soft-match knob; overlap forensics on data"),
    (r"L_C", "Team congestion", r"mean_j \sigma(\gamma(\hat{A}_{j} - \theta)) per roster"),
    (r"\lambda", "Congestion weight", r"Score S_i = A_i - \lambda L_C (sim / later on data)"),
    (r"\theta", "Viability cutline", r"F^{-1}_{\hat{A}}(1-K/N) on empirical \hat{A}_{i}"),
    (r"\gamma", "Sigmoid sharpness", r"Reshapes team L_C — 539 placeholder until lock"),
    (r"K/N", "Draft rate", r"Accepted / total in filtered panel (\approx 1.8\%)"),
]

NOTES_HEAD = r"How to read slides 2–7"
NOTES = [
    r"Pipeline on real MBB 2011–2021: ASSIGN (overlap) \rightarrow SCORE (L_C, \gamma) \rightarrow later SELECT.",
    r"Phase B deck (separate): same pipeline in a fake league — characterization, not NCAA fit.",
    r"Filters: PPM z within season, min 20 min, poolq winsor 0.01–0.99, all teams.",
    r"Not curve-fitting to the hero yet — descriptive inputs Alex asked for before pinning \rho, \theta, \gamma, \lambda.",
]

SLIDE_MAP_HEAD = r"Deck map (assign \rightarrow score)"
SLIDE_MAP = [
    r"Slide 2 — \hat{A}_{i}, \hat{T}_{j} inputs",
    r"Slide 3 — team interval overlap (\rho diagnostic, 530 CELL 8)",
    r"Slide 4 — team L_C distribution",
    r"Slide 5 — Sketch A: \hat{T}_{j} vs L_C",
    r"Slide 6 — \gamma sweep on real rosters",
    r"Slide 7 — sim \rho calibration capstone (empirical vs 539 coverage)",
]

CLAIM = (
    "Claim (PD17): Map empirical properties of real rosters — prerequisite for pinning "
    r"\rho and \theta from data, then \gamma and \lambda without hero-curve matching."
)


def _set_cell_raw_latex(cell, text: str, *, font_size: int = 9, bold: bool = False) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    populate_paragraph_raw_latex(tf.paragraphs[0], text, font_size=font_size, bold=bold)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_glossary_table(slide, *, left, top, width, height) -> None:
    head_box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    populate_paragraph_raw_latex(
        head_box.text_frame.paragraphs[0], GLOSSARY_HEAD, font_size=12, bold=True
    )

    table_top = top + Inches(0.3)
    table_h = height - Inches(0.3)
    n_rows = 1 + len(GLOSSARY_ROWS)
    shape = slide.shapes.add_table(n_rows, 3, left, table_top, width, table_h)
    table = shape.table

    col_fracs = (0.16, 0.30, 0.54)
    for i, frac in enumerate(col_fracs):
        table.columns[i].width = int(width * frac)

    for j, header in enumerate(GLOSSARY_HEADERS):
        _set_cell_raw_latex(table.cell(0, j), header, font_size=9, bold=True)

    for i, (sym, plain, controls) in enumerate(GLOSSARY_ROWS, start=1):
        _set_cell_raw_latex(table.cell(i, 0), sym, font_size=9)
        _set_cell_raw_latex(table.cell(i, 1), plain, font_size=9)
        _set_cell_raw_latex(table.cell(i, 2), controls, font_size=8)

    for j in range(3):
        cell = table.cell(0, j)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)


def _add_bullet_section(
    slide,
    *,
    left,
    top,
    width,
    height,
    head_text: str,
    bullets: list[str],
    head_size: int = 11,
    bullet_size: int = HAND_BODY_PT,
) -> None:
    from pptx.util import Pt

    head_h = Inches(0.28)
    head_box = slide.shapes.add_textbox(left, top, width, head_h)
    populate_paragraph_raw_latex(
        head_box.text_frame.paragraphs[0], head_text, font_size=head_size, bold=True
    )

    bullet_box = slide.shapes.add_textbox(left, top + head_h, width, height - head_h)
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        populate_paragraph_raw_latex(para, HAND_BULLET_LEAD + line.lstrip(), font_size=bullet_size)
        para.line_spacing = HAND_BULLET_LINE_SPACING
        _set_char_bullet(para)


def main() -> None:
    ensure_hero_dirs()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0], TITLE, font_size=HAND_TITLE_PT, bold=True
    )

    body_top = title_top + title_h + Inches(0.08)
    footer_h = Inches(0.42)
    body_h = SLIDE_H - MARGIN - footer_h - body_top

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_glossary_table(slide, left=left_x, top=body_top, width=LEFT_W, height=body_h)

    notes_h = Inches(1.85)
    map_h = body_h - notes_h - Inches(0.08)

    _add_bullet_section(
        slide,
        left=right_x,
        top=body_top,
        width=RIGHT_W,
        height=notes_h,
        head_text=NOTES_HEAD,
        bullets=NOTES,
    )
    _add_bullet_section(
        slide,
        left=right_x,
        top=body_top + notes_h + Inches(0.06),
        width=RIGHT_W,
        height=map_h,
        head_text=SLIDE_MAP_HEAD,
        bullets=SLIDE_MAP,
    )

    foot = slide.shapes.add_textbox(MARGIN, SLIDE_H - MARGIN - footer_h, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — PD17 empirical intro / glossary)")


if __name__ == "__main__":
    main()
