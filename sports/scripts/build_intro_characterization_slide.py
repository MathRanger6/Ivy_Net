#!/usr/bin/env python3
"""Build Phase B intro slide (text-only deck opener).

Literal LaTeX in Calibri — no $...$; highlight → Insert → Equation → LaTeX to math.

Run (repo root):
  python sports/scripts/build_intro_characterization_slide.py

Outputs:
  HEROs_and_PASSes/slides/auto/CHAR_intro_characterization_AUTO.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from gallery_knobs import (
    HERO_K_OVER_N,
    HERO_N_SELECTED,
    HERO_SEED,
    LAMBDA_FIXED_RHO,
    LAMBDA_MODERATE,
    PRESET,
    THETA_MODE,
    hero_league_n,
    resolve_viability_theta,
)
from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_BULLET_LEAD,
    HAND_BULLET_LINE_SPACING,
    HAND_CLAIM_PT,
    HAND_TITLE_PT,
    populate_paragraph_raw_latex,
    _set_char_bullet,
)
from hero_gallery_paths import AUTO_INTRO_DECK, ensure_hero_dirs

OUT_PPTX = AUTO_INTRO_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.28)
LEFT_W = Inches(6.35)
RIGHT_W = CONTENT_W - LEFT_W - COL_GAP

TITLE = r"Phase B — Fake-league characterization (not curve fitting)"

GLOSSARY_HEAD = r"Knobs (mini glossary)"
GLOSSARY_HEADERS = ("Symbol", "Plain name", "What it controls")
GLOSSARY_ROWS = [
    (r"A_{i}", "Player ability", r"Innate talent draw — Beta(2,2) on [0,1] (539 preset)"),
    (
        r"T_{j^*}",
        "Sim assignment target",
        r"Synthetic iid Uniform[0,1] per team; soft-assign attractor — not realized roster talent",
    ),
    (
        r"T_{j}",
        "Realized team talent",
        r"Mean A_{i} on team j roster after ASSIGN (not T_{j^*})",
    ),
    (
        r"\rho",
        "Assignment assortativity",
        r"Soft match strength: players \rightarrow teams with T_{j^*} \approx A_{i}",
    ),
    (
        r"L_C",
        "Team congestion",
        r"Team-level mean \sigma(\gamma(A_{j}-\theta)) on roster — same L_C for all players on team j",
    ),
    (
        r"\lambda",
        "Congestion weight in score",
        r"How hard L_C penalizes ranking — not congestion itself",
    ),
    (r"\theta", "Viability cutline", r"Center of the sigmoid (F_A^{-1}(1-K/N) in k_over_n mode)"),
    (r"\gamma", "Sigmoid sharpness", r"Steepness of viability step around \theta"),
    (r"K/N", "Selectivity", r"Fraction of league selected (K \div N)"),
]

NOTES_HEAD = r"How to read this deck"
NOTES = [
    r"Not curve fitting: which rules bend synthetic selection curves — not NCAA fit yet.",
    r"Three steps: ASSIGN (\rho) \rightarrow SCORE (S_i = A_i - \lambda L_C) \rightarrow SELECT (top-K); then VISUALIZE by pool-mean bin.",
    r"L_C in this deck: team smooth (crowding_smooth_team) — mean peer viability over the full roster.",
    r"One-at-a-time (OAT): each slide moves one knob; other settings stay at benchmark values.",
    rf"Seed {HERO_SEED} on stochastic steps: draw A_i, T_{{j^*}}; soft-\rho assignment. Score, top-K, bins are deterministic.",
    rf"League: N={hero_league_n()}, K={HERO_N_SELECTED}, K/N={HERO_K_OVER_N:.0%} ({PRESET} preset).",
]

BENCHMARK_HEAD = r"Benchmark values (while sweeping other knobs)"
_theta_bench = resolve_viability_theta(preset=0.72)
BENCHMARK_BULLETS = [
    r"A_i: Beta(2,2) on [0,1] — SELECTION_539_ABILITY_DRAW",
    r"T_{j^*}: Uniform[0,1] iid per team — draw_target_means; synthetic targets, not NCAA data",
    r"T_j: realized roster mean after ASSIGN — not drawn; Sketch A 2D y-axis",
    rf"\lambda={LAMBDA_MODERATE:g} — tier1_539_reference_settings.json / playground default",
    (
        rf"\theta={_theta_bench:.3f} — F_A^{{-1}}(1-K/N) (naive draft on sim draw)"
        if THETA_MODE != "preset"
        else r"\theta=0.72, \gamma=10 — same 539 reference JSON"
    ),
    rf"\rho={LAMBDA_FIXED_RHO:g} — gallery fixed assign for \lambda / \theta slides (moderate-high mixing)",
    rf"K/N={HERO_K_OVER_N:.0%} — characterization default (MBB draft \approx 1% is a separate calibration point)",
]
if THETA_MODE != "preset":
    BENCHMARK_BULLETS.insert(5, r"\gamma=10 — 539 placeholder until lock on real rosters")

EMPIRICAL_HEAD = r"Empirical data (Layer A) vs this deck (Layer C)"
EMPIRICAL_BULLETS = [
    r"Real hero (Layer A): draft rate vs teammate-quality ventiles — bend at the top bins (Pass A / So_Far_).",
    r"This deck: fraction selected vs pool mean in a fake league — same kind of question, different axis and league.",
    r"Benchmarks anchor to the Alex 539 playground track; not re-fit to NCAA data in Phase B.",
]

CLAIM = (
    r"Claim: Phase B maps which generative knobs move who gets selected — "
    r"prerequisite before Phase C calibration to the hero."
)


def _set_cell_raw_latex(cell, text: str, *, font_size: int = 9, bold: bool = False) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    populate_paragraph_raw_latex(tf.paragraphs[0], text, font_size=font_size, bold=bold)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_glossary_table(slide, *, left, top, width, height) -> None:
    head_box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    populate_paragraph_raw_latex(
        head_box.text_frame.paragraphs[0], GLOSSARY_HEAD, font_size=12, bold=True
    )

    table_top = top + Inches(0.3)
    table_h = height - Inches(0.3)
    n_rows = 1 + len(GLOSSARY_ROWS)
    shape = slide.shapes.add_table(n_rows, 3, left, table_top, width, table_h)
    table = shape.table

    col_fracs = (0.14, 0.28, 0.58)
    for i, frac in enumerate(col_fracs):
        table.columns[i].width = int(width * frac)

    for j, header in enumerate(GLOSSARY_HEADERS):
        _set_cell_raw_latex(table.cell(0, j), header, font_size=9, bold=True)

    for i, (sym, plain, controls) in enumerate(GLOSSARY_ROWS, start=1):
        _set_cell_raw_latex(table.cell(i, 0), sym, font_size=9)
        _set_cell_raw_latex(table.cell(i, 1), plain, font_size=9)
        _set_cell_raw_latex(table.cell(i, 2), controls, font_size=8)

    for j in range(3):
        cell = table.cell(0, j)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)


def _add_bullet_section(
    slide,
    *,
    left,
    top,
    width,
    height,
    head_text: str,
    bullets: list[str],
    head_size: int = 11,
    bullet_size: int = HAND_BODY_PT,
) -> None:
    head_h = Inches(0.28)
    head_box = slide.shapes.add_textbox(left, top, width, head_h)
    populate_paragraph_raw_latex(
        head_box.text_frame.paragraphs[0], head_text, font_size=head_size, bold=True
    )

    bullet_box = slide.shapes.add_textbox(left, top + head_h, width, height - head_h)
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        populate_paragraph_raw_latex(para, HAND_BULLET_LEAD + line.lstrip(), font_size=bullet_size)
        para.line_spacing = HAND_BULLET_LINE_SPACING
        _set_char_bullet(para)


def main() -> None:
    ensure_hero_dirs()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0], TITLE, font_size=HAND_TITLE_PT, bold=True
    )

    body_top = title_top + title_h + Inches(0.08)
    footer_h = Inches(0.42)
    body_h = SLIDE_H - MARGIN - footer_h - body_top

    left_x = MARGIN
    right_x = MARGIN + LEFT_W + COL_GAP

    _add_glossary_table(slide, left=left_x, top=body_top, width=LEFT_W, height=body_h)

    notes_h = Inches(2.05)
    bench_h = Inches(1.55)
    emp_h = body_h - notes_h - bench_h - Inches(0.12)

    _add_bullet_section(
        slide,
        left=right_x,
        top=body_top,
        width=RIGHT_W,
        height=notes_h,
        head_text=NOTES_HEAD,
        bullets=NOTES,
    )
    _add_bullet_section(
        slide,
        left=right_x,
        top=body_top + notes_h + Inches(0.04),
        width=RIGHT_W,
        height=bench_h,
        head_text=BENCHMARK_HEAD,
        bullets=BENCHMARK_BULLETS,
    )
    _add_bullet_section(
        slide,
        left=right_x,
        top=body_top + notes_h + bench_h + Inches(0.08),
        width=RIGHT_W,
        height=emp_h,
        head_text=EMPIRICAL_HEAD,
        bullets=EMPIRICAL_BULLETS,
    )

    foot = slide.shapes.add_textbox(MARGIN, SLIDE_H - MARGIN - footer_h, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — Phase B intro / glossary, raw LaTeX)")


if __name__ == "__main__":
    main()
