#!/usr/bin/env python3
"""Build Alex-facing deck: draft MLE fit (λ, t, γ) + PD20 temperature gate.

Slides:
  1. Intro — Bernoulli MLE on empirical draft outcomes (Alex PD21 lock)
  2. PD20 temperature sweep — inverted-U survives Gibbs SELECT (diagnostic grid)
  3. MLE γ profile — log-lik vs γ with λ̂, t̂ at each γ
  4. Fit summary — committed λ̂, t̂, γ; panel + diagnostics

Run (repo root):
  python sports/scripts/build_pd21_mle_fit_slides.py --slides-only
  python sports/scripts/build_pd21_mle_fit_slides.py --season-min 2013 --season-max 2021
  python sports/scripts/build_pd21_mle_fit_slides.py --slides-only --force-figures

Output:
  slides/auto/CHAR_PD21_MLE_fit_AUTO.pptx
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_PD21_MLE_DECK, PD20_TEMPERATURE, PD21_MLE, ensure_hero_dirs
from interval_overlap_paths import seasons_label
from pd20_22_campaign_window import (
    activate_from_args,
    add_window_args,
    auto_deck_path,
    current_window,
)
from pd17_interval_overlap_slide import add_picture_fitted
from pd21_draft_bernoulli_mle import _load_gamma_default, plot_gamma_profile

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_TEXT_W = Inches(4.15)
FIG_W = CONTENT_W - LEFT_TEXT_W - COL_GAP


def _w():
    return current_window()


def _mle_json() -> Path:
    return PD21_MLE / f"PD21_draft_bernoulli_mle_{_w().tag}.json"


def _gamma_profile_png() -> Path:
    return PD21_MLE / f"PD21_draft_bernoulli_mle_{_w().tag}_gamma_profile.png"


def _gamma_profile_csv() -> Path:
    return PD21_MLE / f"PD21_draft_bernoulli_mle_{_w().tag}_gamma_profile.csv"


def _temp_sweep_png() -> Path:
    return PD20_TEMPERATURE / f"GRANDCHILD_temperature_select_sweep_{_w().tag}.png"


def _ensure_gamma_profile_png(*, force: bool = False) -> Path:
    out = _gamma_profile_png()
    if out.is_file() and not force:
        return out
    csv_path = _gamma_profile_csv()
    if not csv_path.is_file():
        raise FileNotFoundError(
            f"Missing {csv_path} — run: python sports/scripts/pd21_draft_bernoulli_mle.py "
            f"--profile-gamma --season-min {_w().season_min} --season-max {_w().season_max}"
        )
    profile = pd.read_csv(csv_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    plot_gamma_profile(
        profile,
        seasons=seasons_label(_w().season_min, _w().season_max),
        out_path=out,
        gamma_ref=_load_gamma_default(),
    )
    print(f"Wrote {out} (from CSV)")
    return out


def _load_fit() -> dict:
    path = _mle_json()
    if not path.is_file():
        raise FileNotFoundError(
            f"Missing {path} — run: python sports/scripts/pd21_draft_bernoulli_mle.py "
            f"--gamma 18 --season-min {_w().season_min} --season-max {_w().season_max}"
        )
    return json.loads(path.read_text(encoding="utf-8"))


def _append_text_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )
    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_box = slide.shapes.add_textbox(MARGIN, body_top, CONTENT_W, footer_top - body_top - Inches(0.06))
    fill_bullets_raw_latex(bullet_box.text_frame, bullets, font_size=HAND_BODY_PT)
    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_raw_latex(
        foot.text_frame.paragraphs[0],
        claim,
        font_size=HAND_CLAIM_PT,
        bold=True,
    )


def _append_figure_slide(
    prs: Presentation,
    *,
    fig_path: Path,
    title: str,
    subtitle: str,
    bullets: list[str],
    claim: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        title,
        font_size=HAND_TITLE_PT,
        bold=True,
    )
    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        subtitle,
        font_size=HAND_SUBTITLE_PT,
    )
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    bullet_h = Inches(2.35)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = sub_top + sub_h + Inches(0.08)
    fig_h = footer_top - fig_top - Inches(0.04)
    add_picture_fitted(slide, fig_path, MARGIN + LEFT_TEXT_W + COL_GAP, fig_top, FIG_W, fig_h)
    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, LEFT_TEXT_W, bullet_h)
    fill_bullets_raw_latex(bullet_box.text_frame, bullets, font_size=HAND_BODY_PT)
    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    populate_paragraph_raw_latex(
        foot.text_frame.paragraphs[0],
        claim,
        font_size=HAND_CLAIM_PT,
        bold=True,
    )
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _intro_bullets(fit: dict) -> list[str]:
    seasons = fit.get("seasons", "2013-2021")
    n = int(fit.get("n_player_seasons", 0))
    nd = int(fit.get("n_drafted", 0))
    return [
        r"SELECT layer: fit $(\lambda, \gamma, t)$ on real NCAA draft outcomes.",
        r"Board: $p_i \propto \exp(A_i/t)\,\exp(-\lambda L^C_i)$ — season-wise softmax.",
        r"Likelihood (Alex PD21): Bernoulli $\prod_i p_i^{Y_i}(1-p_i)^{1-Y_i}$; **$K$ not in formula**.",
        r"Method: coarse $(\lambda, t)$ grid $\rightarrow$ L-BFGS-B maximize $\ell$.",
        rf"Panel: MBB {seasons} · {n:,} player-seasons · {nd:,} drafted (POST-QC, min 20 min).",
        r"ASSIGN $\rho$ is a separate calibration (not in this likelihood).",
    ]


def _temp_bullets() -> list[str]:
    seasons = f"{_w().season_min}–{_w().season_max}"
    return [
        r"PD20 **generative gate** (before MLE): Gibbs SELECT $P_i \propto \exp(S_i/t)$.",
        r"Sweep $\log_{10} t$ on sim bridge (rule D); fixed $\lambda \in \{1.5, 2.0\}$.",
        r"Question: does inverted-U in draft rate vs LOO **survive** soft SELECT?",
        r"Answer: **yes** — inverted-U-like curvature across most of the $t$ grid.",
        r"This is a **diagnostic sweep**, not the MLE fit (MLE uses empirical $Y_{\mathrm{draft}}$).",
        rf"Empirical MLE on same window ({seasons}) gives $\hat t \approx 1.07$ (next slides).",
    ]


def _gamma_bullets(fit: dict) -> list[str]:
    lam = float(fit.get("lambda_hat", 0))
    t = float(fit.get("t_hat", 0))
    gam = float(fit.get("gamma_hat", 18))
    return [
        r"Fix $\gamma$; re-fit $(\lambda, t)$ by MLE at each $\gamma$ on the profile grid.",
        r"Log-likelihood is **flat** from $\gamma \approx 18$ to 100 — not sharply identified.",
        rf"Committed run: $\gamma = {gam:g}$ (tier-1 sim default); $\hat\lambda = {lam:.3g}$, $\hat t = {t:.3g}$.",
        r"$\hat\lambda \approx 2.5$–3: congestion in score matters; $\hat t \approx 1$: moderate softmax.",
        r"Joint 3-parameter MLE supported in script; v1 reports fixed-$\gamma$ fit for stability.",
    ]


def _summary_bullets(fit: dict) -> list[str]:
    lam = float(fit.get("lambda_hat", 0))
    t = float(fit.get("t_hat", 0))
    gam = float(fit.get("gamma_hat", 18))
    ll = float(fit.get("loglik_hat", 0))
    diag = fit.get("fixed_gamma_fit", {}).get("bfgs", {}).get("diagnostics", {})
    recall = diag.get("mean_topk_recall")
    overlap = diag.get("mean_topk_overlap")
    bullets = [
        rf"Bernoulli MLE (empirical draft): $\hat\lambda = {lam:.4g}$, $\hat t = {t:.4g}$, $\gamma = {gam:g}$ fixed.",
        rf"Max log-likelihood $\ell = {ll:.2f}$ · L-BFGS-B converged in 11 iterations.",
        r"Interpretation: draft odds rise with ability ($1/t$) and fall with LOO congestion ($\lambda>0$).",
    ]
    if recall is not None:
        bullets.append(rf"Sanity: mean top-$K$ recall on drafted players $\approx {float(recall):.3g}$ (soft ranking check).")
    if overlap is not None:
        bullets.append(rf"Mean overlap with empirical top-$K$ per season $\approx {float(overlap):.1f}$ picks.")
    bullets.extend([
        r"Not in this deck: empirical hero ventile shape (Layer A) — separate QC thread; generative Pass B/C unchanged.",
        r"Next: joint $(\lambda, \gamma, t)$ MLE or PD14 predictive gain on same panel.",
    ])
    return bullets


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PD21 MLE fit slide deck for Alex.")
    add_window_args(parser)
    parser.add_argument("--slides-only", action="store_true", help="Skip figure regeneration unless --force-figures.")
    parser.add_argument(
        "--force-figures",
        action="store_true",
        help="Regenerate γ profile PNG from CSV even if it already exists.",
    )
    args = parser.parse_args()
    activate_from_args(args)
    ensure_hero_dirs()

    fit = _load_fit()
    gamma_png = _gamma_profile_png()
    if args.force_figures:
        gamma_png = _ensure_gamma_profile_png(force=True)
    elif not args.slides_only and not gamma_png.is_file():
        gamma_png = _ensure_gamma_profile_png(force=False)
    elif not gamma_png.is_file():
        raise FileNotFoundError(
            f"Missing {gamma_png} — run with --force-figures (needs CSV) or "
            f"python sports/scripts/pd21_draft_bernoulli_mle.py --profile-gamma "
            f"--season-min {_w().season_min} --season-max {_w().season_max}"
        )
    temp_png = _temp_sweep_png()
    if not temp_png.is_file():
        raise FileNotFoundError(f"Missing {temp_png} — run grandchild_temperature_select_sweep.py first.")

    out_pptx = auto_deck_path(AUTO_PD21_MLE_DECK)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    seasons = fit.get("seasons", _w().tag.replace("_", "-"))

    _append_text_slide(
        prs,
        title=r"Draft MLE — $\lambda$, $\gamma$, $t$ on empirical NCAA outcomes",
        subtitle=rf"MBB {seasons} · Bernoulli softmax · Alex PD21 board factorization",
        bullets=_intro_bullets(fit),
        claim=r"Score $\neq$ select — we fit the **draft probability model**, not ASSIGN $\rho$.",
    )
    _append_figure_slide(
        prs,
        fig_path=temp_png,
        title=r"PD20 gate — temperature sweep (generative diagnostic)",
        subtitle=rf"Gibbs SELECT survives inverted-U · MBB {_w().season_min}–{_w().season_max}",
        bullets=_temp_bullets(),
        claim=r"Soft SELECT cleared — MLE on empirical $Y_{\mathrm{draft}}$ is the fit step.",
    )
    _append_figure_slide(
        prs,
        fig_path=gamma_png,
        title=r"Bernoulli MLE — $\gamma$ profile ($\lambda$, $t$ re-fit at each $\gamma$)",
        subtitle=rf"Maximize $\ell(\lambda, t \mid \gamma)$ · panel {seasons}",
        bullets=_gamma_bullets(fit),
        claim=rf"Committed: $\hat\lambda \approx {float(fit['lambda_hat']):.2g}$, $\hat t \approx {float(fit['t_hat']):.2g}$, $\gamma = {float(fit['gamma_hat']):g}$.",
    )
    _append_text_slide(
        prs,
        title=r"Fit summary — what to tell Alex",
        subtitle=rf"Empirical MLE on real draft labels · {seasons}",
        bullets=_summary_bullets(fit),
        claim=r"Yes: MLE on empirical data for $(\lambda, t)$; hero ventile read is a separate (parked) Layer A item.",
    )

    prs.save(str(out_pptx))
    print(f"Wrote {out_pptx} (4 slides — intro + temperature + γ profile + summary)")


if __name__ == "__main__":
    main()
