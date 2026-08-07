#!/usr/bin/env python3
"""Build one PD17 empirical γ sweep slide (HAND slide ~15 layout — full-width strip).

Slide text = literal LaTeX in Calibri — highlight → Insert → Equation → LaTeX to math.

Run (repo root):
  python sports/scripts/build_empirical_gamma_lc_slide.py
  python sports/scripts/build_empirical_gamma_lc_slide.py --slides-only
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_EMPIRICAL_GAMMA_LC_DECK, EMPIRICAL_PD17, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "empirical_gamma_lc_sweep.py"
FIG_STRIP = EMPIRICAL_PD17 / "EMPIRICAL_L_C_gamma_sweep_strip.png"
META = EMPIRICAL_PD17 / "EMPIRICAL_L_C_gamma_sweep_meta.json"
OUT_PPTX = AUTO_EMPIRICAL_GAMMA_LC_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.2)
READOUT_COL_W = (CONTENT_W - 2 * COL_GAP) / 3

CLAIM = (
    "Claim (PD17): On real rosters, γ reshapes team L_C — high γ compresses "
    r"viability near 0 on the PPM-z scale; low γ spreads L_C into the interior of [0,1]."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figure() -> None:
    print("Regenerating empirical γ sweep strip ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def _summary_by_gamma(meta: dict) -> dict[float, dict]:
    return {float(row["gamma"]): row for row in meta.get("summary", [])}


def _readout_bullets(meta: dict) -> list[str]:
    by_g = _summary_by_gamma(meta)
    bullets = [
        r"OAT on real panel: fixed \theta (K/N cutline on PPM z); only \gamma in \sigma(\gamma(\hat{A}-\theta)) changes.",
        r"Arms: \gamma \in {10, 5, 1, 0.5, 0.001} — same team-seasons, same rosters.",
    ]
    g10 = by_g.get(10.0)
    g0 = by_g.get(0.001)
    if g10 and g0:
        bullets.append(
            rf"At \gamma=10: mean L_C={g10.get('L_C_mean', 0):.3f}, "
            rf"{100 * g10.get('frac_L_C_below_0.05', 0):.0f}\% teams below 0.05."
        )
        bullets.append(
            rf"At \gamma \approx 0: mean L_C={g0.get('L_C_mean', 0):.3f}, "
            rf"max L_C={g0.get('L_C_max', 0):.3f} — spread opens on z-scale."
        )
    bullets.append(
        r"Pairs with Phase B \gamma slide — same viability knob, empirical rosters not sim."
    )
    return bullets


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slides-only",
        action="store_true",
        help="Skip PNG regen (text/layout only)",
    )
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figure()
    meta = _load_meta()
    kn = meta.get("theta_K_over_N", {})
    theta = meta.get("theta")
    k_over_n = kn.get("K_over_N")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        r"PD17 — Empirical \gamma sweep: team L_C on real rosters",
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    seasons = meta.get("seasons", "2011-2021")
    sub_parts = [rf"MBB {seasons} · team smooth L_C · PPM z · min 20 min"]
    if theta is not None and k_over_n is not None:
        sub_parts.append(
            rf"\theta = {theta:.3f} z-units (K/N={k_over_n:.4f}) fixed across \gamma arms"
        )
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        " · ".join(sub_parts),
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    readout_h = Inches(0.85)
    readout_top = footer_top - readout_h - Inches(0.08)
    fig_top = sub_top + sub_h + Inches(0.08)
    fig_h = readout_top - fig_top - Inches(0.06)

    _add_picture_fitted(slide, FIG_STRIP, MARGIN, fig_top, CONTENT_W, fig_h)

    readouts = _readout_bullets(meta)
    for i, text in enumerate(readouts[:3]):
        x = MARGIN + i * (READOUT_COL_W + COL_GAP)
        box = slide.shapes.add_textbox(x, readout_top, READOUT_COL_W, readout_h)
        tf = box.text_frame
        tf.word_wrap = True
        populate_paragraph_raw_latex(tf.paragraphs[0], text, font_size=9)

    if len(readouts) > 3:
        extra_box = slide.shapes.add_textbox(
            MARGIN, readout_top + readout_h - Inches(0.02), CONTENT_W * 0.65, Inches(0.28)
        )
        populate_paragraph_raw_latex(
            extra_box.text_frame.paragraphs[0], readouts[3], font_size=9
        )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — empirical γ strip, HAND slide-15 layout)")


if __name__ == "__main__":
    main()
