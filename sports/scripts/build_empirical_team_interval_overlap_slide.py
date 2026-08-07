#!/usr/bin/env python3
"""Build one PD17 empirical team interval overlap slide (530 CELL 8 / ρ diagnostic).

Bullets bottom-left; 2×2 figure enlarged on the right (L_C sidebar layout, swapped).

Slide text = literal LaTeX in Calibri — highlight → Insert → Equation → LaTeX to math.
PD17 bullets: use fill_bullets_raw_latex (1.5 line spacing); never $...$ in slide strings.

Run (repo root):
  python sports/scripts/build_empirical_team_interval_overlap_slide.py
  python sports/scripts/build_empirical_team_interval_overlap_slide.py --slides-only
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_EMPIRICAL_OVERLAP_DECK, EMPIRICAL_PD17, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "empirical_team_interval_overlap.py"
FIG = EMPIRICAL_PD17 / "EMPIRICAL_team_interval_overlap.png"
META = EMPIRICAL_PD17 / "EMPIRICAL_team_interval_overlap_meta.json"
OUT_PPTX = AUTO_EMPIRICAL_OVERLAP_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_TEXT_W = Inches(4.15)
FIG_W = CONTENT_W - LEFT_TEXT_W - COL_GAP

CLAIM = (
    "Claim (PD17): Real NCAA rosters leave massively overlapping team talent windows — "
    r"the empirical target \rho was meant to match in sim (530 CELL 8 / 538 Plot A)."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figure() -> None:
    print("Regenerating empirical team interval overlap figure ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def _readout_bullets(meta: dict) -> list[str]:
    cov_max = meta.get("coverage_max")
    cov_frac = meta.get("coverage_frac_gt_1")
    dis_max = meta.get("coverage_disjoint_max")
    span = meta.get("perf_span", {})
    n_ts = meta.get("n_team_seasons")

    bullets = [
        r"Each team-season: talent window [\min \hat{A}_{i}, \max \hat{A}_{i}] on PPM z.",
        r"Coverage = how many windows cover a point on the spectrum (530 CELL 8).",
        r"Red dashed: disjoint sort-and-chop on same player-seasons (537 B analog).",
    ]
    if cov_max is not None and dis_max is not None:
        bullets.append(
            rf"Actual max coverage={cov_max:,}; sort-and-chop max={dis_max}."
        )
    if cov_frac is not None:
        bullets.append(rf"{100 * cov_frac:.0f}\% of grid points have >1 team covering.")
    if span and n_ts:
        bullets.append(
            rf"Roster span: mean={span.get('mean', 0):.2f} z, "
            rf"median={span.get('median', 0):.2f} z ({n_ts:,} team-seasons)."
        )
    bullets.append(
        r"Pairs with Phase B \rho slide — sim \rho dials overlap between these curves."
    )
    return bullets


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slides-only",
        action="store_true",
        help="Skip PNG regen (text/layout only)",
    )
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figure()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        r"PD17 — Empirical team \hat{A}_{i} interval overlap (\rho diagnostic)",
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    seasons = meta.get("seasons", "2011-2021")
    n_ts = meta.get("n_team_seasons")
    sub_parts = [
        rf"MBB {seasons} · PPM z within season · min 20 min · poolq winsor 0.01–0.99",
    ]
    if n_ts:
        sub_parts.append(rf"{n_ts:,} team-seasons (530 CELL 8 port)")
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        " · ".join(sub_parts),
        font_size=HAND_SUBTITLE_PT,
    )

    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_top = sub_top + sub_h + Inches(0.08)
    bullet_h = Inches(2.35)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_top = body_top
    fig_h = footer_top - fig_top - Inches(0.04)
    fig_left = MARGIN + LEFT_TEXT_W + COL_GAP

    _add_picture_fitted(slide, FIG, fig_left, fig_top, FIG_W, fig_h)

    readouts = _readout_bullets(meta)
    bullet_box = slide.shapes.add_textbox(MARGIN, bullet_top, LEFT_TEXT_W, bullet_h)
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True
    fill_bullets_raw_latex(bullet_tf, readouts, font_size=HAND_BODY_PT)

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — empirical overlap / ρ diagnostic)")


if __name__ == "__main__":
    main()
