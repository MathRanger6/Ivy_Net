#!/usr/bin/env python3
"""Build one-slide soft-assign γ × λ sweep deck (overlap regime).

Companion to sort-and-chop γ slides — ADD to HAND deck; does not replace them.

Run (repo root):
  python sports/scripts/build_soft_assign_gamma_sweep_slide.py
  python sports/scripts/build_soft_assign_gamma_sweep_slide.py --slides-only

Outputs:
  HEROs_and_PASSes/slides/auto/CHAR_soft_assign_gamma_sweep_AUTO.pptx
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_knobs import LAMBDA_FIXED_RHO
from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import AUTO_SOFT_ASSIGN_GAMMA_DECK, SOFT_ASSIGN_LAMBDA, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "soft_assign_gamma_sweep_diagnostic.py"
FIG = SOFT_ASSIGN_LAMBDA / "GAMMA_sweep_soft_assign_lambda_curves_key_arms.png"
META = SOFT_ASSIGN_LAMBDA / "GAMMA_sweep_soft_assign_meta.json"
OUT_PPTX = AUTO_SOFT_ASSIGN_GAMMA_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.45)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.22)
LEFT_W = Inches(4.55)
FIG_W = CONTENT_W - LEFT_W - COL_GAP

CLAIM = (
    r"Claim: under soft assign (overlapping rosters), \gamma reshapes team L_C and "
    r"\lambda in score bends selection — the regime where we will pin \lambda on data. "
    r"No \lambda_{\mathrm{crit}} \approx 4/\gamma law here (sort-and-chop only)."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figure() -> None:
    print("Regenerating soft-assign γ × λ sweep figure ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return
    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2


def _readout_bullets(meta: dict) -> list[str]:
    assign = meta.get("assignment", {})
    theta = assign.get("theta")
    rho = assign.get("rho_fixed", LAMBDA_FIXED_RHO)
    bullets = [
        r"Soft assign at fixed \rho — overlapping talent windows (not sort-and-chop).",
        rf"\rho={rho:g} fixed; same roster per \gamma row; team L_C; \theta(K/N) on sim draw.",
        r"Arms: \gamma \in \{5, 10, 20\}; within each panel \lambda \in \{0, 0.25, 0.55, 0.75, 1\}.",
        r"Figure: key \lambda arms only; full grid in GAMMA_sweep_soft_assign_lambda_curves.png.",
        r"Pairs with sort-and-chop \gamma slides (benchmark) — keep both in HAND deck.",
    ]
    if theta is not None:
        bullets.insert(2, rf"\theta = {theta:.3f} z-units (F_A^{{-1}}(1-K/N) on Beta draw).")
    return bullets


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--slides-only", action="store_true")
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figure()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.28), CONTENT_W, Inches(0.46))
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        r"Phase B — \gamma \times \lambda sweep (soft assign, overlapping rosters)",
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    body_top = Inches(0.84)
    footer_h = Inches(0.52)
    footer_top = SLIDE_H - MARGIN - footer_h
    bullet_h = Inches(2.35)
    bullet_top = footer_top - bullet_h - Inches(0.06)
    fig_h = bullet_top - body_top - Inches(0.04)
    fig_left = MARGIN + LEFT_W + COL_GAP

    _add_picture_fitted(slide, FIG, fig_left, body_top, FIG_W, fig_h)

    bullet_box = slide.shapes.add_textbox(MARGIN, body_top + Inches(0.32), LEFT_W, bullet_top - body_top - Inches(0.32))
    fill_bullets_raw_latex(bullet_box.text_frame, _readout_bullets(meta), font_size=HAND_BODY_PT)

    head_box = slide.shapes.add_textbox(MARGIN, body_top, LEFT_W, Inches(0.28))
    populate_paragraph_raw_latex(
        head_box.text_frame.paragraphs[0],
        r"Overlap regime — same readout as sort-and-chop grid, different ASSIGN world",
        font_size=11,
        bold=True,
    )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (1 slide — soft-assign γ × λ; add to HAND, keep sort-and-chop)")


if __name__ == "__main__":
    main()
