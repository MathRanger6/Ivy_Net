#!/usr/bin/env python3
"""Build one-slide deck — actual $A_i$ and $T_{j^*}$ draws (gallery seed).

Run (repo root):
  python sports/scripts/build_sim_input_distributions_slide.py

Writes:
  HEROs_and_PASSes/slides/auto/CHAR_sim_input_distributions_AUTO.pptx
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_knobs import HERO_N_TEAMS, HERO_ROSTER_SIZE, HERO_SEED, PRESET, hero_league_n
from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex
from hero_gallery_paths import AUTO_SIM_INPUTS_DECK, SIM_INPUTS, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "sim_league_input_distributions.py"
FIG = SIM_INPUTS / "SIM_league_Ai_Tj_distributions.png"
META = SIM_INPUTS / "SIM_league_Ai_Tj_meta.json"
OUT_PPTX = AUTO_SIM_INPUTS_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.2)
SIDEBAR_W = Inches(3.35)
FIG_W = CONTENT_W - SIDEBAR_W - COL_GAP


def _claim_text() -> str:
    return (
        rf"Claim: Phase B fixes one stochastic league draw (seed {HERO_SEED}) — "
        r"$A_i$ and $T_{j^*}$ are the inputs every OAT arm reuses; "
        "only ASSIGN / SCORE knobs change across slides."
    )


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate() -> None:
    print("Regenerating sim input distribution figure ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def main() -> None:
    ensure_hero_dirs()
    _regenerate()
    meta = _load_meta()
    ai = meta.get("A_i", {})
    tj = meta.get("T_j_star", meta.get("T_j", {}))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.22)
    title_h = Inches(0.48)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_with_latex(
        title_box.text_frame.paragraphs[0],
        r"Phase B — Simulated league inputs: actual $A_i$ and $T_{j^*}$ draws",
        font_size=18,
    )
    title_box.text_frame.paragraphs[0].font.bold = True

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.26)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    populate_paragraph_with_latex(
        sub_box.text_frame.paragraphs[0],
        rf"{PRESET} preset · seed {HERO_SEED} · "
        rf"$N={hero_league_n():,}$ players · {HERO_N_TEAMS} teams × roster {HERO_ROSTER_SIZE}",
        font_size=9,
    )

    body_top = sub_top + sub_h + Inches(0.1)
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_h = footer_top - body_top - Inches(0.08)

    bullets = [
        r"One draw reused across all OAT arms — only placement / score rules change.",
        r"Left: $A_i \sim$ Beta(2,2) on [0,1] — unimodal, mean $\approx 0.5$.",
        r"Right: $T_{j^*} \sim$ iid Uniform[0,1] per team — sim assignment targets.",
        r"Orange curve = theoretical pdf scaled to histogram bin width.",
    ]
    if ai and tj:
        bullets.append(
            rf"This run: mean $A_i={ai.get('mean', 0):.3f}$, mean $T_{{j^*}}={tj.get('mean', 0):.3f}$ "
            rf"(sd {ai.get('std', 0):.3f} / {tj.get('std', 0):.3f})."
        )
    bullets.append(r"Not NCAA data — fake league inputs for characterization.")

    side_box = slide.shapes.add_textbox(MARGIN, body_top, SIDEBAR_W, body_h)
    side_tf = side_box.text_frame
    side_tf.word_wrap = True
    fill_bullets_latex(side_tf, bullets, font_size=9)

    _add_picture_fitted(
        slide,
        FIG,
        MARGIN + SIDEBAR_W + COL_GAP,
        body_top,
        FIG_W,
        body_h,
    )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_with_latex(foot_tf.paragraphs[0], _claim_text(), font_size=10)
    foot_tf.paragraphs[0].font.bold = True
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
