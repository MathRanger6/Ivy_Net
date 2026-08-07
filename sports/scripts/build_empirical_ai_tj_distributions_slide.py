#!/usr/bin/env python3
"""Build one-slide deck — empirical \\hat{A}_{i} and \\hat{T}_{j} (PD17).

Slide text = literal LaTeX in Calibri (braced subscripts: \\hat{A}_{i}).
Highlight → Insert → Equation → LaTeX to math.

Run (repo root):
  python sports/scripts/build_empirical_ai_tj_distributions_slide.py
  python sports/scripts/build_empirical_ai_tj_distributions_slide.py --slides-only
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from gallery_mathtext import (
    HAND_BODY_PT,
    HAND_CLAIM_PT,
    HAND_SUBTITLE_PT,
    HAND_TITLE_PT,
    fill_bullets_raw_latex,
    populate_paragraph_raw_latex,
)
from hero_gallery_paths import (
    AUTO_EMPIRICAL_AI_TJ_DECK,
    EMPIRICAL_PD17,
    ensure_hero_dirs,
)

DIAG_SCRIPT = SCRIPTS / "empirical_ai_tj_distributions.py"
FIG = EMPIRICAL_PD17 / "EMPIRICAL_Ai_Tj_distributions.png"
META = EMPIRICAL_PD17 / "EMPIRICAL_Ai_Tj_meta.json"
OUT_PPTX = AUTO_EMPIRICAL_AI_TJ_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.42)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.2)
SIDEBAR_W = Inches(3.35)
FIG_W = CONTENT_W - SIDEBAR_W - COL_GAP

CLAIM = (
    "Claim (PD17): World-first inputs from the real MBB panel — "
    r"\hat{A}_{i} per player-season and \hat{T}_{j} per team-season. "
    r"Not T_{j^*} (sim assignment targets)."
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {}


def _regenerate_figure() -> None:
    print("Regenerating empirical A_i / T_j distribution figure ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slides-only",
        action="store_true",
        help="Skip PNG regen (text/layout only)",
    )
    args = parser.parse_args()

    ensure_hero_dirs()
    if not args.slides_only:
        _regenerate_figure()
    meta = _load_meta()
    ai = meta.get("A_i_hat", meta.get("A_i", {}))
    tj = meta.get("T_j_hat", meta.get("T_j", {}))
    kn = meta.get("theta_K_over_N", {})

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_top = Inches(0.24)
    title_h = Inches(0.5)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    populate_paragraph_raw_latex(
        title_box.text_frame.paragraphs[0],
        r"PD17 — Empirical MBB inputs: \hat{A}_{i} and \hat{T}_{j}",
        font_size=HAND_TITLE_PT,
        bold=True,
    )

    sub_top = title_top + title_h + Inches(0.04)
    sub_h = Inches(0.27)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, sub_h)
    seasons = meta.get("seasons", "2011-2021")
    populate_paragraph_raw_latex(
        sub_box.text_frame.paragraphs[0],
        rf"MBB {seasons} · PPM z within season · min 20 min · poolq winsor 0.01–0.99",
        font_size=HAND_SUBTITLE_PT,
    )

    body_top = sub_top + sub_h + Inches(0.1)
    footer_h = Inches(0.46)
    footer_top = SLIDE_H - MARGIN - footer_h
    body_h = footer_top - body_top - Inches(0.08)

    bullets = [
        r"Realized rosters — descriptive inputs, not sim draws.",
        r"Left: \hat{A}_{i} — player ability (one value per player-season).",
        r"Right: \hat{T}_{j} — mean \hat{A}_{i} on roster (team-season).",
        r"\hat{T}_{j} is tighter than \hat{A}_{i} (averaging shrinks spread).",
    ]
    if ai and tj:
        bullets.append(
            rf"This run: mean \hat{{A}}_{{i}}={ai.get('mean', 0):.3f}, "
            rf"mean \hat{{T}}_{{j}}={tj.get('mean', 0):.3f} "
            rf"(sd {ai.get('std', 0):.3f} / {tj.get('std', 0):.3f})."
        )
    if kn:
        bullets.append(
            rf"PD17 \theta proxy: K/N={kn.get('K_over_N', 0):.4f} "
            rf"({kn.get('n_accepted', 0):,} drafted / {kn.get('n_total', 0):,} player-seasons)."
        )

    side_box = slide.shapes.add_textbox(MARGIN, body_top, SIDEBAR_W, body_h)
    side_tf = side_box.text_frame
    side_tf.word_wrap = True
    fill_bullets_raw_latex(side_tf, bullets, font_size=HAND_BODY_PT)

    _add_picture_fitted(
        slide,
        FIG,
        MARGIN + SIDEBAR_W + COL_GAP,
        body_top,
        FIG_W,
        body_h,
    )

    foot = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    foot_tf = foot.text_frame
    foot_tf.word_wrap = True
    populate_paragraph_raw_latex(foot_tf.paragraphs[0], CLAIM, font_size=HAND_CLAIM_PT, bold=True)
    foot_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
