#!/usr/bin/env python3
"""Build two-slide PD16 Sketch A deck — Charles hand-deck layout (Aug 2026).

Run (repo root):
  python sports/scripts/build_lc_congestion_characterization_slide.py

PD16 figures (team L_C + naive-draft θ):
  export GALLERY_LC_MODE=team_smooth GALLERY_THETA_MODE=k_over_n GALLERY_OUTPUT_SUFFIX=_pd16
  python sports/scripts/build_lc_congestion_characterization_slide.py

Writes:
  HEROs_and_PASSes/slides/auto/CHAR_lc_congestion_characterization_AUTO.pptx
    slide 1 — L_C strip + two-column readout + mean-L_C note
    slide 2 — 2D heatmap (figure right) + left bullet column
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

REPO = Path(__file__).resolve().parents[2]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))
from gallery_knobs import (
    HERO_N_TEAMS,
    HERO_SEED,
    LC_MODE,
    OUTPUT_SUFFIX,
    PRESET,
    gallery_mode_subtitle,
)
from gallery_mathtext import fill_bullets_latex, populate_paragraph_with_latex
from hero_gallery_paths import AUTO_LC_CONGESTION_DECK, PASS_C_RHO, ensure_hero_dirs

DIAG_SCRIPT = SCRIPTS / "lc_distribution_vs_rho_diagnostic.py"
META = PASS_C_RHO / f"LC_distribution_vs_rho_meta{OUTPUT_SUFFIX}.json"
FIG_STRIP = PASS_C_RHO / f"LC_distribution_vs_rho_1d_strip{OUTPUT_SUFFIX}.png"
FIG_2D = PASS_C_RHO / f"LC_distribution_vs_rho_2d{OUTPUT_SUFFIX}.png"
OUT_PPTX = AUTO_LC_CONGESTION_DECK

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.35)
CONTENT_W = SLIDE_W - 2 * MARGIN
COL_GAP = Inches(0.2)
READOUT_COL_W = (CONTENT_W - 2 * COL_GAP) / 3
SIDEBAR_W = Inches(3.15)
FIG_2D_W = CONTENT_W - SIDEBAR_W - COL_GAP

CLAIM_STRIP = (
    "Claim: $\\rho$ stratifies team congestion across the league — low-$\\rho$ mixing "
    "yields one narrow $L_C$ hump; high-$\\rho$ assortativity piles weak teams at "
    "$L_C\\approx 0$ and elite rosters at higher $L_C$."
)

CLAIM_2D = (
    "Claim: at high $\\rho$, better teams systematically face more peer pressure — "
    "the upward cloud in realized $T_j$ vs $L_C$. That gives $\\rho$ a model-internal "
    "read beyond “the Pass C selection plot moved.”"
)


def _load_meta() -> dict:
    if META.is_file():
        return json.loads(META.read_text(encoding="utf-8"))
    return {
        "preset": PRESET,
        "seed": HERO_SEED,
        "theta": 0.72,
        "gamma": 10.0,
        "lc_bins": 48,
        "summary": [],
    }


def _regenerate_figures() -> None:
    print("Regenerating L_C distribution vs ρ figures ...")
    subprocess.run([sys.executable, str(DIAG_SCRIPT)], cwd=str(REPO), check=True)


def _add_picture_fitted(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        box = slide.shapes.add_textbox(left, top, max_width, Inches(0.4))
        box.text_frame.text = f"[Missing figure: {img_path.name}]"
        return

    pic = slide.shapes.add_picture(str(img_path), left, top, width=max_width)
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2


def _mean_lc_note(meta: dict) -> str:
    summary = meta.get("summary", [])
    if summary:
        mean_lc = float(summary[0].get("L_C_mean", 0.15))
        lc_str = f"{mean_lc:.2f}".lstrip("0")  # ".15" not "0.15"
        return rf"Of Note: The mean $L_C$ is almost identical for all plots, hovering $\approx {lc_str}$"
    return r"Of Note: The mean $L_C$ is almost identical for all plots, hovering $\approx .15$"


def _rho32_bullet(meta: dict) -> str:
    summary = {row["arm"]: row for row in meta.get("summary", [])}
    rho32 = summary.get("rho_very_high", {})
    frac32 = float(rho32.get("frac_L_C_below_0.05", 0.48))
    return (
        rf"At $\rho=32$: ${100*frac32:.0f}\%$ of teams below $L_C=0.05$ - "
        r"weak rosters - almost no peer pressure."
    )


def _add_title_block(slide, *, title: str, meta: dict) -> float:
    """Title + PD16 subtitle; return y where figure body starts."""
    theta = float(meta.get("theta", 0.72))
    lc_bins = int(meta.get("lc_bins", 48))

    title_top = Inches(0.16)
    title_h = Inches(0.52)
    title_box = slide.shapes.add_textbox(MARGIN, title_top, CONTENT_W, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    populate_paragraph_with_latex(tf.paragraphs[0], title, font_size=16)
    tf.paragraphs[0].font.bold = True

    sub_top = title_top + title_h + Inches(0.02)
    sub_box = slide.shapes.add_textbox(MARGIN, sub_top, CONTENT_W, Inches(0.26))
    populate_paragraph_with_latex(
        sub_box.text_frame.paragraphs[0],
        gallery_mode_subtitle(theta_value=theta)
        + rf" · {lc_bins} bins · {HERO_N_TEAMS} teams · seed {HERO_SEED}",
        font_size=9,
    )
    return float(sub_top + Inches(0.28))


def _add_claim_footer(slide, *, claim: str, top: float) -> None:
    foot_h = Inches(0.5)
    foot = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, foot_h)
    populate_paragraph_with_latex(foot.text_frame.paragraphs[0], claim, font_size=10)
    foot.text_frame.paragraphs[0].font.bold = True
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _add_strip_slide(prs: Presentation, meta: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fig_top = _add_title_block(
        slide,
        title=(
            r"Phase B — Team $L_C$ vs $\rho$ (PD16 Sketch A) -- "
            r"How Congestion alone ($L_C$) varies with change in $\rho$"
        ),
        meta=meta,
    )

    footer_top = float(SLIDE_H - MARGIN - Inches(0.5))
    note_h = Inches(0.28)
    readout_h = Inches(0.72)
    note_top = footer_top - note_h - Inches(0.04)
    readout_top = note_top - readout_h - Inches(0.05)
    fig_h = readout_top - fig_top - Inches(0.06)

    _add_picture_fitted(slide, FIG_STRIP, MARGIN, Inches(fig_top), CONTENT_W, Inches(fig_h))

    readout_texts = [
        r"low $\rho$ ($\rho=0.001$, $1$) = one narrow spike near $L_C\approx 0.15$",
        r"High $\rho$ = mass at $L_C\approx 0$ (weak rosters) + second hump at elite $L_C$",
        _rho32_bullet(meta),
    ]
    for i, text in enumerate(readout_texts):
        x = MARGIN + i * (READOUT_COL_W + COL_GAP)
        box = slide.shapes.add_textbox(x, Inches(readout_top), READOUT_COL_W, readout_h)
        tf = box.text_frame
        tf.word_wrap = True
        populate_paragraph_with_latex(tf.paragraphs[0], text, font_size=9)

    note_box = slide.shapes.add_textbox(MARGIN, Inches(note_top), CONTENT_W * 0.55, note_h)
    populate_paragraph_with_latex(
        note_box.text_frame.paragraphs[0],
        _mean_lc_note(meta),
        font_size=9,
    )
    note_box.text_frame.paragraphs[0].font.bold = True

    _add_claim_footer(slide, claim=CLAIM_STRIP, top=footer_top)


def _add_2d_sidebar_slide(prs: Presentation, meta: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    body_top = _add_title_block(
        slide,
        title=(
            r"Phase B — Realized $T_j$ vs $L_C$ (PD16 Sketch A — co-movement) -- "
            r"Team talent $T_j$ vs congestion ($L_C$) as $\rho$ varies."
        ),
        meta=meta,
    )

    footer_top = float(SLIDE_H - MARGIN - Inches(0.5))
    body_h = footer_top - body_top - Inches(0.08)

    sidebar_x = MARGIN
    fig_x = MARGIN + SIDEBAR_W + COL_GAP

    bullets = [
        r"Each cell: # teams at (realized $T_j$, team $L_C$).",
        r"Color = count.",
        r"Low $\rho$: compact blob — $T_j$ and congestion barely co-vary.",
        r"High $\rho$: upward tail — as $T_j$ rises, $L_C$ rises.",
        r"Weak rosters at high $\rho$: low $T_j$ + low $L_C$",
        r"Elite rosters: high $T_j$ + high $L_C$.",
    ]
    side_box = slide.shapes.add_textbox(sidebar_x, Inches(body_top), SIDEBAR_W, Inches(body_h))
    fill_bullets_latex(side_box.text_frame, bullets, font_size=9)

    _add_picture_fitted(
        slide,
        FIG_2D,
        fig_x,
        Inches(body_top),
        FIG_2D_W,
        Inches(body_h),
    )

    _add_claim_footer(slide, claim=CLAIM_2D, top=footer_top)


def main() -> None:
    ensure_hero_dirs()
    _regenerate_figures()
    meta = _load_meta()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_strip_slide(prs, meta)
    _add_2d_sidebar_slide(prs, meta)

    prs.save(str(OUT_PPTX))
    print(f"Wrote {OUT_PPTX} (2 slides — Charles hand-deck layout)")
    if LC_MODE != "team_smooth" or OUTPUT_SUFFIX != "_pd16":
        print(
            "Note: for Alex PD16 figures, run with "
            "GALLERY_LC_MODE=team_smooth GALLERY_THETA_MODE=k_over_n GALLERY_OUTPUT_SUFFIX=_pd16"
        )


if __name__ == "__main__":
    main()
