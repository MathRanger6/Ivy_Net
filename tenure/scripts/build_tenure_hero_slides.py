#!/usr/bin/env python3
"""Build PowerPoint for tenure HERO runs (MBB slide layout).

Run (repo root, after tenure_pass_a_hero.py):
  python tenure/scripts/build_tenure_hero_slides.py
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SPORTS_SCRIPTS = REPO / "sports" / "scripts"
HERO_DIR = REPO / "3-Master_Plan/re_entry/HEROs_and_PASSes/tenure_sandbox/hero"
DEFAULT_OUT = HERO_DIR / "TENURE_HERO_slides_AUTO.pptx"
SCRIPTS_DIR = REPO / "tenure" / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))
from tenure_grain_labels import (  # noqa: E402
    normalize_stat,
    normalize_window,
    perf_display_label,
    window_badge,
)

sys.path.insert(0, str(SPORTS_SCRIPTS))
from build_hero_permutation_slides import (  # noqa: E402
    CONTENT_W,
    MARGIN,
    SLIDE_H,
    SLIDE_W,
    _add_mono_box,
    _add_picture,
    _add_title,
    _shape_line,
)


LOCK_LINES = [
    "Tenure HERO: HIGH + MEDIUM OpenAlex inference.",
    "ASST-PS · mean · peer LOO (annum) = v0 default · N≈796.",
    "LAST-PS · cum · peer LOO = final assistant year (MBB exit cross-section) · N≈732–805.",
    "Rates among resolved only (Option A); yellow badge = window on each plot.",
]


def _load_json(path: Path) -> dict[str, Any]:
    if path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def _spec_window_stat(spec: dict[str, Any]) -> tuple[str, str]:
    if "window" in spec:
        return normalize_window(spec["window"]), normalize_stat(spec.get("stat", "annum"))
    return normalize_window(spec.get("grain", "asst_ps")), normalize_stat(
        spec.get("pool_perf", "annum")
    )


def _hero_catalog() -> list[dict[str, Any]]:
    return [
        {
            "tag": "q16_infHM_resolved_v0",
            "title": "Tenure HERO v0 — Q16 ASST-PS mean peer LOO",
            "subtitle": "Baseline lock · mean poolq_LOO over assistant years",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--output-tag q16_infHM_resolved_v0"
            ),
        },
        {
            "tag": "q20_loo_infHM",
            "title": "Tenure HERO — Q20 ASST-PS mean (dip robustness)",
            "subtitle": "Higher bin count · peer LOO (annum) · ASST-PS mean",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--n-bins 20 --output-tag q20_loo_infHM"
            ),
        },
        {
            "tag": "q16_lastps_loo_cum_infHM",
            "title": "Tenure HERO — Q16 LAST-PS cum peer LOO",
            "subtitle": "Final assistant year · peer cumulative pubs · MBB last-ps analog",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--window last_ps --stat cum "
                "--output-tag q16_lastps_loo_cum_infHM"
            ),
        },
        {
            "tag": "q20_lastps_loo_cum_infHM",
            "title": "Tenure HERO — Q20 LAST-PS cum peer LOO",
            "subtitle": "Dip / shoulder robustness · peer cum LOO at exit year",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--window last_ps --stat cum --n-bins 20 "
                "--output-tag q20_lastps_loo_cum_infHM"
            ),
        },
        {
            "tag": "q16_lastps_own_cum_infHM",
            "title": "Tenure HERO — Q16 LAST-PS own cum pubs",
            "subtitle": "Ability slice · own cum pubs at exit (F-HERO porch)",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--window last_ps --x-metric own_cum "
                "--output-tag q16_lastps_own_cum_infHM"
            ),
        },
        {
            "tag": "ew16_lastps_loo_cum_infHM",
            "title": "Tenure HERO — EW16 LAST-PS cum peer LOO",
            "subtitle": "Equal-width bins · hue = bin n · on-bar counts",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--window last_ps --stat cum --bin-method equal_width "
                "--output-tag ew16_lastps_loo_cum_infHM"
            ),
        },
        {
            "tag": "ew20_lastps_loo_cum_infHM",
            "title": "Tenure HERO — EW20 LAST-PS cum peer LOO",
            "subtitle": "Equal-width robustness · peer cum LOO at exit year",
            "command": (
                "python tenure/scripts/tenure_pass_a_hero.py "
                "--window last_ps --stat cum --bin-method equal_width --n-bins 20 "
                "--output-tag ew20_lastps_loo_cum_infHM"
            ),
        },
    ]


def _decision_hero_catalog() -> list[dict[str, Any]]:
    return [
        {
            "tag": "q16_decision_dept_loo_infHM",
            "title": "Decision HERO — Q16 dept pond LOO",
            "subtitle": "All resolved · dept career-rate LOO at decision year",
            "command": "python tenure/scripts/tenure_pass_a_decision_hero.py",
        },
        {
            "tag": "q8_decision_dept_loo_infHM",
            "title": "Decision HERO — Q8 dept pond LOO",
            "subtitle": "Coarser bins · dept career-rate LOO at decision year",
            "command": (
                "python tenure/scripts/tenure_pass_a_decision_hero.py "
                "--n-bins 8 --output-tag q8_decision_dept_loo_infHM"
            ),
        },
        {
            "tag": "q10_decision_dept_loo_infHM",
            "title": "Decision HERO — Q10 dept pond LOO",
            "subtitle": "Mid bin count · dept career-rate LOO at decision year",
            "command": (
                "python tenure/scripts/tenure_pass_a_decision_hero.py "
                "--n-bins 10 --output-tag q10_decision_dept_loo_infHM"
            ),
        },
        {
            "tag": "q16_decision_own_career_infHM",
            "title": "Decision HERO — Q16 own career rate",
            "subtitle": "Ability slice · pubs_per_career_year at decision",
            "command": (
                "python tenure/scripts/tenure_pass_a_decision_hero.py "
                "--x-metric own_career --output-tag q16_decision_own_career_infHM"
            ),
        },
        {
            "tag": "q8_decision_own_career_infHM",
            "title": "Decision HERO — Q8 own career rate",
            "subtitle": "Coarser bins · own pubs_per_career_year at decision",
            "command": (
                "python tenure/scripts/tenure_pass_a_decision_hero.py "
                "--x-metric own_career --n-bins 8 --output-tag q8_decision_own_career_infHM"
            ),
        },
        {
            "tag": "q10_decision_own_career_infHM",
            "title": "Decision HERO — Q10 own career rate",
            "subtitle": "Mid bin count · own pubs_per_career_year at decision",
            "command": (
                "python tenure/scripts/tenure_pass_a_decision_hero.py "
                "--x-metric own_career --n-bins 10 --output-tag q10_decision_own_career_infHM"
            ),
        },
    ]


def _subtitle_from_provenance(prov: dict[str, Any], entry: dict[str, Any]) -> str:
    spec = prov.get("spec") or {}
    summ = prov.get("stage9_summary") or {}
    window, stat = _spec_window_stat(spec)
    x_metric = str(spec.get("x_metric", "loo"))
    n_p = summ.get("n_persons_with_loo", "?")
    n_res = summ.get("n_resolved", "?")
    return (
        f"{window_badge(window)} · "
        f"{perf_display_label(x_metric=x_metric, stat=stat)} · "
        f"N={n_p} · res={n_res}"
    )


def _add_hero_run_slide(
    prs: Presentation,
    entry: dict[str, Any],
    *,
    slide_num: int,
) -> bool:
    base = f"HERO_tenure_{entry['tag']}"
    slide_png = HERO_DIR / f"{base}_slide.png"
    stage9_png = HERO_DIR / f"{base}.png"
    prov = _load_json(HERO_DIR / f"{base}_provenance.json")

    png = slide_png if slide_png.is_file() else stage9_png
    if not png.is_file():
        print(f"Skip (missing PNG): {base}")
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = _subtitle_from_provenance(prov, entry)
    y_after = _add_title(
        slide,
        f"Slide {slide_num} — {entry['title']}",
        subtitle=subtitle,
    )

    cmd_h = Inches(1.15)
    shape_h = Inches(0.35)
    footer_reserve = cmd_h + shape_h + Inches(0.2)
    max_img_h = SLIDE_H - y_after - footer_reserve - MARGIN
    img_bottom = _add_picture(slide, png, MARGIN, y_after, CONTENT_W, max_img_h)

    cmd_top = img_bottom + Inches(0.1)
    _add_mono_box(slide, MARGIN, cmd_top, CONTENT_W, cmd_h, entry["command"], font_size=8)

    shape_top = cmd_top + cmd_h + Inches(0.05)
    shape_box = slide.shapes.add_textbox(MARGIN, shape_top, CONTENT_W, shape_h)
    sf = shape_box.text_frame
    sf.paragraphs[0].text = _shape_line(prov.get("shape") or {})
    sf.paragraphs[0].font.size = Pt(11)
    return True


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        slide,
        "Tenure empirical HERO — inference panel",
        subtitle=f"{date.today().isoformat()} · infHM · MBB slide format",
    )
    box = slide.shapes.add_textbox(MARGIN, Inches(1.05), CONTENT_W, Inches(1.25))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(LOCK_LINES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(12)
        p.space_after = Pt(3)

    _add_mono_box(
        slide,
        MARGIN,
        Inches(2.55),
        CONTENT_W,
        Inches(0.85),
        "python tenure/scripts/tenure_pass_a_hero.py …\n"
        "python tenure/scripts/build_tenure_hero_slides.py",
        font_size=9,
    )

    slide_num = 1
    n_added = 0
    for entry in _hero_catalog():
        if _add_hero_run_slide(prs, entry, slide_num=slide_num):
            slide_num += 1
            n_added += 1

    decision_specs = _decision_hero_catalog()
    if any((HERO_DIR / f"HERO_tenure_{s['tag']}_slide.png").is_file() for s in decision_specs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title(
            slide,
            "Decision cohort HERO — dept pond (PD29)",
            subtitle="Resolved exits · whole department at decision year",
        )
        box = slide.shapes.add_textbox(MARGIN, Inches(1.05), CONTENT_W, Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate([
            "Peer X = LOO mean pubs_per_career_year among dept faculty at decision year.",
            "Cohort = all resolved infHM (391); excludes transferred.",
            "Own-career slide = ability slice (not peer context).",
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(12)

    for entry in decision_specs:
        if _add_hero_run_slide(prs, entry, slide_num=slide_num):
            slide_num += 1
            n_added += 1

    if n_added == 0:
        raise SystemExit(
            f"No HERO slide PNGs found in {HERO_DIR.relative_to(REPO)}\n"
            "Run tenure_pass_a_hero.py for at least one catalog tag."
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path.relative_to(REPO)} · {len(prs.slides)} slides (intro + {n_added} runs)")


def regenerate_all_slide_pngs() -> None:
    """Re-render *_slide.png for every catalog run (uses existing binned CSVs)."""
    script = REPO / "tenure" / "scripts" / "tenure_pass_a_hero.py"
    for entry in _hero_catalog():
        cmd = entry["command"].replace("python tenure/scripts/tenure_pass_a_hero.py ", "").split()
        full = [sys.executable, str(script), *cmd, "--slides-only"]
        print(f"Regenerating slide: {entry['tag']}")
        subprocess.run(full, cwd=REPO, check=True)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    parser.add_argument(
        "--regenerate-slides",
        action="store_true",
        help="re-render all *_slide.png from existing CSVs, then build deck",
    )
    args = parser.parse_args()
    if args.regenerate_slides:
        regenerate_all_slide_pngs()
    build_deck(args.out)


if __name__ == "__main__":
    main()
