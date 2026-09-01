#!/usr/bin/env python3
"""Build PowerPoint from tenure basic data plots (porch BDPs).

Run (repo root, after tenure_basic_plots.py):
  python tenure/scripts/build_tenure_basic_plots_slides.py
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[2]
SPORTS_SCRIPTS = REPO / "sports" / "scripts"
BDP = REPO / "3-Master_Plan/re_entry/HEROs_and_PASSes/tenure_sandbox/basic_data_plots"
DEFAULT_OUT = BDP / "TENURE_BDP_slides_AUTO.pptx"
TAG = "infHM"

sys.path.insert(0, str(SPORTS_SCRIPTS))
from build_hero_permutation_slides import (  # noqa: E402
    CONTENT_W,
    MARGIN,
    SLIDE_H,
    SLIDE_W,
    _add_mono_box,
    _add_picture,
    _add_title,
)

LOCK_LINES = [
    "Tenure HERO population: HIGH + MEDIUM OpenAlex inference · LOO computable on assistant rows.",
    "HERO lock: person-level mean poolq_LOO · Q16 quantile (Q20 for dip robustness).",
    "Rates: tenure among resolved (Option A); censored excluded from denominator.",
]


def _add_prose_box(slide, top: float, lines: list[str], *, height: float = 0.78) -> float:
    h = Inches(height)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}" if line and not line.startswith(" ") else line
        p.font.size = Pt(11)
        p.space_after = Pt(2)
    return top + h + Inches(0.06)


def _load_json(path: Path) -> dict[str, Any]:
    if path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def _summary_overlap(meta: dict) -> str:
    h = meta.get("H_sort")
    h_txt = f"{h:.3f}" if h is not None else "n/a"
    n_pools = meta.get("n_uni_year_pools", meta.get("n_team_seasons", "?"))
    n_asst = meta.get("n_assistant_years", meta.get("n_player_seasons", "?"))
    return (
        f"uni-year pools={n_pools:,} · "
        f"asst-years={n_asst:,} · "
        f"max coverage={meta.get('coverage_max', '?')} · "
        f"H_sort={h_txt}"
    )


def _summary_dist(meta: dict, key: str) -> str:
    block = meta.get(key) or {}
    return (
        f"n={block.get('n', '?'):,} · median={block.get('median', 0):.2f} · "
        f"mean={block.get('mean', 0):.2f} · sd={block.get('std', 0):.2f}"
    )


def _summary_pubs_year(meta: dict) -> str:
    all_ = meta.get("pubs_year") or {}
    pos = meta.get("pubs_year_gt0") or {}
    z = meta.get("pct_zero")
    z_txt = f"{z:.0f}% zero" if z is not None else "zero n/a"
    return (
        f"all n={all_.get('n', '?'):,} ({z_txt}) · "
        f"inset n={pos.get('n', '?'):,} · med={pos.get('median', 0):.1f} · mean={pos.get('mean', 0):.1f}"
    )


def _summary_outcome_groups(meta: dict, keys: tuple[str, ...]) -> str:
    parts = []
    for k in keys:
        block = meta.get(k) or {}
        if block.get("n", 0):
            parts.append(f"{k} n={block['n']:,} med={block.get('median', 0):.1f}")
    return " · ".join(parts) if parts else "(no meta)"


def _summary_tenured_pubs(meta: dict) -> str:
    all_ = meta.get("tenured_pubs_mean") or {}
    pos = meta.get("tenured_pubs_mean_gt0") or {}
    return (
        f"tenured n={all_.get('n', '?'):,} · med={all_.get('median', 0):.1f} · "
        f"mean={all_.get('mean', 0):.1f} · inset n={pos.get('n', '?'):,}"
    )


def _plot_catalog() -> list[dict[str, Any]]:
    return [
        {
            "key": "overlap",
            "title": "Uni-year peer pool interval overlap",
            "subtitle": "MBB PD17 / reigning overlap analog · pubs_year z within year",
            "prose": [
                "Pool unit = university × year (OpenAlex assistant peer pool).",
                "Each interval = [min, max] own pubs (z within year) across OA assistants in pool.",
                "Overlap along the spectrum motivates LOO vs full-pool pressure (F-HERO later).",
            ],
            "png": BDP / f"TENURE_pool_interval_overlap_{TAG}.png",
            "meta": BDP / f"TENURE_pool_interval_overlap_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only overlap",
            "summary_fn": _summary_overlap,
        },
        {
            "key": "overlap_gt0",
            "title": "Uni-year peer pool interval overlap (pubs > 0)",
            "subtitle": "Sensitivity — intensive margin only · compare to prior slide",
            "prose": [
                "Same pool unit; assistant-years with pubs_year = 0 dropped before z and intervals.",
                "z within year recomputed on pubs > 0 rows only (not identical to full-panel z).",
                "Compare H_sort and coverage to full-panel slide — zeros inflate overlap mass at bottom.",
            ],
            "png": BDP / f"TENURE_pool_interval_overlap_{TAG}_pubs_gt0.png",
            "meta": BDP / f"TENURE_pool_interval_overlap_{TAG}_pubs_gt0_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only overlap_gt0",
            "summary_fn": _summary_overlap,
        },
        {
            "key": "poolq_loo",
            "title": "poolq_LOO distribution (HERO grain)",
            "subtitle": "Person-level mean · N≈796 matches Q16/Q20 HERO",
            "prose": [
                "Same collapse as tenure_pass_a_hero: mean poolq_loo_mean over assistant years.",
                "Right-skew drives equal-width bin sparsity at high LOO.",
            ],
            "png": BDP / f"TENURE_poolq_loo_distribution_{TAG}.png",
            "meta": BDP / f"TENURE_poolq_loo_distribution_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only poolq_loo",
            "summary_fn": lambda m: _summary_dist(m, "loo_mean"),
        },
        {
            "key": "pubs_year",
            "title": "Own pubs_year distribution",
            "subtitle": "Unconditional + inset (pubs > 0) · inference panel",
            "prose": [
                "Main: all assistant-years (36% zero pub years in OpenAlex count).",
                "Inset: intensive margin — shape among those with ≥1 pub (N≈1,544).",
                "Main x-axis capped near p99; tail max=237 noted in legend box.",
            ],
            "png": BDP / f"TENURE_pubs_year_distribution_{TAG}.png",
            "meta": BDP / f"TENURE_pubs_year_distribution_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only pubs_year",
            "summary_fn": _summary_pubs_year,
        },
        {
            "key": "pool_size",
            "title": "LOO peer pool size distribution",
            "subtitle": "pool_size_oa_loo on assistant person-years",
            "prose": [
                "Count of OA-matched co-assistants after leave-one-out.",
                "Thin pools → noisy LOO; complements overlap slide.",
            ],
            "png": BDP / f"TENURE_pool_size_loo_distribution_{TAG}.png",
            "meta": BDP / f"TENURE_pool_size_loo_distribution_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only pool_size",
            "summary_fn": lambda m: _summary_dist(m, "pool_size_oa_loo"),
        },
        {
            "key": "pubs_by_outcome",
            "title": "Own pubs by outcome (person-level)",
            "subtitle": "Tenured vs attrition vs censored · same HERO N",
            "prose": [
                "Person-level mean pubs_year over assistant spell — descriptive, not HERO.",
                "Blue = promoted (tenure_event); orange = left as assistant; gray = still asst ≈ 2024.",
                "Compares own output to the LOO environment HERO bins on.",
            ],
            "png": BDP / f"TENURE_pubs_mean_by_outcome_{TAG}.png",
            "meta": BDP / f"TENURE_pubs_mean_by_outcome_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only pubs_by_outcome",
            "summary_fn": lambda m: _summary_outcome_groups(m, ("tenured", "attrition", "censored")),
        },
        {
            "key": "pubs_tenured",
            "title": "Promoted instructors — own pubs",
            "subtitle": "Tenured only · person-level mean + inset (mean > 0)",
            "prose": [
                "Subset to tenure_event = True (N≈174 on inference panel).",
                "Mean annual pubs over assistant years — not cumulative career total.",
                "Pairs with HERO high-bin tenure rates; inset drops zero-mean spells.",
            ],
            "png": BDP / f"TENURE_pubs_mean_tenured_{TAG}.png",
            "meta": BDP / f"TENURE_pubs_mean_tenured_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only pubs_tenured",
            "summary_fn": _summary_tenured_pubs,
        },
        {
            "key": "loo_by_outcome",
            "title": "poolq_LOO by outcome (person-level)",
            "subtitle": "Peer environment split · complements HERO x-axis",
            "prose": [
                "Same outcome colors as pubs-by-outcome slide.",
                "Shows whether promoted faculty carried higher LOO peer pressure on average.",
                "HERO asks rate(T|LOO bin); this is marginal LOO by outcome.",
            ],
            "png": BDP / f"TENURE_poolq_loo_by_outcome_{TAG}.png",
            "meta": BDP / f"TENURE_poolq_loo_by_outcome_{TAG}_meta.json",
            "command": "python tenure/scripts/tenure_basic_plots.py --only loo_by_outcome",
            "summary_fn": lambda m: _summary_outcome_groups(m, ("tenured", "attrition", "censored")),
        },
    ]


def _add_plot_slide(prs: Presentation, entry: dict, meta: dict, *, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = _add_title(slide, f"Slide {slide_num} — {entry['title']}", subtitle=entry["subtitle"])
    y = _add_prose_box(slide, y, entry["prose"])

    cmd_h = Inches(0.55)
    footer_h = Inches(0.42)
    footer_reserve = cmd_h + footer_h + Inches(0.15)
    max_img_h = SLIDE_H - y - footer_reserve - MARGIN

    img_bottom = _add_picture(slide, entry["png"], MARGIN, y, CONTENT_W, max_img_h)

    cmd_top = img_bottom + Inches(0.08)
    _add_mono_box(slide, MARGIN, cmd_top, CONTENT_W, cmd_h, entry["command"], font_size=8)

    footer_top = cmd_top + cmd_h + Inches(0.04)
    summary_fn = entry["summary_fn"]
    stats = summary_fn(meta) if meta else "(no meta JSON — re-run plot)"
    box = slide.shapes.add_textbox(MARGIN, footer_top, CONTENT_W, footer_h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = stats
    p.font.size = Pt(10)


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = _add_title(slide, "Tenure basic data plots — inference porch", subtitle=date.today().isoformat())
    _add_prose_box(slide, Inches(1.1), LOCK_LINES, height=1.2)
    _add_mono_box(
        slide,
        MARGIN,
        Inches(2.6),
        CONTENT_W,
        Inches(0.7),
        "python tenure/scripts/tenure_basic_plots.py\n"
        "python tenure/scripts/build_tenure_basic_plots_slides.py",
        font_size=9,
    )

    slide_num = 1
    for spec in _plot_catalog():
        if not spec["png"].is_file():
            print(f"Skip (missing PNG): {spec['png'].name}")
            continue
        meta = _load_json(spec["meta"])
        _add_plot_slide(prs, spec, meta, slide_num=slide_num)
        slide_num += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path.relative_to(REPO)}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()
    build_deck(args.out)


if __name__ == "__main__":
    main()
